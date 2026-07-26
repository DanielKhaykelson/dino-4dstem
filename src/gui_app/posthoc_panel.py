"""posthoc_panel.py -- Tab 4 (Post-hoc analysis).

Sidebar of one-click renders that show up in the main canvas:

    Class map (adaptive K)
    Class averages (grid, raw Cartesian, no central mask)
    Virtual BF
    Virtual HAADF
    Overlay class set on HAADF
    Centroid cosine matrix
    UMAP of embeddings
    GradCAM (per class)
    Integrated Gradients (per class)

The post-hoc tab inherits the run dir + sample from the eval tab (auto
linked) or the user can load a different run dir. Heavy renders
(GradCAM, IG) run in a worker thread so the GUI stays responsive.
"""
from __future__ import annotations
import os, sys, json, threading, time
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import numpy as np
import customtkinter as ctk
import tkinter as tk
from tkinter import filedialog, messagebox

import matplotlib
matplotlib.use("TkAgg", force=True)
from matplotlib.figure import Figure
from matplotlib.backends.backend_tkagg import (FigureCanvasTkAgg,
                                                  NavigationToolbar2Tk)
from matplotlib.colors import ListedColormap, BoundaryNorm, to_rgb

from data import SAMPLES, LoadPRZ
from gui_app.runner import list_ckpts


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------

def _radial_mask(H, W, r_in, r_out):
    yy, xx = np.ogrid[:H, :W]
    cy, cx = H / 2.0, W / 2.0
    r2 = (yy - cy) ** 2 + (xx - cx) ** 2
    return (r2 >= r_in ** 2) & (r2 < r_out ** 2)


def _open_lazy(path, scan_shape=None):
    """Lazy 4D-cube loader — supports .prz/.npz/.npy/.h5 (+ Dectris
    masters with external-linked data files). Delegates to
    data.open_lazy_cube so every panel shares the same logic + always
    passes allow_pickle=True."""
    from data import open_lazy_cube
    return open_lazy_cube(path, scan_shape=scan_shape)


def _adaptive_cmap(K_act, base_name="tab10"):
    import matplotlib.pyplot as plt
    if K_act <= 10:
        base = list(plt.get_cmap("tab10").colors[:K_act])
    elif K_act <= 20:
        base = list(plt.get_cmap("tab20").colors[:K_act])
    else:
        base = [plt.get_cmap("turbo")(i / max(K_act - 1, 1))
                for i in range(K_act)]
    return ListedColormap(base, name=f"K{K_act}")


# ---------------------------------------------------------------------------

class PostHocPanel(ctk.CTkFrame):

    def __init__(self, master, app=None):
        super().__init__(master)
        self.app = app                       # for shared resolution fields
        # state pulled in from eval/training tabs (or via "Load run dir")
        self.outdir: "str | None" = None
        self.sample: "str | None" = None
        self._inf: "dict | None" = None      # {soft_probs, assigns, embeds}
        self._cube_path: "str | None" = None
        self._scan_shape = None
        self._busy = False
        # cached intermediates
        self._BF = None
        self._HA = None
        # interactive cursor state
        self._crosshair_lines = []
        self._pick_cids = []
        self._click_cid = None
        # ACOM cache: a built+planned Crystal keyed by (cif, k_max, plan_mode).
        self._acom_crystal = None
        self._acom_crystal_key = None
        # Multi-phase ACOM state.
        self._acom_phases: list = []          # [(name, cif_path), ...]
        self._acom_phase_crystals: dict = {}  # name → cached Crystal
        self._acom_phase_keys: dict = {}      # name → cache key
        self._acom_last_mp = None             # most recent multiphase result
        self._build()

    def _recip_per_px(self):
        try:
            return float(self.app.recip_res.get()) if self.app else 0.0
        except Exception:
            return 0.0

    @staticmethod
    def _event_shift(event) -> bool:
        """Reliable shift-held test for a matplotlib mouse event.

        matplotlib's `event.key` is only populated when the canvas holds
        keyboard focus, so shift+click is frequently MISSED.  The raw Tk
        event (`event.guiEvent`) carries a modifier bitmask that is
        always correct (Shift = 0x0001 on Tk), so prefer that.
        """
        ge = getattr(event, "guiEvent", None)
        if ge is not None:
            try:
                return bool(int(ge.state) & 0x0001)
            except Exception:
                pass
        return bool(event.key) and ("shift" in str(event.key).lower())

    def _real_per_px(self):
        try:
            return float(self.app.real_res.get()) if self.app else 0.0
        except Exception:
            return 0.0

    # ---- API for main app ----
    def on_runtime_sample_added(self, key):
        try:
            keys = sorted(SAMPLES.keys())
            if hasattr(self, "_sample_menu"):
                self._sample_menu.configure(values=keys)
        except Exception:
            pass

    def link_run(self, outdir, sample, native_sample="__keep__"):
        self.outdir = outdir
        self.sample = sample
        # The run's trained dataset.  Default: the linked sample IS the
        # trained one (single-dataset runs).  `_load_dir_dialog` passes
        # native_sample=None for MULTI runs (no single trained dataset),
        # which keys every dataset's inference cache separately.
        if native_sample != "__keep__":
            self._run_native_sample = native_sample
        else:
            self._run_native_sample = sample
        self._inf = None
        self._BF = None; self._HA = None
        cfg = SAMPLES[sample]
        self._cube_path = cfg.get("path") or (cfg.get("paths") or [None])[0]
        self._scan_shape = cfg["scan_shape"]
        self._info.configure(
            text=f"linked run: {outdir}\nsample = {sample}  (auto-linked)"
                  f"    scan = {self._scan_shape}")
        self._maybe_load_inference()
        self._refresh_class_dropdown()
        # Phase B: refresh the fine-tune status panel for this sample.
        try: self._refresh_ft_status()
        except Exception: pass
        # Drive the global session — replaces the old push-style
        # `acom2.refresh_from_posthoc()` chain.  Every panel that
        # subscribes will react automatically.
        try:
            sess = getattr(self.app, "session", None)
            if sess is not None:
                sess.set(run_dir=outdir, sample=sample,
                            inference=self._inf)
        except Exception: pass
        # Legacy back-compat: still call refresh_from_posthoc on
        # acom2 so older session-unaware panels keep working.
        try:
            ac = getattr(self.app, "acom2", None)
            if ac is not None:
                ac.refresh_from_posthoc()
        except Exception: pass

    def _inference_cache_path(self):
        """Per-dataset inference cache path.

        The run's NATIVE (trained) dataset keeps the canonical
        ``eval/inference.npz`` (back-compat).  Running the model on ANY
        other dataset writes to ``eval/inference__<dataset>.npz`` so we
        never clobber the trained dataset's cached inference.  A
        multi-trained run has NO single native dataset, so every
        constituent cube is keyed.
        """
        eval_dir = os.path.join(self.outdir, "eval")
        native = getattr(self, "_run_native_sample", None)
        if native is not None and self.sample == native:
            return os.path.join(eval_dir, "inference.npz")
        safe = "".join(c if (c.isalnum() or c in "._-") else "_"
                        for c in str(self.sample))
        return os.path.join(eval_dir, f"inference__{safe}.npz")

    def _maybe_load_inference(self):
        """Load the cached inference for the CURRENT dataset if present;
        otherwise leave None and let the heavy renders run inference on
        demand."""
        if self.outdir is None:
            return
        ip = self._inference_cache_path()
        if os.path.exists(ip):
            inf = np.load(ip, allow_pickle=True)
            self._inf = dict(soft_probs=inf["soft_probs"],
                             assigns=inf["assigns"],
                             embeds=inf["embeds"])

    # ---- UI ----
    def _build(self):
        # Top row: Load run dir + sample selector + status
        top = ctk.CTkFrame(self)
        top.pack(side="top", fill="x", padx=6, pady=6)
        ctk.CTkButton(top, text="Load run dir…", width=110,
                       command=self._load_dir_dialog).pack(side="left",
                                                             padx=4)
        # The dataset comes from the Dataset tab at the top of the app —
        # no dataset dropdown / cube browser here (that caused a run-vs-cube
        # deadlock). Load the run dir; the cube follows the Dataset tab.
        self._sample_var = ctk.StringVar(value="")
        self._info = ctk.CTkLabel(top, text="(no run linked)",
                                    font=("Consolas", 10), anchor="w",
                                    justify="left")
        self._info.pack(side="left", padx=8)

        # Body: sidebar + canvas
        body = ctk.CTkFrame(self)
        body.pack(side="top", fill="both", expand=True, padx=6, pady=4)

        # Scrollable so all options stay reachable when content overflows
        # the window height (the bottom buttons used to get clipped).
        sidebar = ctk.CTkScrollableFrame(body, width=280)
        sidebar.pack(side="left", fill="y")

        # ── Fine-tune section (Phase B) ───────────────────────────
        ctk.CTkLabel(sidebar, text="Fine-tune with labels",
                      font=("Segoe UI", 12, "bold")).pack(
            anchor="w", padx=10, pady=(8, 2))
        # Active-label source
        self._ft_proposer_var = ctk.StringVar(value="cross+intra")
        ftrow = ctk.CTkFrame(sidebar, fg_color="transparent")
        ftrow.pack(fill="x", padx=10, pady=1)
        ctk.CTkLabel(ftrow, text="source:", width=60,
                      anchor="w").pack(side="left")
        ctk.CTkOptionMenu(ftrow, variable=self._ft_proposer_var,
            values=["cross+intra", "cross_class", "intra_class",
                     "scan_edge", "low_margin", "random"],
            width=140).pack(side="left", padx=2)

        # Class filters: blank → random across all classes; integer →
        # restrict to that class (intra) or pair (inter).
        self._ft_intra_class = ctk.StringVar(value="")
        self._ft_inter_a     = ctk.StringVar(value="")
        self._ft_inter_b     = ctk.StringVar(value="")
        f_intra = ctk.CTkFrame(sidebar, fg_color="transparent")
        f_intra.pack(fill="x", padx=10, pady=(2, 0))
        ctk.CTkLabel(f_intra, text="Intra:", width=60,
                      anchor="w").pack(side="left")
        ctk.CTkEntry(f_intra, textvariable=self._ft_intra_class,
                       width=40, placeholder_text="—"
                       ).pack(side="left", padx=2)
        ctk.CTkLabel(f_intra,
            text="(class id; blank = all)",
            font=("Segoe UI", 9), text_color=("#666", "#888")
            ).pack(side="left", padx=4)
        f_inter = ctk.CTkFrame(sidebar, fg_color="transparent")
        f_inter.pack(fill="x", padx=10, pady=(0, 2))
        ctk.CTkLabel(f_inter, text="Inter:", width=60,
                      anchor="w").pack(side="left")
        ctk.CTkEntry(f_inter, textvariable=self._ft_inter_a,
                       width=40, placeholder_text="—"
                       ).pack(side="left", padx=2)
        ctk.CTkLabel(f_inter, text=",",
            font=("Segoe UI", 11, "bold")).pack(side="left", padx=2)
        ctk.CTkEntry(f_inter, textvariable=self._ft_inter_b,
                       width=40, placeholder_text="—"
                       ).pack(side="left", padx=2)
        ctk.CTkLabel(f_inter,
            text="(pair; A/B order randomised)",
            font=("Segoe UI", 9), text_color=("#666", "#888")
            ).pack(side="left", padx=4)

        ctk.CTkButton(sidebar, text="Active-label pairs…",
                       width=240,
                       command=self._open_active_labeler
                       ).pack(anchor="w", padx=10, pady=2)
        # Fine-tune knobs
        self._ft_epochs        = ctk.IntVar(value=5)
        self._ft_lambda_pair   = ctk.DoubleVar(value=0.5)
        self._ft_lambda_c1d    = ctk.DoubleVar(value=0.05)
        self._ft_subsample_n   = ctk.IntVar(value=2000)
        self._ft_freeze_enc    = ctk.BooleanVar(value=True)
        ft_knobs = ctk.CTkFrame(sidebar, fg_color="transparent")
        ft_knobs.pack(fill="x", padx=10, pady=(2, 0))
        for label, var, w in [
            ("epochs",      self._ft_epochs,      60),
            ("λ_pair",      self._ft_lambda_pair, 60),
            ("λ_cluster1d", self._ft_lambda_c1d,  60),
            ("subsample",   self._ft_subsample_n, 60),
        ]:
            row = ctk.CTkFrame(ft_knobs, fg_color="transparent")
            row.pack(fill="x")
            ctk.CTkLabel(row, text=label, width=80,
                          anchor="w").pack(side="left")
            ctk.CTkEntry(row, textvariable=var, width=w
                          ).pack(side="left", padx=2)
        ctk.CTkCheckBox(sidebar, text="freeze encoder (faster, "
                          "head-only)", variable=self._ft_freeze_enc
                          ).pack(anchor="w", padx=10, pady=2)
        # Default OFF → use cumulative labels (option A). When ON,
        # only pairs whose `t` timestamp is later than the parent
        # run's `timestamp` are sent into the loss. Useful when the
        # user wants the fine-tune to act as a CORRECTION pass on
        # top of what the parent already learned.
        self._ft_only_post = ctk.BooleanVar(value=False)
        ctk.CTkCheckBox(sidebar,
                          text="only labels added since parent run",
                          variable=self._ft_only_post
                          ).pack(anchor="w", padx=10, pady=2)
        ctk.CTkButton(sidebar, text="Fine-tune with my labels  ▶",
                       width=240,
                       fg_color=("#2D7A2D", "#1F7A1F"),
                       font=("Segoe UI", 11, "bold"),
                       command=self._start_finetune
                       ).pack(anchor="w", padx=10, pady=(2, 4))
        self._ft_status = ctk.CTkLabel(sidebar, text="(no fine-tune yet)",
            font=("Consolas", 9), wraplength=240, justify="left",
            text_color=("#666", "#aaa"))
        self._ft_status.pack(anchor="w", padx=10, pady=(0, 8))

        ctk.CTkLabel(sidebar, text="Renders",
                      font=("Segoe UI", 12, "bold")).pack(
            anchor="w", padx=10, pady=(8, 4))

        # global renders
        for label, cmd in [
            ("Class map (adaptive K)",   self._render_classmap),
            ("Class averages grid",      self._render_class_avgs),
            ("Class distribution (occupancy)",
                                          self._render_class_distribution),
            ("1D radial avg (log-y)",    self._render_radial_1d),
            ("1D advanced…  (toggle / baseline / CIF)",
                                          self._open_radial_1d_popup),
            ("1D residuals (loss-side, SAXS)",
                                          self._render_radial_1d_saxs),
            ("Centroid cosine matrix",   self._render_centroid_matrix),
            ("UMAP of embeddings",       self._render_umap),
            ("Virtual BF",               self._render_bf),
            ("Virtual HAADF",            self._render_haadf),
            ("Virtual annular (r, dr)…", self._open_annular_popup),
        ]:
            ctk.CTkButton(sidebar, text=label, width=240,
                           command=cmd).pack(anchor="w", padx=10, pady=2)

        # Batch report — renders everything that doesn't need user input
        # (the buttons above + GradCAM/IG triptych for every class), then
        # bundles into PDF + PPTX + a png/ subfolder.
        ctk.CTkButton(sidebar, text="Report All  →  PPTX + PNGs",
                       width=240,
                       fg_color=("#2D7A2D", "#1F7A1F"),
                       command=self._report_all
                       ).pack(anchor="w", padx=10, pady=(8, 2))

        # ── Merge classes ────────────────────────────────────────────
        ctk.CTkLabel(sidebar, text="Merge classes",
                      font=("Segoe UI", 12, "bold")).pack(
            anchor="w", padx=10, pady=(12, 4))
        self._merge_method_var = ctk.StringVar(value="manual")
        ctk.CTkOptionMenu(sidebar, variable=self._merge_method_var,
            values=["manual", "SSIM threshold",
                       "cosine threshold (centroids)"],
            width=240,
            command=lambda _v: self._on_merge_method_change()
            ).pack(anchor="w", padx=10, pady=2)

        # Manual: X → Y
        self._merge_x_var = ctk.StringVar(value="")
        self._merge_y_var = ctk.StringVar(value="")
        self._merge_manual_row = ctk.CTkFrame(sidebar,
                                                 fg_color="transparent")
        self._merge_manual_row.pack(fill="x", padx=10, pady=2)
        ctk.CTkLabel(self._merge_manual_row, text="merge",
                       width=50).pack(side="left")
        ctk.CTkEntry(self._merge_manual_row,
                       textvariable=self._merge_x_var,
                       width=40, placeholder_text="X"
                       ).pack(side="left", padx=2)
        ctk.CTkLabel(self._merge_manual_row, text="→ keep as",
                       width=70).pack(side="left", padx=4)
        ctk.CTkEntry(self._merge_manual_row,
                       textvariable=self._merge_y_var,
                       width=40, placeholder_text="Y"
                       ).pack(side="left", padx=2)

        # Threshold (used by SSIM / cosine modes). Default is set per
        # method when the dropdown changes: cosine centroids cluster
        # tightly so 0.95 is a sensible "merge if very similar" cut;
        # SSIM on class-average diffraction patterns spans a much
        # wider range so 0.70 is a better default.
        self._merge_thresh_var = ctk.DoubleVar(value=0.95)
        self._merge_thresh_row = ctk.CTkFrame(sidebar,
                                                 fg_color="transparent")
        thr_top = ctk.CTkFrame(self._merge_thresh_row,
                                  fg_color="transparent")
        thr_top.pack(fill="x")
        ctk.CTkLabel(thr_top, text="threshold",
                       width=80).pack(side="left")
        ctk.CTkEntry(thr_top,
                       textvariable=self._merge_thresh_var,
                       width=60).pack(side="left", padx=2)
        ctk.CTkButton(thr_top, text="Preview pairs",
                        width=110,
                        command=self._preview_merge_pairs
                        ).pack(side="left", padx=4)
        # Method-aware hint line; updated by _on_merge_method_change.
        self._merge_thresh_hint = ctk.CTkLabel(self._merge_thresh_row,
            text="(pairs above this merge)",
            font=("Segoe UI", 9), text_color=("#666", "#aaa"),
            justify="left", wraplength=260)
        self._merge_thresh_hint.pack(anchor="w", padx=4, pady=(2, 0))

        btn_row = ctk.CTkFrame(sidebar, fg_color="transparent")
        btn_row.pack(fill="x", padx=10, pady=(2, 4))
        ctk.CTkButton(btn_row, text="Apply merge", width=110,
                       fg_color=("#2D7A2D", "#1F7A1F"),
                       command=self._apply_merge
                       ).pack(side="left", padx=2)
        ctk.CTkButton(btn_row, text="Reset", width=70,
                       command=self._reset_merge).pack(side="left",
                                                          padx=2)

        # View toggle — only meaningful after a merge has been applied.
        self._merge_view_var = ctk.StringVar(value="merged")
        view_row = ctk.CTkFrame(sidebar, fg_color="transparent")
        view_row.pack(fill="x", padx=10, pady=(2, 0))
        ctk.CTkLabel(view_row, text="View:", width=50,
                       anchor="w").pack(side="left")
        ctk.CTkOptionMenu(view_row,
            variable=self._merge_view_var,
            values=["merged", "original", "compare (side-by-side)"],
            width=200,
            command=lambda _v: self._switch_merge_view()
            ).pack(side="left", padx=2)

        self._merge_status = ctk.CTkLabel(sidebar, text="",
            font=("Consolas", 9), justify="left",
            text_color=("#444", "#aaa"), wraplength=260)
        self._merge_status.pack(anchor="w", padx=10, pady=(2, 4))
        # Show only the relevant control row for the chosen method.
        self._on_merge_method_change()

        ctk.CTkLabel(sidebar, text="Attribution (Avg + GradCAM + IG)",
                      font=("Segoe UI", 12, "bold")).pack(
            anchor="w", padx=10, pady=(12, 4))

        # class dropdown
        cls_row = ctk.CTkFrame(sidebar, fg_color="transparent")
        cls_row.pack(fill="x", padx=10, pady=2)
        ctk.CTkLabel(cls_row, text="class:").pack(side="left")
        self._class_var = ctk.StringVar(value="—")
        self._class_menu = ctk.CTkOptionMenu(cls_row,
                                                variable=self._class_var,
                                                values=["—"], width=80)
        self._class_menu.pack(side="left", padx=4)

        # source radio: class average vs single frame
        self._attr_source = ctk.StringVar(value="class")
        src_row = ctk.CTkFrame(sidebar, fg_color="transparent")
        src_row.pack(fill="x", padx=10, pady=(2, 0))
        ctk.CTkRadioButton(src_row, text="class avg (top 200)",
                            variable=self._attr_source,
                            value="class").pack(anchor="w")
        ctk.CTkRadioButton(src_row, text="single frame",
                            variable=self._attr_source,
                            value="frame").pack(anchor="w")

        # frame index entry (used when source = "frame")
        idx_row = ctk.CTkFrame(sidebar, fg_color="transparent")
        idx_row.pack(fill="x", padx=10, pady=(2, 4))
        ctk.CTkLabel(idx_row, text="frame idx:").pack(side="left")
        self._frame_idx_var = ctk.StringVar(value="0")
        ctk.CTkEntry(idx_row, textvariable=self._frame_idx_var,
                       width=80).pack(side="left", padx=4)
        ctk.CTkLabel(idx_row,
            text="(or click on the class map to fill)",
            font=("Segoe UI", 8), text_color=("#666", "#888")
            ).pack(side="left", padx=4)

        # GradCAM target layer (matters when the trained run used n_layers > 1).
        cam_row = ctk.CTkFrame(sidebar, fg_color="transparent")
        cam_row.pack(fill="x", padx=10, pady=(2, 4))
        ctk.CTkLabel(cam_row, text="CAM layer:").pack(side="left")
        self._cam_layer_var = ctk.StringVar(value="last")
        ctk.CTkOptionMenu(cam_row, variable=self._cam_layer_var,
            values=["last", "layer1", "layer2", "layer3", "layer4"],
            width=100).pack(side="left", padx=4)
        ctk.CTkLabel(cam_row,
            text="('last' = final stage of the trained backbone)",
            font=("Segoe UI", 8), text_color=("#666", "#888")
            ).pack(side="left", padx=4)

        ctk.CTkButton(sidebar, text="Avg + GradCAM + IG  (triptych)",
                       width=240, command=self._render_triptych
                       ).pack(anchor="w", padx=10, pady=4)

        ctk.CTkLabel(sidebar, text="Multi-class overlay",
                      font=("Segoe UI", 12, "bold")).pack(
            anchor="w", padx=10, pady=(12, 4))
        ctk.CTkLabel(sidebar,
            text="check the classes to combine, then click below",
            font=("Segoe UI", 9), text_color=("#444", "#aaa"),
            wraplength=240, justify="left").pack(anchor="w", padx=10)
        self._cls_check_frame = ctk.CTkFrame(sidebar, fg_color="transparent")
        self._cls_check_frame.pack(fill="x", padx=10, pady=4)
        self._cls_check_vars: dict = {}      # filled after inference

        for label, on_haadf in [
            ("Overlay on HAADF",  True),
            ("Overlay on BF",     False),
        ]:
            ctk.CTkButton(sidebar, text=label, width=240,
                           command=lambda h=on_haadf:
                              self._render_overlay(on_haadf=h)
                           ).pack(anchor="w", padx=10, pady=2)

        # ── ACOM has moved ──────────────────────────────────────────
        # The CIF / batch / multi-phase ACOM controls are now in their
        # own dedicated tab ("ACOM (validated)") so the user can go
        # through the step-by-step validation (spot detection → pixel
        # calibration vs CIF → goodness of fit) BEFORE running on the
        # whole set.  The right-click grain popup still exposes a
        # one-click "Run ACOM on grain" button that uses the cached
        # CIF + calibration from the ACOM tab.
        ctk.CTkLabel(sidebar, text="ACOM (validated)",
                      font=("Segoe UI", 12, "bold")).pack(
            anchor="w", padx=10, pady=(12, 4))
        _hint = ("Open the ACOM tab to validate spot detection, "
                  "calibrate 1/Å/px against the CIF, and check goodness "
                  "of fit BEFORE running on the whole dataset.  This "
                  "panel just exposes the per-grain drill-in via "
                  "right-click on the class map.")
        ctk.CTkLabel(sidebar, text=_hint,
                      font=("Segoe UI", 9),
                      text_color=("#666", "#aaa"),
                      wraplength=270, justify="left", anchor="w"
                      ).pack(anchor="w", padx=10, pady=(0, 4))
        # Keep DoubleVars used by the right-click grain popup helper
        # (so _acom_run_on_pattern still works without the sidebar).
        rp = self._recip_per_px() or 0.0185
        self._acom_cif_var = ctk.StringVar(value="")
        self._acom_inv_ang = ctk.DoubleVar(value=round(rp * 0.1, 6))
        self._acom_kmax    = ctk.DoubleVar(value=2.0)
        self._acom_plan_mode = ctk.StringVar(value="corners")

        # status (sidebar bottom)
        self._sb = ctk.CTkLabel(sidebar, text="",
                                  font=("Consolas", 9), wraplength=260,
                                  justify="left")
        self._sb.pack(side="bottom", anchor="w", padx=10, pady=8)

        # main canvas
        canvas_holder = ctk.CTkFrame(body)
        canvas_holder.pack(side="left", fill="both", expand=True,
                            padx=(8, 0))
        # toolbar uses tk widgets so wrap in a tk frame; pack ABOVE the canvas
        toolbar_frame = tk.Frame(canvas_holder, bg="#f4f4f4")
        toolbar_frame.pack(side="top", fill="x")
        self._fig = Figure(figsize=(10, 7), dpi=95, facecolor="#f4f4f4")
        self._fig.text(0.5, 0.5, "(load a run and pick a render)",
                        ha="center", va="center", fontsize=12,
                        color="#888")
        self._canvas = FigureCanvasTkAgg(self._fig, master=canvas_holder)
        self._canvas.get_tk_widget().pack(side="top", fill="both",
                                            expand=True)
        # zoom / pan / save / reset toolbar
        self._toolbar = NavigationToolbar2Tk(self._canvas, toolbar_frame,
                                               pack_toolbar=False)
        self._toolbar.update()
        self._toolbar.pack(side="left")
        # explicit "Save panel" button -> writes to <outdir>/eval/posthoc/
        ctk.CTkButton(toolbar_frame, text="Save panel (PDF+PNG)",
                       width=170, height=28,
                       command=self._save_panel).pack(side="left", padx=8)
        # Hover-coordinates label (right side of toolbar)
        self._hover_lbl = tk.Label(toolbar_frame,
            text="", font=("Consolas", 9), bg="#f4f4f4",
            anchor="e", width=46)
        self._hover_lbl.pack(side="right", padx=8)
        self._canvas.mpl_connect(
            "motion_notify_event", self._on_motion)

        # color legend strip (filled in once K is known)
        self._legend_strip = tk.Frame(canvas_holder, bg="#f4f4f4", height=28)
        self._legend_strip.pack(side="bottom", fill="x")

    def _on_motion(self, event):
        if event.inaxes and event.xdata is not None and event.ydata is not None:
            self._hover_lbl.configure(
                text=f"x={event.xdata:.4g}   y={event.ydata:.4g}")
            # update crosshair if active
            if self._crosshair_lines:
                v, h = self._crosshair_lines
                v.set_xdata([event.xdata, event.xdata])
                h.set_ydata([event.ydata, event.ydata])
                v.set_visible(True); h.set_visible(True)
                self._canvas.draw_idle()
        else:
            self._hover_lbl.configure(text="")
            if self._crosshair_lines:
                v, h = self._crosshair_lines
                v.set_visible(False); h.set_visible(False)
                self._canvas.draw_idle()

    def _enable_crosshair_on(self, ax):
        v = ax.axvline(0, color="#888", lw=0.7, ls="--", visible=False)
        h = ax.axhline(0, color="#888", lw=0.7, ls="--", visible=False)
        self._crosshair_lines = [v, h]

    def _disconnect_clicks(self):
        if self._click_cid is not None:
            try: self._canvas.mpl_disconnect(self._click_cid)
            except Exception: pass
            self._click_cid = None
        self._crosshair_lines = []
        for cid in list(getattr(self, "_pick_cids", []) or []):
            try: self._canvas.mpl_disconnect(cid)
            except Exception: pass
        self._pick_cids = []

    def _save_panel(self):
        if self.outdir is None:
            messagebox.showinfo("save", "No run linked yet."); return
        out = os.path.join(self.outdir, "eval", "posthoc")
        os.makedirs(out, exist_ok=True)
        from datetime import datetime
        stamp = datetime.now().strftime("%H%M%S")
        png = os.path.join(out, f"panel_{stamp}.png")
        pdf = os.path.join(out, f"panel_{stamp}.pdf")
        self._fig.savefig(png, dpi=200, bbox_inches="tight",
                            facecolor="white")
        self._fig.savefig(pdf, bbox_inches="tight", facecolor="white")
        self._set_status(f"saved -> {png}")

    # ---- helpers ----
    def _set_status(self, msg):
        self._sb.configure(text=msg)
        self.update_idletasks()

    def _browse_dataset(self):
        """Pick ANY cube to run the loaded model on (decouples the run
        dir / model from the dataset).  Registers it, selects it, and
        links it just like a dropdown change."""
        if self.outdir is None:
            messagebox.showinfo("dataset",
                "Load a run dir (model) first, then pick a dataset.")
            return
        p = filedialog.askopenfilename(
            title="Pick a dataset cube to run the model on",
            filetypes=[("Cube files", "*.prz *.npz *.npy *.h5 *.hdf5"),
                        ("All files", "*.*")])
        if not p:
            return
        scan_override = None
        if p.lower().endswith((".h5", ".hdf5")):
            try:
                import h5py
                from data import (_h5_find_data_path, _h5_infer_scan_shape)
                with h5py.File(p, "r") as fh:
                    dpath, ndim = _h5_find_data_path(fh)
                    s = tuple(fh[dpath].shape)
                    if ndim == 3:
                        scan_override = _h5_infer_scan_shape(fh, s[0])
                        if scan_override is None:
                            from gui_app._dialogs import ask_scan_shape
                            scan_override = ask_scan_shape(
                                self, s[0], s[1], s[2])
                            if scan_override is None:
                                return
            except Exception as e:
                messagebox.showerror("dataset",
                    f"could not read h5 shape:\n{e}"); return
        try:
            from data import register_runtime_sample
            key = register_runtime_sample(
                p, scan_shape=(tuple(scan_override) if scan_override
                                else None))
        except Exception as e:
            messagebox.showerror("dataset",
                f"could not register cube:\n{e}"); return
        try:
            self._sample_menu.configure(values=sorted(SAMPLES.keys()))
        except Exception:
            pass
        self._sample_var.set(key)
        self._on_sample_change()

    def _on_sample_change(self):
        s = self._sample_var.get()
        if s and s in SAMPLES:
            self.sample = s
            cfg = SAMPLES[s]
            self._cube_path = cfg.get("path") or (cfg.get("paths") or [None])[0]
            self._scan_shape = cfg["scan_shape"]
            # Dataset changed -> drop the old inference (the previous bug:
            # the stale class map was reused).  Load this dataset's cache
            # if present; otherwise the next render computes it.
            self._inf = None
            self._BF = None; self._HA = None
            self._maybe_load_inference()
            native = getattr(self, "_run_native_sample", None)
            if native is not None and s == native:
                link = "native (trained) dataset — auto-linked"
            else:
                link = ("OTHER dataset — model will be RUN on it "
                         "(not the trained set)")
            cached = "  [inference cached]" if self._inf is not None \
                else "  [click a render to run inference]"
            self._info.configure(
                text=f"linked run: {self.outdir}\nsample = {s}    "
                      f"scan = {self._scan_shape}\n{link}{cached}")
            try:
                self._refresh_class_dropdown()
            except Exception:
                pass
            # Keep the global session in step so ACOM/Blob/etc. follow.
            try:
                sess = getattr(self.app, "session", None)
                if sess is not None and self.outdir:
                    sess.set(run_dir=self.outdir, sample=s)
            except Exception:
                pass

    @staticmethod
    def _find_sample_lock(start_dir, max_walk=5):
        """Walk up from `start_dir` looking for a `SAMPLE_LOCK.json`.

        sweep_m_K.py drops one at the sample-level dir (i.e.
        <sweep_root>/<sample>/), so a stage2 run dir like
        ``<sweep_root>/<sample>/stage2/<name>/`` finds it two levels up.
        """
        cur = os.path.abspath(start_dir)
        for _ in range(int(max_walk)):
            p = os.path.join(cur, "SAMPLE_LOCK.json")
            if os.path.exists(p):
                return p
            parent = os.path.dirname(cur)
            if parent == cur:
                break
            cur = parent
        return None

    def _load_dir_dialog(self):
        d = filedialog.askdirectory(initialdir="runs",
            title="Pick a run dir (with eval/inference.npz preferred)")
        if not d:
            return
        # 1) Try run_summary.json for the sample name.
        rs = os.path.join(d, "run_summary.json")
        sample = self._sample_var.get()
        sample_inferred = None
        if os.path.exists(rs):
            try:
                js = json.load(open(rs))
                s = js.get("sample") or js.get("cfg", {}).get("sample") \
                    or js.get("cfg", {}).get("target_sample")
                if s:
                    sample_inferred = s
            except Exception:
                pass

        # 2a) Sweep runs (tools/sweep_m_K.py) live at
        #     <root>/<sweep_root>/<sample>/stage{1,2}/<run_name>
        # and drop SAMPLE_LOCK.json at <root>/<sweep_root>/<sample>/.
        # That file carries the EXACT cube_path + vmax + crop + polar
        # mask the sweep used for training, which may differ from any
        # static entry in `data.py` (e.g. the sweep uses
        # "IMC_150nm_SI5_nbed.cube.npy" while the static
        # `IMC_150nm_SI5` points at the .prz).  Use SAMPLE_LOCK when
        # present so the GUI's pre-processing matches the sweep's.
        # Also re-registers built-in keys (Na007b, EuInAs_B100) under
        # the sweep's cube path so all three samples are consistent.
        if sample_inferred:
            lock_path = self._find_sample_lock(d)
            if lock_path:
                try:
                    spec = json.load(open(lock_path, encoding="utf-8"))
                    cube_p = spec.get("cube_path")
                    if cube_p and os.path.exists(cube_p):
                        from data import register_runtime_sample
                        vmax = float(spec.get("vmax", 2.0))
                        pmc = int(spec.get("polar_mask_cols", 0))
                        # Match sweep convention: cart CenterMask
                        # radius = polar_mask_cols // 2 so we don't
                        # double-mask the low-r region.
                        derived_cmr = pmc // 2
                        register_runtime_sample(
                            cube_p, vmax=vmax,
                            center_mask_radius=derived_cmr,
                            key=sample_inferred,
                        )
                        try:
                            self._sample_menu.configure(
                                values=sorted(SAMPLES.keys()))
                        except Exception:
                            pass
                except Exception as e:
                    print(f"[posthoc] SAMPLE_LOCK auto-register "
                           f"failed: {e!r}", flush=True)

        # 2b) Fallback for runs created in the GUI (not the sweep):
        # `_train_kwargs.json` carries the same info under
        # `_sample_config`.  MULTI runs list every constituent cube.
        multi_keys = []
        tk_path = os.path.join(d, "_train_kwargs.json")
        if os.path.exists(tk_path):
            try:
                kw = json.load(open(tk_path, encoding="utf-8"))
                cfg = kw.get("_sample_config")
                from data import register_runtime_sample
                if cfg and cfg.get("is_multi") and cfg.get("paths"):
                    for pth in cfg["paths"]:
                        try:
                            multi_keys.append(register_runtime_sample(
                                pth, vmax=float(cfg.get("vmax", 2.0)),
                                center_mask_radius=int(
                                    cfg.get("center_mask_radius", 15))))
                        except Exception as e:
                            print(f"[posthoc] multi cube register "
                                   f"failed for {pth}: {e!r}", flush=True)
                elif (sample_inferred and sample_inferred not in SAMPLES
                        and cfg and cfg.get("path")):
                    scan_shape = cfg.get("scan_shape")
                    register_runtime_sample(
                        cfg["path"],
                        scan_shape=(tuple(scan_shape)
                                      if scan_shape else None),
                        vmax=float(cfg.get("vmax", 2.0)),
                        center_mask_radius=int(
                            cfg.get("center_mask_radius", 15)),
                        key=sample_inferred,
                    )
                try:
                    self._sample_menu.configure(
                        values=sorted(SAMPLES.keys()))
                except Exception:
                    pass
            except Exception as e:
                print(f"[posthoc] auto-register from "
                       f"_train_kwargs failed: {e!r}", flush=True)

        # If the inferred sample is itself a multi entry already present,
        # expand its constituents too.
        if (not multi_keys and sample_inferred in SAMPLES
                and SAMPLES[sample_inferred].get("is_multi")):
            from data import register_runtime_sample
            mc = SAMPLES[sample_inferred]
            for pth in mc.get("paths", []):
                try:
                    multi_keys.append(register_runtime_sample(
                        pth, vmax=float(mc.get("vmax", 2.0)),
                        center_mask_radius=int(
                            mc.get("center_mask_radius", 15))))
                except Exception:
                    pass
            try:
                self._sample_menu.configure(values=sorted(SAMPLES.keys()))
            except Exception:
                pass

        # 3) Final selection
        if multi_keys:
            # MULTI run: NO single trained dataset, so we do NOT silently
            # auto-link to one.  Populate the dropdown with the cubes and
            # link the first, but state clearly it's not auto-linked.
            first = multi_keys[0]
            self._sample_var.set(first)
            self.link_run(d, first, native_sample=None)
            self._info.configure(
                text=f"linked run: {d}\nMULTI-dataset run ({len(multi_keys)} "
                      f"cubes) — analyzing «{first}».")
            return
        if sample_inferred and sample_inferred in SAMPLES:
            sample = sample_inferred
            self._sample_var.set(sample)
        if sample not in SAMPLES:
            # Couldn't infer the run's dataset — use the cube loaded in the
            # Dataset tab (top of the app).
            pre = getattr(self.app, "pre", None)
            k = pre.get_sample_key() if pre is not None else None
            if k and k in SAMPLES:
                sample = k
                self._sample_var.set(sample)
            else:
                messagebox.showerror("No dataset",
                    "Load a cube in the Dataset tab (top of the window) "
                    "first, then load the run dir here.")
                return
        self.link_run(d, sample)

    def _ensure_inference(self) -> bool:
        """If we don't have inference cached, try to load it. If still
        missing, run infer_scan now (slow). Return True on success."""
        if self._inf is not None:
            return True
        if self.outdir is None or self.sample is None:
            messagebox.showinfo("no run", "Load a run dir first.")
            return False
        # try existing inference.npz (we may have only the training run dir)
        self._maybe_load_inference()
        if self._inf is not None:
            return True
        # else run inference fresh
        ckpt = self._best_ckpt()
        if ckpt is None:
            messagebox.showerror("no checkpoint",
                f"No best.pth or ckpt_ep*.pth found in {self.outdir}")
            return False
        self._set_status(f"running inference from {os.path.basename(ckpt)} …")
        try:
            self._inf = self._infer(ckpt)
        except Exception as e:
            messagebox.showerror("inference failed", str(e))
            return False
        self._set_status("inference done.")
        return True

    def _best_ckpt(self) -> "str | None":
        cand = os.path.join(self.outdir, "best.pth")
        if os.path.exists(cand):
            return cand
        cks = list_ckpts(self.outdir)
        if cks:
            return cks[-1][1]
        return None

    def _infer(self, ckpt_path):
        import torch
        from dino_sr_contrastive_model import load_contrastive_checkpoint
        from contrastive_eval import infer_scan
        device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        model, _, _, _ = load_contrastive_checkpoint(ckpt_path, device=device)
        model.eval()
        cfg = SAMPLES[self.sample]
        ds = LoadPRZ(cfg["path"], resize=192, vmax=cfg["vmax"])
        mask_r, mask_cols, ccrop, com = \
            self._read_polar_pipeline_cfg()
        inf = infer_scan(model, ds, device, dense_remap=True,
                          polar_size=192, polar_mask_cols=mask_cols,
                          center_crop_size=ccrop,
                          com_centering=com, center_mask_radius=mask_r,
                          eval_temp=0.06, batch_size=128)
        # cache to disk for future reuse (per-dataset key so running on a
        # non-trained dataset never overwrites the trained inference).
        eval_dir = os.path.join(self.outdir, "eval")
        os.makedirs(eval_dir, exist_ok=True)
        np.savez(self._inference_cache_path(),
                  soft_probs=inf["soft_probs"], assigns=inf["assigns"],
                  embeds=inf["embeds"],
                  K_original_ids=np.asarray(inf.get("K_original_ids", []),
                                            dtype=np.int64),
                  K_original=np.int64(inf.get("K_original", 0)))
        return dict(soft_probs=inf["soft_probs"],
                    assigns=inf["assigns"], embeds=inf["embeds"],
                    K_original_ids=inf.get("K_original_ids", []))

    def _refresh_class_dropdown(self):
        """After inference, populate class dropdown + multi-select +
        color legend strip."""
        if self._inf is None:
            return
        K = int(self._inf["soft_probs"].shape[1])
        ids = [str(c) for c in range(K)]
        self._class_menu.configure(values=ids)
        self._class_var.set("0")
        # rebuild multi-class checkboxes
        for w in self._cls_check_frame.winfo_children():
            w.destroy()
        self._cls_check_vars.clear()
        cols = 4
        for c in range(K):
            r, cc = divmod(c, cols)
            v = ctk.BooleanVar(value=False)
            self._cls_check_vars[c] = v
            ctk.CTkCheckBox(self._cls_check_frame, text=f"p{c}",
                             variable=v, width=60).grid(
                row=r, column=cc, sticky="w", padx=2, pady=2)
        # rebuild color legend strip below the canvas
        for w in self._legend_strip.winfo_children():
            w.destroy()
        cmap = _adaptive_cmap(K)
        counts = np.bincount(self._inf["assigns"], minlength=K).tolist()
        for c in range(K):
            r, g, b = cmap(c)[:3]
            hex_color = "#%02x%02x%02x" % (int(r * 255), int(g * 255),
                                              int(b * 255))
            sw = tk.Frame(self._legend_strip, bg="#f4f4f4")
            sw.pack(side="left", padx=4)
            tk.Frame(sw, bg=hex_color, width=18, height=14, relief="solid",
                      borderwidth=1).pack(side="left")
            tk.Label(sw, text=f" p{c} (N={counts[c]})  ",
                      font=("Consolas", 9), bg="#f4f4f4").pack(side="left")

    # ----------------------------------------------------------------
    # individual renders
    # ----------------------------------------------------------------
    def _new_fig(self):
        self._fig.clear()
        self._disconnect_clicks()
        return self._fig

    def _redraw(self):
        self._fig.tight_layout()
        self._canvas.draw_idle()

    # ---- Parameters table (text-only image for the report) ----------
    def _render_parameters(self):
        """Render `_train_kwargs.json` (the run's training config) as a
        formatted text image.  Used by Report All and also clickable
        on its own.  No interaction.
        """
        fig = self._new_fig()
        ax = fig.add_subplot(111)
        ax.set_axis_off()

        kw_path = os.path.join(self.outdir or "", "_train_kwargs.json")
        rs_path = os.path.join(self.outdir or "", "run_summary.json")
        blob = None
        src = "?"
        if os.path.exists(kw_path):
            try:
                with open(kw_path) as f:
                    blob = json.load(f); src = "_train_kwargs.json"
            except Exception:
                blob = None
        if blob is None and os.path.exists(rs_path):
            try:
                with open(rs_path) as f:
                    blob = json.load(f); src = "run_summary.json"
            except Exception:
                blob = None
        if blob is None:
            ax.text(0.5, 0.5,
                    "No _train_kwargs.json or run_summary.json found.",
                    ha="center", va="center", fontsize=11)
            self._redraw(); return

        # Pretty-print: top-level scalars first, then nested dicts.
        def _fmt(v):
            if isinstance(v, (str, bool)) or v is None:
                return repr(v)
            if isinstance(v, float):
                return f"{v:g}"
            if isinstance(v, list):
                if len(v) > 8:
                    return f"[…{len(v)} items]"
                return repr(v)
            return repr(v)

        # Hoist 'sample' to the top if present.
        order_keys = []
        if "sample" in blob:
            order_keys.append("sample")
        # Inner kwargs may live under 'kwargs' (gui worker format).
        inner = blob.get("kwargs", None)

        lines: list[str] = [f"# source: {src}"]
        for k in order_keys:
            lines.append(f"{k:30s}: {_fmt(blob[k])}")

        # Top-level scalars (excluding the keys we'll dive into).
        skip = {"kwargs", "_sample_config"}
        scalars = {k: v for k, v in blob.items()
                     if k not in skip and k not in order_keys
                     and not isinstance(v, (dict, list))}
        for k in sorted(scalars):
            lines.append(f"{k:30s}: {_fmt(scalars[k])}")

        if isinstance(inner, dict):
            lines.append("")
            lines.append("# kwargs:")
            for k in sorted(inner):
                v = inner[k]
                if isinstance(v, dict):
                    lines.append(f"{k}:")
                    for kk in sorted(v):
                        lines.append(f"  {kk:28s}: {_fmt(v[kk])}")
                else:
                    lines.append(f"{k:30s}: {_fmt(v)}")

        # Render in two columns if too long for one.
        text = "\n".join(lines)
        n_lines = len(lines)
        if n_lines <= 60:
            ax.text(0.01, 0.99, text, ha="left", va="top",
                     fontsize=8, family="monospace",
                     transform=ax.transAxes)
        else:
            half = (n_lines + 1) // 2
            left = "\n".join(lines[:half])
            right = "\n".join(lines[half:])
            ax.text(0.01, 0.99, left, ha="left", va="top",
                     fontsize=7.5, family="monospace",
                     transform=ax.transAxes)
            ax.text(0.51, 0.99, right, ha="left", va="top",
                     fontsize=7.5, family="monospace",
                     transform=ax.transAxes)
        fig.suptitle(f"{self.sample}  —  training parameters", fontsize=11)
        self._redraw()

    # ====== Batch report  ===========================================
    # "Report All" runs every render that doesn't need user input
    # (class map, class averages grid, 1D radial, 1D SAXS residuals,
    #  centroid cosine, UMAP, virtual BF/HAADF, GradCAM+IG triptych
    #  per class) and writes:
    #   <out>/report/png/NN_<name>.png   one png per render
    #   <out>/report/report.pdf          all pngs combined (Pillow)
    #   <out>/report/report.pptx         one slide per png (python-pptx)
    # State machine via self.after so the GUI stays responsive.
    # GradCAM compute for each class runs in a worker thread; the
    # state machine waits on a callback before stepping forward.

    def _report_all(self):
        if not self._ensure_inference():
            messagebox.showinfo(
                "Report", "Run inference first."); return
        if not self.outdir:
            messagebox.showinfo(
                "Report", "Load a run dir first."); return
        out_root = filedialog.askdirectory(
            initialdir=self.outdir,
            title="Pick output folder for report  "
                   "(a 'report/' sub-folder will be created)",
            mustexist=False)
        if not out_root:
            return
        out_dir = os.path.join(out_root, "report")
        png_dir = os.path.join(out_dir, "png")
        os.makedirs(png_dir, exist_ok=True)

        K = int(self._inf["soft_probs"].shape[1])
        # (name, payload, kind) tuples.  kind="fig" → call payload();
        # kind="triptych" → payload is a class id, run async.
        tasks = [
            ("parameters",           self._render_parameters,      "fig"),
            ("class_map",            self._render_classmap,        "fig"),
            ("class_averages_grid",  self._render_class_avgs,      "fig"),
            ("class_distribution",
             self._render_class_distribution, "fig"),
            ("radial_1d",            self._render_radial_1d,       "fig"),
            ("radial_1d_saxs",       self._render_radial_1d_saxs,  "fig"),
            ("centroid_cosine",      self._render_centroid_matrix, "fig"),
            ("umap",                 self._render_umap,            "fig"),
            ("virtual_bf",           self._render_bf,              "fig"),
            ("virtual_haadf",        self._render_haadf,           "fig"),
        ]
        for c in range(K):
            tasks.append((f"gradcam_p{c}", c, "triptych"))

        self._report_state = dict(
            tasks=tasks, idx=0,
            out_dir=out_dir, png_dir=png_dir,
            png_paths=[],
            t0=time.perf_counter(),
        )
        self._set_status(f"report: starting {len(tasks)} renders…")
        self.after(50, self._report_step)

    def _report_step(self):
        st = getattr(self, "_report_state", None)
        if st is None:
            return
        idx = st["idx"]
        tasks = st["tasks"]
        if idx >= len(tasks):
            self._report_finalize()
            return
        name, payload, kind = tasks[idx]
        self._set_status(
            f"report: ({idx + 1}/{len(tasks)}) {name} …")
        self.update_idletasks()
        if kind == "fig":
            try:
                payload()
                # Force layout/draw before capture so tight_layout has run.
                self._fig.tight_layout()
                self._canvas.draw()
            except Exception as e:
                print(f"[report] {name} failed: {e!r}", flush=True)
            self._save_report_png(name)
            st["idx"] += 1
            self.after(50, self._report_step)
        elif kind == "triptych":
            c = payload
            threading.Thread(
                target=self._report_triptych_worker,
                args=(c, name), daemon=True).start()

    def _report_triptych_worker(self, c, name):
        try:
            ckpt = self._best_ckpt()
            if ckpt is None:
                raise RuntimeError("no checkpoint")
            raw_full, avg, cam, ig = self._compute_attribution_pair(
                ckpt, c, "class", None)
            tag = f"p{c}  (top-200 weighted avg)"
            def _draw_and_step():
                try:
                    self._draw_triptych(c, "class", None, tag,
                                          avg, cam, ig,
                                          raw_full=raw_full)
                    self._fig.tight_layout()
                    self._canvas.draw()
                except Exception as e:
                    print(f"[report] {name} draw failed: {e!r}",
                          flush=True)
                self._save_report_png(name)
                self._report_state["idx"] += 1
                self.after(50, self._report_step)
            self.after(0, _draw_and_step)
        except Exception as e:
            err = str(e)
            print(f"[report] {name} compute failed: {err}", flush=True)
            def _skip():
                self._report_state["idx"] += 1
                self.after(50, self._report_step)
            self.after(0, _skip)

    def _save_report_png(self, name):
        st = self._report_state
        path = os.path.join(
            st["png_dir"], f"{st['idx']:02d}_{name}.png")
        try:
            self._fig.savefig(path, dpi=140, bbox_inches="tight",
                                facecolor=self._fig.get_facecolor())
            st["png_paths"].append((name, path))
        except Exception as e:
            print(f"[report] save {name} failed: {e!r}", flush=True)

    def _report_finalize(self):
        st = self._report_state
        paths = st["png_paths"]
        out_dir = st["out_dir"]
        elapsed = time.perf_counter() - st["t0"]
        if not paths:
            self._set_status("report: nothing rendered")
            self._report_state = None
            return

        # ---- PPTX (python-pptx) ----
        pptx_path = None
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            for name, p in paths:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                # Title strip at top
                tb = slide.shapes.add_textbox(
                    Inches(0.3), Inches(0.05),
                    Inches(12.7), Inches(0.4))
                tb.text_frame.text = name
                for para in tb.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(16)
                        run.font.bold = True
                # Image, fitted by height (preserves aspect).
                slide.shapes.add_picture(
                    p, Inches(0.3), Inches(0.55),
                    height=Inches(6.85))
            pptx_path = os.path.join(out_dir, "report.pptx")
            prs.save(pptx_path)
        except ImportError:
            print("[report] python-pptx not installed; "
                  "skipping pptx", flush=True)
        except Exception as e:
            print(f"[report] PPTX build failed: {e!r}", flush=True)

        bits = [f"{len(paths)} pngs"]
        if pptx_path: bits.append("PPTX")
        self._set_status(
            f"report: {' + '.join(bits)} → {out_dir}  "
            f"({elapsed:.0f}s)")
        self._report_state = None
    # ====== end batch report ========================================

    def _render_classmap(self):
        if not self._ensure_inference(): return
        if self._refresh_class_dropdown_lazy(): pass
        assigns = self._inf["assigns"]
        K = int(self._inf["soft_probs"].shape[1])
        Ny, Nx = self._scan_shape
        cm_array = assigns.reshape(Ny, Nx)
        fig = self._new_fig()
        ax = fig.add_subplot(111)
        cmap = _adaptive_cmap(K)
        norm = BoundaryNorm(np.arange(K + 1) - 0.5, K)
        im = ax.imshow(cm_array, cmap=cmap, norm=norm, aspect="equal",
                        interpolation="nearest")
        counts = np.bincount(assigns, minlength=K).tolist()
        # axis units: scan pixels → nm if calibrated; add 100 nm scale bar.
        nm_per_px = self._real_per_px()
        if nm_per_px > 0:
            ax.set_xlabel(f"x  ({nm_per_px:.3g} nm/px)")
            ax.set_ylabel(f"y  ({nm_per_px:.3g} nm/px)")
            from gui_app._calib_utils import add_real_scalebar
            add_real_scalebar(ax, nm_per_px, length_nm=100, color="white")
        ax.set_title(
            f"{self.sample} — class map (K_active = {K})  counts = {counts}\n"
            f"left-click → single pattern   |   "
            f"right-click → grain average   |   "
            f"shift+right-click → add grain to stack",
            fontsize=10)
        cb = fig.colorbar(im, ax=ax, fraction=0.025, pad=0.02,
                           ticks=list(range(K)))
        cb.set_label("class id")

        # Click handler: button==1 → single pattern; button==3 → grain
        # average (connected component of same-class pixels). Both also
        # populate the single-frame idx field so the triptych runs on
        # the same picked pattern.
        def _on_click(event):
            if event.button not in (1, 3) or event.inaxes is not ax:
                return
            if event.xdata is None or event.ydata is None:
                return
            x = int(round(event.xdata)); y = int(round(event.ydata))
            x = max(0, min(Nx - 1, x)); y = max(0, min(Ny - 1, y))
            idx = y * Nx + x
            try:
                self._frame_idx_var.set(str(idx))
                self._attr_source.set("frame")
            except Exception:
                pass
            shift = self._event_shift(event)
            if event.button == 3 and shift:
                # Shift+right-click accumulates grains into a stacked
                # comparison window (one grain per row).
                self._add_grain_to_stack(y, x)
            elif event.button == 1:
                self._show_pattern_popup(y, x)
            elif event.button == 3:
                self._show_grain_popup(y, x)
        self._click_cid = self._canvas.mpl_connect("button_press_event",
                                                     _on_click)
        self._redraw()
        self._set_status(
            "class map rendered (left=pattern, right=grain average).")

    def _show_pattern_popup(self, y, x):
        """Single-pattern popup at scan position (y, x).

        Left  : the raw diffraction pattern (vmax / log controls).
        Right : on-demand GradCAM for the assigned class (target = the
                class the model picked here).  Computed only when the
                user clicks 'Compute GradCAM' so opening the popup
                stays cheap.

        Mirrors `_show_grain_popup` minus the class-map panel.
        """
        from matplotlib.colors import ListedColormap  # noqa: F401 (parallel to grain popup)
        cfg = SAMPLES[self.sample]
        ds = LoadPRZ(cfg["path"], resize=192, vmax=cfg["vmax"])
        Nx = self._scan_shape[1]
        idx = y * Nx + x
        raw = ds.get_raw(int(idx)).astype(np.float32)
        cls = int(self._inf["assigns"][idx])
        prob = float(self._inf["soft_probs"][idx, cls])
        train_vmax = float(cfg["vmax"])

        # ---- popup ----
        win = tk.Toplevel(self)
        win.title(f"pattern (y={y}, x={x})  →  class p{cls}  "
                    f"p={prob:.3f}")
        win.geometry("980x600")

        # Display-control row.
        ctrl = ctk.CTkFrame(win, fg_color="transparent")
        ctrl.pack(side="top", fill="x", padx=6, pady=4)
        log_var = ctk.BooleanVar(value=False)
        vmax_var = ctk.DoubleVar(value=train_vmax)
        ctk.CTkLabel(ctrl, text="display vmax:").pack(side="left",
                                                          padx=(2, 2))
        vmax_entry = ctk.CTkEntry(ctrl, textvariable=vmax_var,
                                      width=70)
        vmax_entry.pack(side="left", padx=2)
        ctk.CTkButton(ctrl, text="reset", width=60,
                       command=lambda: (vmax_var.set(train_vmax),
                                          _redraw())
                       ).pack(side="left", padx=2)
        ctk.CTkCheckBox(ctrl, text="log stretch",
                          variable=log_var,
                          command=lambda: _redraw()
                          ).pack(side="left", padx=8)
        gradcam_btn = ctk.CTkButton(ctrl, text="Compute GradCAM",
                                       width=160,
                                       fg_color=("#2D7A2D", "#1F7A1F"),
                                       command=lambda: _kickoff_gradcam())
        gradcam_btn.pack(side="left", padx=8)

        f = Figure(figsize=(9.6, 5.0), dpi=110, facecolor="white")
        ax_pat = f.add_subplot(1, 2, 1)
        ax_cam = f.add_subplot(1, 2, 2)
        rp = self._recip_per_px()

        # ---- RIGHT: GradCAM (on-demand) ----
        gc_state = {"avg_cart": None, "cam_cart": None,
                      "computing": False}

        def _draw_cam():
            ax_cam.clear()
            if gc_state["computing"]:
                ax_cam.text(0.5, 0.5,
                              "computing GradCAM …  (~5–10 s)",
                              ha="center", va="center", fontsize=10)
                ax_cam.set_axis_off()
                return
            if (gc_state["avg_cart"] is None
                    or gc_state["cam_cart"] is None):
                ax_cam.text(0.5, 0.5,
                              "(click 'Compute GradCAM' above)",
                              ha="center", va="center", fontsize=10,
                              color="#666")
                ax_cam.set_axis_off()
                return
            cam = gc_state["cam_cart"]        # 192×192 cart-cropped
            avg = gc_state["avg_cart"]        # 192×192 cart-cropped
            try:
                from data import SAMPLES as _S
                bm_r_192 = int((_S.get(self.sample) or {}).get(
                    "center_mask_radius", 0))
            except Exception:
                bm_r_192 = 0
            H = avg.shape[0]
            if bm_r_192 > 0:
                yy, xx = np.ogrid[:H, :H]; cy = cx = H / 2.0
                bm = ((yy - cy) ** 2 + (xx - cx) ** 2) > bm_r_192 ** 2
            else:
                bm = np.ones((H, H), dtype=bool)
            disp = avg * bm
            cn = (cam - cam.min()) * bm
            mx = float(cn.max())
            cn = (cn / mx) if mx > 0 else cn
            ax_cam.imshow(disp, cmap="gray",
                            aspect="equal", interpolation="nearest")
            ax_cam.imshow(cn, cmap="jet", alpha=0.55,
                            aspect="equal", interpolation="nearest")
            ax_cam.set_title(
                f"GradCAM  (target = p{cls},  model view)",
                fontsize=10)
            ax_cam.set_xticks([]); ax_cam.set_yticks([])
            try:
                if rp > 0:
                    from gui_app._calib_utils import (
                        q_per_polar_bin, get_raw_detector_size,
                        add_recip_scalebar)
                    qpx = q_per_polar_bin(
                        rp, get_raw_detector_size(self.sample))
                    add_recip_scalebar(ax_cam, q_per_disp_px=qpx,
                                          length_q=0.2)
            except Exception:
                pass

        def _kickoff_gradcam():
            if gc_state["computing"]:
                return
            ckpt = self._best_ckpt()
            if ckpt is None:
                messagebox.showinfo("GradCAM",
                    "No checkpoint available."); return
            gc_state["computing"] = True
            gradcam_btn.configure(state="disabled")
            _draw_cam()
            c.draw_idle()
            import threading
            def _worker():
                try:
                    avg_cart, cam_cart = self._compute_gradcam_from_raw(
                        ckpt, cls, raw)
                    def _done():
                        gc_state["avg_cart"] = avg_cart
                        gc_state["cam_cart"] = cam_cart
                        gc_state["computing"] = False
                        gradcam_btn.configure(state="normal")
                        _draw_cam()
                        c.draw_idle()
                    self.after(0, _done)
                except Exception as e:
                    err = repr(e)
                    def _fail():
                        gc_state["computing"] = False
                        gradcam_btn.configure(state="normal")
                        _draw_cam()
                        c.draw_idle()
                        messagebox.showerror("GradCAM", err)
                    self.after(0, _fail)
            threading.Thread(target=_worker, daemon=True).start()

        def _redraw():
            ax_pat.clear()
            try:
                vm = max(float(vmax_var.get()), 1e-6)
            except Exception:
                vm = train_vmax
            img = np.clip(raw / vm, 0.0, 1.0)
            if log_var.get():
                img = np.log1p(img * 50)
            ax_pat.imshow(img, cmap="inferno",
                            aspect="equal", interpolation="nearest")
            if rp > 0:
                ax_pat.set_xlabel(f"k_x ({rp:.3g} nm⁻¹/px)")
                ax_pat.set_ylabel(f"k_y ({rp:.3g} nm⁻¹/px)")
                from gui_app._calib_utils import add_recip_scalebar
                add_recip_scalebar(ax_pat, q_per_disp_px=rp,
                                      length_q=0.2)
            stretch_tag = "  log1p×50" if log_var.get() else ""
            tag_vmax = ("vmax=train"
                         if abs(float(vmax_var.get()) - train_vmax) < 1e-9
                         else f"vmax={float(vmax_var.get()):.3g}")
            ax_pat.set_title(
                f"pattern (y={y}, x={x})  →  p{cls}  "
                f"p={prob:.3f}  [{tag_vmax}{stretch_tag}]",
                fontsize=10)
            _draw_cam()
            c.draw_idle()

        c = FigureCanvasTkAgg(f, master=win)
        c.get_tk_widget().pack(fill="both", expand=True)
        tb = NavigationToolbar2Tk(c, win, pack_toolbar=False)
        tb.update(); tb.pack(side="bottom", fill="x")
        vmax_entry.bind("<Return>", lambda _e: _redraw())
        _redraw()

    def _send_to_acom_tab(self, grain_avg, y, x, cls, n_pix, mean_conf):
        """Push a grain average to the dedicated ACOM tab as the test
        source, run detection + 1D-radial preview, then switch tabs."""
        try:
            ac = getattr(self.app, "acom2", None)
            if ac is None:
                messagebox.showinfo("ACOM",
                    "ACOM tab not initialised."); return
            ac._test_pattern = np.asarray(grain_avg, dtype=np.float32)
            H, W = ac._test_pattern.shape
            ac._test_center = (H / 2.0, W / 2.0)
            ac._test_origin = (f"grain @ ({y}, {x})  class p{cls}  "
                                f"{n_pix}px ⟨p⟩={mean_conf:.3f}")
            ac._source_status.configure(text=f"loaded: {ac._test_origin}")
            ac._source_var.set("grain")
            ac._detect_and_redraw()
            ac.refresh_from_posthoc()
            # Switch the top-level tabview to the ACOM tab.
            try:
                self.app._tabs.set("ACOM")
            except Exception:
                pass
        except Exception as e:
            messagebox.showerror("Send to ACOM", repr(e))

    # ------------------------------------------------------------------
    # ACOM (CIF → crystal → orientation match) — single + batch
    # ------------------------------------------------------------------
    def _acom_browse_cif(self):
        from tkinter import filedialog
        p = filedialog.askopenfilename(
            filetypes=[("CIF", "*.cif"), ("All", "*.*")])
        if p:
            self._acom_cif_var.set(p)
            self._acom_state.configure(
                text=f"CIF set: {os.path.basename(p)}  (build the crystal)")

    def _acom_sync_calib_from_topbar(self):
        """nm⁻¹/px (topbar) → 1/Å/px (ACOM, ×0.1)."""
        rp = self._recip_per_px()
        if rp > 0:
            self._acom_inv_ang.set(round(rp * 0.1, 6))
            self._set_status(
                f"ACOM calibration synced: {rp:.5g} nm⁻¹/px → "
                f"{rp * 0.1:.5g} 1/Å/px")

    def _acom_get_crystal(self):
        """Return a cached or freshly-built Crystal.  Raises with a
        user-readable error if the CIF isn't set, py4DSTEM is missing,
        or the build/plan steps fail."""
        cif = self._acom_cif_var.get().strip()
        if not cif:
            raise RuntimeError("No CIF selected — Browse to pick one.")
        if not os.path.exists(cif):
            raise RuntimeError(f"CIF not found: {cif}")
        kmax = float(self._acom_kmax.get())
        plan = self._acom_plan_mode.get()
        key = (cif, kmax, plan)
        if self._acom_crystal is not None and self._acom_crystal_key == key:
            return self._acom_crystal
        from gui_app.acom_core import load_crystal, prepare_crystal
        cr = load_crystal(cif)
        prepare_crystal(cr, k_max=kmax, plan_mode=plan)
        self._acom_crystal = cr
        self._acom_crystal_key = key
        return cr

    def _acom_build_crystal(self):
        """User-facing build button: do the heavy CIF + plan work in a
        worker so the GUI stays alive."""
        def _w():
            try:
                self.after(0, lambda: self._acom_state.configure(
                    text="building crystal + orientation plan… "
                         "(~30 s)"))
                cr = self._acom_get_crystal()
                line = "crystal built."
                # Most py4DSTEM versions store unit cell in `cr.cell`
                # as a length-6 array (a, b, c, α, β, γ).  Older
                # versions expose `cr.a`, etc.
                try:
                    cell = getattr(cr, "cell", None)
                    if cell is not None and len(cell) == 6:
                        a, b, c, al, be, ga = [float(v) for v in cell]
                        line += (f"  a={a:.3f}  b={b:.3f}  c={c:.3f} Å"
                                  f"   α={al:.1f}  β={be:.1f}  γ={ga:.1f}°")
                    else:
                        a = float(getattr(cr, "a"))
                        b = float(getattr(cr, "b"))
                        c = float(getattr(cr, "c"))
                        line += f"  a={a:.3f}  b={b:.3f}  c={c:.3f} Å"
                except Exception:
                    pass
                sg = getattr(cr, "spaceGroupNumber", None)
                if sg is not None:
                    line += f"   SG #{sg}"
                self.after(0, lambda: self._acom_state.configure(text=line))
                self.after(0, lambda: self._set_status(line))
            except Exception as e:
                err = repr(e)
                self.after(0, lambda: self._acom_state.configure(
                    text=f"build failed: {err[:200]}"))
                self.after(0, lambda: messagebox.showerror(
                    "ACOM build", err))
        threading.Thread(target=_w, daemon=True).start()

    # ---- Multi-phase ACOM helpers ----------------------------------
    def _acom_phase_refresh_listbox(self):
        try:
            lb = self._acom_phase_list
            lb.delete(0, "end")
            for (n, p) in self._acom_phases:
                cached = " ✓" if n in self._acom_phase_crystals else ""
                lb.insert("end", f"{n}  ({os.path.basename(p)}){cached}")
        except Exception:
            pass

    def _acom_phase_add(self):
        from tkinter import filedialog, simpledialog
        p = filedialog.askopenfilename(
            filetypes=[("CIF", "*.cif"), ("All", "*.*")])
        if not p:
            return
        default = os.path.splitext(os.path.basename(p))[0]
        name = simpledialog.askstring(
            "Phase name",
            f"Short name for this phase (e.g. 'alpha', 'gamma'):",
            initialvalue=default, parent=self)
        if not name:
            return
        # Replace if same name exists.
        self._acom_phases = [(n, q) for (n, q) in self._acom_phases
                                  if n != name]
        self._acom_phases.append((name, p))
        self._acom_phase_refresh_listbox()
        self._acom_mp_status.configure(
            text=f"{len(self._acom_phases)} phase(s) listed. "
                  f"Click 'Build all' before running.")

    def _acom_phase_remove(self):
        try:
            sel = self._acom_phase_list.curselection()
            if not sel:
                return
            i = int(sel[0])
            n, _ = self._acom_phases.pop(i)
            self._acom_phase_crystals.pop(n, None)
            self._acom_phase_keys.pop(n, None)
            self._acom_phase_refresh_listbox()
        except Exception:
            pass

    def _acom_phase_build_all(self):
        """Build (or cache-hit) a Crystal+plan for every phase in the
        list.  Heavy: each is ~30 s.  Runs in a worker."""
        if not self._acom_phases:
            messagebox.showinfo("Multi-phase",
                "Add at least one CIF first."); return
        def _w():
            from gui_app.acom_core import load_crystal, prepare_crystal
            kmax = float(self._acom_kmax.get())
            plan = self._acom_plan_mode.get()
            t0 = time.time()
            for i, (name, cif) in enumerate(self._acom_phases):
                key = (cif, kmax, plan)
                if self._acom_phase_keys.get(name) == key \
                        and name in self._acom_phase_crystals:
                    continue
                try:
                    self.after(0,
                        lambda n=name, i=i: self._acom_mp_status.configure(
                            text=f"building [{i+1}/{len(self._acom_phases)}] "
                                  f"{n}…  ({time.time()-t0:.0f}s)"))
                    cr = load_crystal(cif)
                    prepare_crystal(cr, k_max=kmax, plan_mode=plan)
                    self._acom_phase_crystals[name] = cr
                    self._acom_phase_keys[name] = key
                except Exception as e:
                    err = repr(e)
                    self.after(0, lambda n=name, e=err:
                        messagebox.showerror("Multi-phase build",
                            f"{n}: {e}"))
            self.after(0, self._acom_phase_refresh_listbox)
            self.after(0,
                lambda: self._acom_mp_status.configure(
                    text=f"all {len(self._acom_phase_crystals)} crystals ready "
                          f"({time.time()-t0:.0f}s)."))
        threading.Thread(target=_w, daemon=True).start()

    def _acom_run_multiphase(self, mode="classes"):
        if not self._acom_phases:
            messagebox.showinfo("Multi-phase",
                "Add at least one CIF first."); return
        missing = [n for (n, _) in self._acom_phases
                    if n not in self._acom_phase_crystals]
        if missing:
            messagebox.showinfo("Multi-phase",
                f"Build crystals first.  Missing: {missing}")
            return
        if self._inf is None or self.sample is None:
            messagebox.showinfo("Multi-phase",
                "Load a run + render the class map first."); return
        threading.Thread(
            target=lambda: self._acom_multiphase_worker(mode),
            daemon=True).start()

    def _acom_multiphase_worker(self, mode):
        from scipy.ndimage import label
        from gui_app.acom_core import (acom_multiphase_full_dataset,
                                            acom_multiphase_batch,
                                            zone_axis_from_matrix)
        try:
            cfg = SAMPLES[self.sample]
            Ny, Nx = self._scan_shape
            assigns = self._inf["assigns"]
            soft = self._inf["soft_probs"]
            K = int(soft.shape[1])
            assigns_grid = assigns.reshape(Ny, Nx)
            inv_a = float(self._acom_inv_ang.get())
            thr = float(self._acom_mp_threshold.get())
            mar = float(self._acom_mp_margin.get())
            crystals = {n: self._acom_phase_crystals[n]
                         for (n, _) in self._acom_phases}
            phase_names = list(crystals.keys())

            t0 = time.time()
            if mode == "full_dataset":
                cube = _open_lazy(cfg["path"],
                                     scan_shape=self._scan_shape)
                stride = max(int(self._acom_full_stride.get()), 1)
                def _prog(done, total, stage):
                    if stage == "detect" and done % 256 == 0:
                        dt = time.time() - t0
                        eta = (dt / max(done, 1)) * (total - done)
                        self.after(0, lambda: self._set_status(
                            f"MP-ACOM detect {done}/{total}  "
                            f"({dt:.0f}s, ETA {eta:.0f}s)"))
                    elif stage.startswith("match"):
                        self.after(0, lambda s=stage: self._set_status(
                            f"MP-ACOM {s}…"))
                mp = acom_multiphase_full_dataset(
                    crystals, cube, inv_ang_per_pixel=inv_a,
                    subsample_stride=stride, progress_cb=_prog,
                    threshold=thr, margin=mar)
                self._acom_last_mp = mp
                dt = time.time() - t0
                self.after(0,
                    lambda: self._render_multiphase_map(mp, dt, stride))
                self.after(0, lambda: self._set_status(
                    f"MP-ACOM full dataset done ({dt:.0f}s)."))
                return

            # ----- batch modes (classes / grains) -----
            patterns: list = []
            labels:   list = []
            classes:  list = []
            if mode == "classes":
                avgs = self._compute_class_averages(top_n=256)
                vm = float(cfg.get("vmax", 5.0))
                for k in range(K):
                    p = (avgs[k] * vm).astype(np.float32)
                    patterns.append(p)
                    n = int((assigns == k).sum())
                    labels.append(f"class p{k}  N={n}")
                    classes.append(k)
            elif mode == "grains":
                for k in range(K):
                    mask = (assigns_grid == k)
                    if not mask.any():
                        continue
                    lab, _n = label(mask)
                    if _n == 0:
                        continue
                    sizes = np.bincount(lab.ravel()); sizes[0] = 0
                    gid = int(np.argmax(sizes))
                    if gid == 0:
                        continue
                    ys, xs = np.where(lab == gid)
                    yi = int(ys[len(ys)//2])
                    xi = int(xs[len(xs)//2])
                    gi = self._compute_grain_average(yi, xi)
                    if gi is None:
                        continue
                    patterns.append(gi["grain_avg"].astype(np.float32))
                    labels.append(
                        f"class p{k}  grain={gi['n_pix']}px  "
                        f"⟨p⟩={gi['mean_conf']:.2f}")
                    classes.append(k)
            else:
                raise ValueError(f"unknown MP mode: {mode!r}")
            if not patterns:
                self.after(0, lambda: messagebox.showinfo(
                    "Multi-phase", "No patterns to run on."))
                return
            mp = acom_multiphase_batch(crystals, patterns,
                                            inv_ang_per_pixel=inv_a,
                                            threshold=thr, margin=mar)
            self._acom_last_mp = mp
            dt = time.time() - t0
            self.after(0, lambda:
                self._render_multiphase_cards(mp, patterns, labels,
                                                  classes, mode, dt))
            self.after(0, lambda: self._set_status(
                f"MP-ACOM {mode} done ({dt:.0f}s)."))
        except Exception as e:
            err = repr(e)
            self.after(0, lambda: messagebox.showerror(
                "Multi-phase ACOM", err))
            self.after(0, lambda: self._set_status(
                f"MP-ACOM failed: {err[:120]}"))

    # ---- Multi-phase rendering -------------------------------------
    def _phase_palette(self, n_phases):
        """Pick a distinct, readable colour per phase.  Phase 0 is blue,
        phase 1 is orange-red (these are the dominant α/γ scenario)."""
        import matplotlib.pyplot as plt
        base = ["#1f77b4", "#d62728", "#2ca02c", "#ff7f0e",
                  "#9467bd", "#8c564b", "#e377c2", "#7f7f7f"]
        return [base[i % len(base)] for i in range(n_phases)]

    def _render_multiphase_map(self, mp, elapsed_s, stride):
        """Full-dataset MP render: 4-panel = phase mask + winning corr +
        per-phase orientation + corr-diff."""
        from matplotlib.colors import ListedColormap
        from gui_app.acom_core import zone_axis_from_matrix
        Ny, Nx = mp["scan_shape"]
        names = mp["phase_names"]
        n_ph = len(names)
        phase_id = mp["phase_id"]
        corr_win = mp["winning_corr"]
        corr_per = mp["corr_per_phase"]
        rmat_win = mp["winning_rmat"]

        # Ternary phase mask: -2 = ambiguous (grey), -1 = neither
        # (black), 0..n-1 = phase colour.
        palette = self._phase_palette(n_ph)
        # Build a discrete RGB image of the phase mask.
        from matplotlib.colors import to_rgb
        phase_rgb = np.zeros((Ny, Nx, 3), dtype=np.float32)
        # ambiguous → mid grey
        phase_rgb[phase_id == -2] = (0.45, 0.45, 0.45)
        # neither stays black
        for pi in range(n_ph):
            m = (phase_id == pi)
            if m.any():
                phase_rgb[m] = to_rgb(palette[pi])

        # Per-phase orientation: combine ZA-RGB only where that phase wins.
        za_rgb = np.zeros((Ny, Nx, 3), dtype=np.float32)
        za_map = {}      # (phase_idx, (u,v,w)) → colour
        import matplotlib.pyplot as plt
        cmap_or = plt.get_cmap("tab20")
        for rx in range(Ny):
            for ry in range(Nx):
                pi = int(phase_id[rx, ry])
                if pi < 0:
                    continue
                R = rmat_win[rx, ry]
                if not np.isfinite(R).all():
                    continue
                za, _ = zone_axis_from_matrix(R)
                key = (pi, tuple(za))
                if key not in za_map:
                    za_map[key] = cmap_or(len(za_map) % 20)[:3]
                za_rgb[rx, ry] = za_map[key]

        # Corr-diff (only meaningful when ≥2 phases): top1 - top2.
        if n_ph >= 2:
            cp_sorted = np.sort(-corr_per, axis=0)
            top1 = -cp_sorted[0]
            top2 = -cp_sorted[1]
            corr_diff = (top1 - top2).astype(np.float32)
        else:
            corr_diff = np.zeros((Ny, Nx), dtype=np.float32)

        self._fig.clf()
        self._fig.set_size_inches(13, 11)
        gs = self._fig.add_gridspec(2, 2, hspace=0.20, wspace=0.10)
        ax_mask = self._fig.add_subplot(gs[0, 0])
        ax_corr = self._fig.add_subplot(gs[0, 1])
        ax_za   = self._fig.add_subplot(gs[1, 0])
        ax_diff = self._fig.add_subplot(gs[1, 1])

        ax_mask.imshow(phase_rgb, interpolation="nearest")
        # phase legend
        from matplotlib.patches import Patch
        leg_handles = []
        for pi, n in enumerate(names):
            frac = float((phase_id == pi).sum()) / max(Ny * Nx, 1)
            leg_handles.append(Patch(color=palette[pi],
                                        label=f"{n}  {frac*100:.1f}%"))
        frac_neither = float((phase_id == -1).sum()) / max(Ny * Nx, 1)
        if frac_neither > 0:
            leg_handles.append(Patch(color="black",
                                        label=f"neither  {frac_neither*100:.1f}%"))
        frac_amb = float((phase_id == -2).sum()) / max(Ny * Nx, 1)
        if frac_amb > 0:
            leg_handles.append(Patch(color="#727272",
                                        label=f"ambiguous  {frac_amb*100:.1f}%"))
        ax_mask.legend(handles=leg_handles, loc="lower right",
                          fontsize=8, framealpha=0.75)
        ax_mask.set_title(f"phase mask  thr={self._acom_mp_threshold.get():.3f}  "
                          f"margin={self._acom_mp_margin.get():.3f}",
                          fontsize=11)
        ax_mask.set_xticks([]); ax_mask.set_yticks([])

        im_c = ax_corr.imshow(corr_win, cmap="viridis",
                                interpolation="nearest")
        ax_corr.set_title("winning correlation (top-1)", fontsize=11)
        ax_corr.set_xticks([]); ax_corr.set_yticks([])
        self._fig.colorbar(im_c, ax=ax_corr, fraction=0.045, pad=0.02)

        ax_za.imshow(za_rgb, interpolation="nearest")
        ax_za.set_title("winning zone-axis  (phase × ZA hash)",
                          fontsize=11)
        ax_za.set_xticks([]); ax_za.set_yticks([])

        im_d = ax_diff.imshow(corr_diff, cmap="magma",
                                interpolation="nearest")
        ax_diff.set_title("corr(top1) − corr(top2)  (decisiveness)",
                            fontsize=11)
        ax_diff.set_xticks([]); ax_diff.set_yticks([])
        self._fig.colorbar(im_d, ax=ax_diff, fraction=0.045, pad=0.02)

        nm_per_px = self._real_per_px()
        if nm_per_px > 0:
            from gui_app._calib_utils import add_real_scalebar
            for a in (ax_mask, ax_corr, ax_za, ax_diff):
                add_real_scalebar(a, nm_per_px, length_nm=100,
                                       color="white")

        self._fig.suptitle(
            f"Multi-phase ACOM — {self.sample}  "
            f"phases={names}  stride={stride}  "
            f"calib={float(self._acom_inv_ang.get()):.5g} 1/Å/px  "
            f"({elapsed_s:.0f}s)",
            fontsize=11)
        self._fig.tight_layout(rect=[0, 0, 1, 0.96])
        self._canvas.draw_idle()

    def _render_multiphase_cards(self, mp, patterns, labels, classes,
                                       mode, elapsed_s):
        """Class/grain card grid for multiphase: each card shows the
        pattern, its detected peaks, plus a label `phase: ZA  corr=...`
        showing what was assigned."""
        from gui_app.acom_core import zone_axis_from_matrix
        names = mp["phase_names"]
        phase_id = mp["phase_id"]                 # (N,)
        corr_per = mp["corr_per_phase"]           # (P, N)
        rmat_per = mp["rmat_per_phase"]           # (P, N, 3, 3)
        N = len(patterns)
        cols = 4
        rows = (N + cols - 1) // cols
        self._fig.clf()
        self._fig.set_size_inches(2.7 * cols, 3.0 * rows + 0.5)
        gs = self._fig.add_gridspec(rows, cols)
        palette = self._phase_palette(len(names))
        for k in range(N):
            ax = self._fig.add_subplot(gs[k // cols, k % cols])
            p = patterns[k]
            ax.imshow(np.log1p(np.clip(p, 0, None)),
                       cmap="inferno", aspect="equal",
                       interpolation="nearest")
            ax.set_xticks([]); ax.set_yticks([])
            pi = int(phase_id[k])
            if pi < 0:
                # show all-phase corrs in subtitle
                corr_str = "  ".join(
                    f"{names[p]}={corr_per[p, k]:.3f}"
                    for p in range(len(names)))
                tag = ("neither" if pi == -1 else "ambiguous")
                ax.set_title(f"{tag}   ({corr_str})\n{labels[k]}",
                              fontsize=8, color="#888")
            else:
                R = rmat_per[pi, k]
                za, mis = zone_axis_from_matrix(R)
                corr_win = corr_per[pi, k]
                # second-best phase + corr
                others = [(names[p], corr_per[p, k])
                           for p in range(len(names)) if p != pi]
                if others:
                    others_str = "  ".join(
                        f"{n}={c:.3f}" for (n, c) in others)
                else:
                    others_str = ""
                title_color = palette[pi]
                ax.set_title(
                    f"{names[pi]}  ZA=[{za[0]} {za[1]} {za[2]}]  "
                    f"corr={corr_win:.3f}\n"
                    f"{labels[k]}   {others_str}",
                    fontsize=8, color=title_color)
                # border-tint each card by phase
                for spine in ax.spines.values():
                    spine.set_edgecolor(title_color)
                    spine.set_linewidth(2.0)
        for k in range(N, rows * cols):
            ax = self._fig.add_subplot(gs[k // cols, k % cols])
            ax.set_axis_off()
        self._fig.suptitle(
            f"Multi-phase ACOM {mode} — {self.sample}  "
            f"phases={names}  thr={float(self._acom_mp_threshold.get()):.3f}  "
            f"margin={float(self._acom_mp_margin.get()):.3f}  "
            f"({elapsed_s:.0f}s)",
            fontsize=11)
        self._fig.tight_layout(rect=[0, 0, 1, 0.97])
        self._canvas.draw_idle()

    def _acom_run_on_pattern(self, pattern, title="ACOM",
                                 subtitle="", center=None):
        """Run ACOM on a single 2D pattern and pop a result window.

        Used by the grain popup ("Run ACOM on this grain") and by the
        ACOM-class-grid as a per-card drill-in.
        """
        try:
            cr = self._acom_get_crystal()
        except Exception as e:
            messagebox.showerror("ACOM", str(e)); return
        # Pop placeholder window immediately so the user sees feedback.
        win = tk.Toplevel(self)
        win.title(title)
        win.geometry("1100x560")
        info_lbl = ctk.CTkLabel(win, text="(running ACOM…)",
            font=("Consolas", 10), anchor="w", justify="left")
        info_lbl.pack(side="top", fill="x", padx=6, pady=4)
        fig = Figure(figsize=(11.0, 5.0), dpi=110, facecolor="white")
        c = FigureCanvasTkAgg(fig, master=win)
        c.get_tk_widget().pack(side="top", fill="both", expand=True)
        tb = NavigationToolbar2Tk(c, win, pack_toolbar=False)
        tb.update(); tb.pack(side="bottom", fill="x")
        ax_pat = fig.add_subplot(1, 2, 1)
        ax_fit = fig.add_subplot(1, 2, 2)
        for ax in (ax_pat, ax_fit):
            ax.text(0.5, 0.5, "(computing…)", ha="center", va="center",
                     fontsize=11, color="#888", transform=ax.transAxes)
            ax.set_xticks([]); ax.set_yticks([])

        def _worker():
            from gui_app.acom_core import (acom_single_pattern,
                                                zone_axis_from_matrix)
            try:
                inv_a = float(self._acom_inv_ang.get())
                res = acom_single_pattern(cr, pattern, center=center,
                                              inv_ang_per_pixel=inv_a)
                # Try to extract ZA from the orientation result.
                rmat = None
                ort = res["orientation"]
                for attr in ("matrix",):
                    v = getattr(ort, attr, None)
                    if v is not None:
                        arr = np.asarray(v)
                        if arr.ndim == 3:
                            rmat = arr[0]
                        elif arr.ndim == 2:
                            rmat = arr
                        break
                za, mis = zone_axis_from_matrix(rmat)

                def _draw():
                    # LEFT: pattern + detected peaks (red rings).
                    ax_pat.clear()
                    H, W = pattern.shape
                    cy, cx = (center if center is not None
                                else (H / 2.0, W / 2.0))
                    img = np.log1p(np.clip(pattern, 0, None))
                    ax_pat.imshow(img, cmap="inferno", aspect="equal",
                                    interpolation="nearest")
                    peaks = res["peaks"]
                    if peaks.size:
                        ax_pat.scatter(peaks[:, 1], peaks[:, 0],
                                          s=80, facecolors="none",
                                          edgecolors="cyan", linewidths=1.3,
                                          label=f"detected ({len(peaks)})")
                        ax_pat.legend(loc="lower right", fontsize=8,
                                         framealpha=0.6)
                    ax_pat.set_xticks([]); ax_pat.set_yticks([])
                    ax_pat.set_title(
                        f"experimental  ({len(peaks)} peaks)",
                        fontsize=10)
                    # RIGHT: experimental peaks (centered, 1/Å) vs fitted.
                    ax_fit.clear()
                    pl = res["calibrated_pl"]
                    try:
                        from py4DSTEM.process.diffraction import (
                            plot_diffraction_pattern)
                        plot_diffraction_pattern(
                            res["fit_pattern"], bragg_peaks_compare=pl,
                            scale_markers=1000,
                            scale_markers_compare=4e4,
                            min_marker_size=1, figsize=(5, 5),
                            input_fig_handle=(fig, [ax_fit]))
                    except Exception as e2:
                        ax_fit.text(0.5, 0.5,
                            f"fit plot failed:\n{e2!r}",
                            ha="center", va="center", fontsize=9,
                            color="#a33", transform=ax_fit.transAxes)
                    ax_fit.set_title(
                        f"matched ZA ≈ [{za[0]} {za[1]} {za[2]}]"
                        f"   miso = {mis:.2f}°",
                        fontsize=10)
                    info_lbl.configure(
                        text=f"{subtitle}    corr = {res['corr']:.4f}"
                              f"    ZA = [{za[0]} {za[1]} {za[2]}]"
                              f"  (miso {mis:.2f}°)    "
                              f"#peaks = {len(res['peaks'])}    "
                              f"calib = {inv_a:.5g} 1/Å/px")
                    fig.tight_layout(); c.draw_idle()
                self.after(0, _draw)
            except Exception as e:
                err = repr(e)
                self.after(0, lambda: info_lbl.configure(
                    text=f"ACOM failed: {err[:200]}"))
                self.after(0, lambda: messagebox.showerror("ACOM", err))
        threading.Thread(target=_worker, daemon=True).start()

    def _acom_run_batch(self, mode="classes"):
        """Run ACOM on a batch of patterns and render a card grid.

        mode = "classes"          → one card per class (class-avg pattern).
        mode = "grains"           → one card per class, using the largest
                                      grain in that class as the source.
        mode = "grains_in_class"  → one card per grain in the class
                                      specified by `self._acom_grain_cls`.
        """
        if self._inf is None or self.sample is None:
            messagebox.showinfo("ACOM",
                "Load a run + render the class map first."); return
        try:
            cr = self._acom_get_crystal()
        except Exception as e:
            messagebox.showerror("ACOM", str(e)); return
        threading.Thread(
            target=lambda: self._acom_batch_worker(mode),
            daemon=True).start()

    def _acom_batch_worker(self, mode):
        from scipy.ndimage import label
        from gui_app.acom_core import (acom_batch, acom_full_dataset,
                                            zone_axis_from_matrix)
        try:
            cr = self._acom_get_crystal()
            cfg = SAMPLES[self.sample]
            Ny, Nx = self._scan_shape
            assigns = self._inf["assigns"]
            soft = self._inf["soft_probs"]
            K = int(soft.shape[1])
            assigns_grid = assigns.reshape(Ny, Nx)
            inv_a = float(self._acom_inv_ang.get())

            # Full-dataset mode goes through a different code path: no
            # card grid, just a 2D orientation/correlation map.
            if mode == "full_dataset":
                cube = _open_lazy(cfg["path"],
                                     scan_shape=self._scan_shape)
                stride = max(int(self._acom_full_stride.get()), 1)
                t0 = time.time()
                def _prog(done, total, stage):
                    if stage == "detect" and done % 256 == 0:
                        dt = time.time() - t0
                        eta = (dt / max(done, 1)) * (total - done)
                        self.after(0, lambda: self._set_status(
                            f"ACOM full ({stride}× stride): "
                            f"detect {done}/{total}  "
                            f"({dt:.0f}s elapsed, ETA {eta:.0f}s)"))
                    elif stage in ("match", "build_vectors"):
                        self.after(0, lambda: self._set_status(
                            f"ACOM full: {stage}…"))
                omap, bv, scan_shape = acom_full_dataset(
                    cr, cube, inv_ang_per_pixel=inv_a,
                    subsample_stride=stride, progress_cb=_prog)
                dt = time.time() - t0
                self.after(0,
                    lambda: self._render_acom_full_map(omap, scan_shape,
                                                            stride, dt))
                return

            patterns: list = []
            labels:   list = []        # per-card subtitle ("class p3" etc.)
            classes:  list = []        # class id per card (for the title row)

            if mode == "classes":
                # Pull the K class-avg patterns from the existing helper
                # (already vmax-normalised, blur/log-applied if set in cfg).
                avgs = self._compute_class_averages(top_n=256)
                # _compute_class_averages returns 192×192 cart-cropped.
                # Re-scale back to vmax for blob detection sensitivity:
                vm = float(cfg.get("vmax", 5.0))
                for k in range(K):
                    p = (avgs[k] * vm).astype(np.float32)
                    patterns.append(p)
                    n = int((assigns == k).sum())
                    labels.append(f"class p{k}  N={n}")
                    classes.append(k)
            elif mode == "grains":
                # one largest grain per class.
                for k in range(K):
                    mask = (assigns_grid == k)
                    if not mask.any():
                        continue
                    lab, _n = label(mask)
                    if _n == 0:
                        continue
                    sizes = np.bincount(lab.ravel())
                    sizes[0] = 0
                    gid = int(np.argmax(sizes))
                    if gid == 0:
                        continue
                    ys, xs = np.where(lab == gid)
                    yi = int(ys[len(ys) // 2]); xi = int(xs[len(xs) // 2])
                    gi = self._compute_grain_average(yi, xi)
                    if gi is None:
                        continue
                    patterns.append(gi["grain_avg"].astype(np.float32))
                    labels.append(
                        f"class p{k}  grain={gi['n_pix']}px  "
                        f"⟨p⟩={gi['mean_conf']:.2f}")
                    classes.append(k)
            elif mode == "grains_in_class":
                try:
                    target_cls = int(self._acom_grain_cls.get())
                except Exception:
                    messagebox.showerror("ACOM",
                        "Enter a class id (integer) first."); return
                if not (0 <= target_cls < K):
                    messagebox.showerror("ACOM",
                        f"Class {target_cls} out of range [0..{K-1}]."); return
                mask = (assigns_grid == target_cls)
                lab, _n = label(mask)
                if _n == 0:
                    messagebox.showinfo("ACOM",
                        f"class p{target_cls} has no grains."); return
                sizes = np.bincount(lab.ravel())
                sizes[0] = 0
                # cap at 16 largest grains so the grid is readable.
                order = np.argsort(-sizes)
                picked = [int(g) for g in order if sizes[g] > 0][:16]
                for gid in picked:
                    ys, xs = np.where(lab == gid)
                    yi = int(ys[len(ys) // 2]); xi = int(xs[len(xs) // 2])
                    gi = self._compute_grain_average(yi, xi)
                    if gi is None:
                        continue
                    patterns.append(gi["grain_avg"].astype(np.float32))
                    labels.append(
                        f"grain {gid}  {gi['n_pix']}px  "
                        f"⟨p⟩={gi['mean_conf']:.2f}")
                    classes.append(target_cls)
            else:
                raise ValueError(f"unknown ACOM batch mode: {mode!r}")

            if not patterns:
                self.after(0, lambda: messagebox.showinfo(
                    "ACOM", "No patterns to run on."))
                return

            self.after(0, lambda: self._set_status(
                f"ACOM batch ({mode}): {len(patterns)} patterns, "
                f"detecting peaks…"))

            results, omap, bv = acom_batch(
                cr, patterns, inv_ang_per_pixel=inv_a)

            # Render the card grid on the main canvas.
            zas = []
            for r in results:
                za, mis = zone_axis_from_matrix(r["rotation_matrix"])
                zas.append((za, mis))
            self.after(0, lambda: self._render_acom_grid(
                patterns, labels, classes, results, zas, mode))
            self.after(0, lambda: self._set_status(
                f"ACOM batch ({mode}) done — {len(patterns)} matched."))
        except Exception as e:
            err = repr(e)
            self.after(0, lambda: messagebox.showerror("ACOM batch", err))
            self.after(0, lambda: self._set_status(
                f"ACOM batch failed: {err[:120]}"))

    def _render_acom_full_map(self, omap, scan_shape, stride, elapsed_s):
        """2-panel render for ACOM-on-full-dataset:

        LEFT  : correlation map (best correlation per scan position).
        RIGHT : zone-axis RGB (each pixel's top-1 ZA mapped to RGB by
                  hashing the integer (u,v,w) for visual distinctness;
                  identical ZAs share the same colour).

        omap is the OrientationMap returned by `match_orientations`.
        """
        from gui_app.acom_core import zone_axis_from_matrix
        Ny, Nx = scan_shape
        # ---- correlation map ----
        corr_field = np.full((Ny, Nx), np.nan, dtype=np.float32)
        cv = None
        for attr in ("corr", "correlation"):
            v = getattr(omap, attr, None)
            if v is not None:
                cv = np.asarray(v); break
        if cv is not None:
            try:
                # cv is typically (Ny, Nx, n_matches) — take top-1.
                if cv.ndim == 3:
                    corr_field[:] = cv[..., 0]
                elif cv.ndim == 2:
                    corr_field[:] = cv
            except Exception:
                pass

        # ---- ZA RGB ----
        mv = None
        for attr in ("matrix",):
            v = getattr(omap, attr, None)
            if v is not None:
                mv = np.asarray(v); break
        za_rgb = np.zeros((Ny, Nx, 3), dtype=np.float32)
        if mv is not None and mv.ndim >= 4:
            # (Ny, Nx, n_matches, 3, 3) — take top-1.
            top = mv[..., 0, :, :] if mv.ndim == 5 else mv
            # cheap ZA hash → colour: integer (u,v,w) → RGB via tab10.
            za_int_map = {}
            for rx in range(Ny):
                for ry in range(Nx):
                    R = top[rx, ry]
                    if not np.isfinite(R).all():
                        continue
                    za, _ = zone_axis_from_matrix(R)
                    # reduce mirror/sign degeneracy: make the first
                    # nonzero component positive
                    z = list(za)
                    for i in range(3):
                        if z[i] != 0:
                            if z[i] < 0:
                                z = [-x for x in z]
                            break
                    key = tuple(z)
                    if key not in za_int_map:
                        idx = len(za_int_map) % 10
                        import matplotlib.pyplot as plt
                        za_int_map[key] = plt.get_cmap("tab10")(idx)[:3]
                    za_rgb[rx, ry] = za_int_map[key]

        self._fig.clf()
        self._fig.set_size_inches(13, 6)
        gs = self._fig.add_gridspec(1, 2, wspace=0.12)
        ax1 = self._fig.add_subplot(gs[0, 0])
        ax2 = self._fig.add_subplot(gs[0, 1])
        im1 = ax1.imshow(corr_field, cmap="viridis",
                          interpolation="nearest")
        ax1.set_title(f"top-1 correlation  ({self.sample})", fontsize=11)
        ax1.set_xticks([]); ax1.set_yticks([])
        self._fig.colorbar(im1, ax=ax1, fraction=0.045, pad=0.02)
        ax2.imshow(za_rgb, interpolation="nearest")
        ax2.set_title("dominant zone axis (RGB-hashed)", fontsize=11)
        ax2.set_xticks([]); ax2.set_yticks([])
        nm_per_px = self._real_per_px()
        if nm_per_px > 0:
            from gui_app._calib_utils import add_real_scalebar
            add_real_scalebar(ax1, nm_per_px, length_nm=100,
                                color="white")
            add_real_scalebar(ax2, nm_per_px, length_nm=100,
                                color="white")
        self._fig.suptitle(
            f"ACOM full dataset — {self.sample}  "
            f"stride={stride}  "
            f"calib={float(self._acom_inv_ang.get()):.5g} 1/Å/px  "
            f"({elapsed_s:.0f}s)",
            fontsize=11)
        self._fig.tight_layout(rect=[0, 0, 1, 0.96])
        self._canvas.draw_idle()
        # Cache for downstream uses (strain, save).
        self._acom_last_omap = omap
        self._acom_last_corr_field = corr_field

    def _render_acom_grid(self, patterns, labels, classes, results,
                              zas, mode):
        """Render an N-card grid: each card shows the pattern with the
        detected peaks (cyan rings) + the predicted Bragg disks from the
        best ZA match (red rings), tagged with corr + ZA + miso."""
        N = len(patterns)
        cols = 4
        rows = (N + cols - 1) // cols
        self._fig.clf()
        self._fig.set_size_inches(2.6 * cols, 2.8 * rows + 0.5)
        gs = self._fig.add_gridspec(rows, cols)
        # Find a global vmax for log1p display so cards are comparable.
        for k in range(N):
            ax = self._fig.add_subplot(gs[k // cols, k % cols])
            p = patterns[k]
            ax.imshow(np.log1p(np.clip(p, 0, None)),
                       cmap="inferno", aspect="equal",
                       interpolation="nearest")
            r = results[k]
            peaks = r["peaks"]
            if peaks.size:
                ax.scatter(peaks[:, 1], peaks[:, 0],
                            s=22, facecolors="none",
                            edgecolors="cyan", linewidths=0.9)
            za, mis = zas[k]
            cls = classes[k]
            ax.set_xticks([]); ax.set_yticks([])
            ax.set_title(
                f"p{cls}  ZA=[{za[0]} {za[1]} {za[2]}]  "
                f"corr={r['corr']:.3f}\n{labels[k]}",
                fontsize=8)
        for k in range(N, rows * cols):
            ax = self._fig.add_subplot(gs[k // cols, k % cols])
            ax.set_axis_off()
        self._fig.suptitle(
            f"ACOM batch ({mode})  —  {self.sample}  "
            f"#patterns={N}  calib={float(self._acom_inv_ang.get()):.5g} 1/Å/px",
            fontsize=11)
        self._fig.tight_layout(rect=[0, 0, 1, 0.97])
        self._canvas.draw_idle()

    # ------------------------------------------------------------------
    # Multi-grain stacking (shift+right-click on the class map)
    # ------------------------------------------------------------------
    def _add_grain_to_stack(self, y, x):
        """Append the grain at (y, x) to the shift+right-click selection
        and (re)draw the stacked-comparison window: one grain per row,
        [class-map w/ grain highlighted | grain-average diffraction
        (+ GradCAM once computed)]."""
        # Immediate feedback — averaging a big grain can take a moment;
        # without this the GUI looks "stuck" while it reads frames.
        self._set_status(f"stacking grain @ (y={y}, x={x}) …")
        try:
            self.update_idletasks()
        except Exception:
            pass
        try:
            gi = self._compute_grain_average(y, x, apply_filters=True)
        except Exception as e:
            self._set_status(f"grain stack failed: {e!r}")
            return
        if gi is None:
            self._set_status("grain lookup failed: pixel not in any grain")
            return
        if not hasattr(self, "_grain_stack") or self._grain_stack is None:
            self._grain_stack = []
        # De-dupe: same grain (same class + overlapping mask at click) is
        # skipped so repeated clicks in one grain don't pile up rows.
        for rec in self._grain_stack:
            try:
                if (rec["cls"] == gi["cls"]
                        and rec["grain_mask"].shape == gi["grain_mask"].shape
                        and rec["grain_mask"][y, x]):
                    self._set_status("grain already in stack")
                    return
            except Exception:
                continue
        self._grain_stack.append(dict(y=y, x=x, cam=None, **gi))
        try:
            self._ensure_grain_stack_window()
            self._redraw_grain_stack()
        except Exception as e:
            self._set_status(f"grain stack render failed: {e!r}")
            return
        self._set_status(
            f"added grain p{gi['cls']} @ (y={y}, x={x}) — "
            f"{len(self._grain_stack)} grain(s) stacked")

    def _ensure_grain_stack_window(self):
        """Create the stacked-grains Toplevel (with vmax / log / GradCAM
        controls) once; reuse it on subsequent shift+right-clicks."""
        win = getattr(self, "_grain_stack_win", None)
        if win is not None and bool(win.winfo_exists()):
            return
        win = tk.Toplevel(self)
        win.title("stacked grains (shift+right-click selection)")
        win.geometry("1200x900")
        bar = ctk.CTkFrame(win, fg_color="transparent")
        bar.pack(side="top", fill="x", padx=6, pady=4)
        self._grain_stack_count = ctk.CTkLabel(bar, text="")
        self._grain_stack_count.pack(side="left", padx=6)

        train_vmax = float(SAMPLES[self.sample]["vmax"])
        self._gs_vmax_var = ctk.DoubleVar(value=train_vmax)
        self._gs_log_var = ctk.BooleanVar(value=False)
        ctk.CTkLabel(bar, text="vmax:").pack(side="left", padx=(12, 2))
        ent = ctk.CTkEntry(bar, textvariable=self._gs_vmax_var, width=70)
        ent.pack(side="left", padx=2)
        ent.bind("<Return>", lambda _e: self._redraw_grain_stack())
        ctk.CTkButton(bar, text="reset", width=56,
                      command=lambda: (self._gs_vmax_var.set(train_vmax),
                                       self._redraw_grain_stack())
                      ).pack(side="left", padx=2)
        ctk.CTkCheckBox(bar, text="log stretch", variable=self._gs_log_var,
                        command=self._redraw_grain_stack).pack(side="left",
                                                               padx=8)
        self._gs_cam_btn = ctk.CTkButton(
            bar, text="Compute GradCAM", width=150,
            fg_color=("#2D7A2D", "#1F7A1F"),
            command=self._compute_grain_stack_gradcam)
        self._gs_cam_btn.pack(side="left", padx=8)
        ctk.CTkButton(bar, text="Clear", width=64,
                      command=self._clear_grain_stack).pack(side="right",
                                                            padx=4)
        holder = ctk.CTkScrollableFrame(win, fg_color="transparent")
        holder.pack(side="top", fill="both", expand=True)
        self._grain_stack_win = win
        self._grain_stack_holder = holder
        self._gs_canvas = None

    def _redraw_grain_stack(self):
        """Render every stacked grain as a row into the scrollable
        holder, honouring the vmax / log controls and showing GradCAM
        columns for any grain that has one computed."""
        import matplotlib.pyplot as plt
        from matplotlib.colors import ListedColormap
        win = getattr(self, "_grain_stack_win", None)
        if win is None or not bool(win.winfo_exists()):
            return
        stack = getattr(self, "_grain_stack", []) or []
        n = len(stack)
        holder = self._grain_stack_holder
        for w in holder.winfo_children():
            w.destroy()
        try:
            self._grain_stack_count.configure(text=f"{n} grain(s) stacked")
        except Exception:
            pass
        if n == 0:
            return

        cfg = SAMPLES[self.sample]
        Ny, Nx = self._scan_shape
        assigns_grid = self._inf["assigns"].reshape(Ny, Nx)
        K = int(self._inf["soft_probs"].shape[1])
        try:
            vm = max(float(self._gs_vmax_var.get()), 1e-6)
        except Exception:
            vm = float(cfg["vmax"])
        log_on = bool(self._gs_log_var.get())
        has_cam = any(rec.get("cam") is not None for rec in stack)
        ncols = 3 if has_cam else 2

        cmap = ListedColormap(plt.cm.tab20(np.linspace(0, 1, max(K, 2))))
        # Larger rows: ~3.6" tall each, ~4.0" per column wide.
        fig = Figure(figsize=(4.0 * ncols, 3.6 * n), dpi=110,
                     facecolor="white")
        rp = self._recip_per_px() or 0.0
        pat_axes = []   # (ax, H, W) for hover-q wiring
        for r, rec in enumerate(stack):
            ax_map = fig.add_subplot(n, ncols, ncols * r + 1)
            ax_pat = fig.add_subplot(n, ncols, ncols * r + 2)
            ax_map.imshow(assigns_grid.astype(float), cmap=cmap,
                          vmin=0, vmax=max(K - 1, 1), interpolation="nearest")
            ax_map.imshow(np.where(rec["grain_mask"], 1.0, np.nan),
                          cmap=ListedColormap(["black"]),
                          interpolation="nearest", alpha=0.9)
            ax_map.set_title(f"p{rec['cls']} @ (y={rec['y']}, x={rec['x']})  "
                             f"{rec['n_pix']}px  conf={rec['mean_conf']:.2f}",
                             fontsize=9)
            ax_map.set_xticks([]); ax_map.set_yticks([])
            img = np.clip(rec["grain_avg"] / vm, 0.0, 1.0)
            if log_on:
                img = np.log1p(img * 50)
            ax_pat.imshow(img, cmap="inferno", interpolation="nearest")
            stag = "  log1p×50" if log_on else ""
            ax_pat.set_title(f"grain-avg diffraction [vmax={vm:.3g}{stag}]",
                             fontsize=9)
            ax_pat.set_xticks([]); ax_pat.set_yticks([])
            H, W = rec["grain_avg"].shape
            pat_axes.append((ax_pat, H, W))
            if has_cam:
                ax_cam = fig.add_subplot(n, ncols, ncols * r + 3)
                cam = rec.get("cam")
                if cam is not None:
                    avg_cart, cam_cart = cam
                    ax_cam.imshow(avg_cart, cmap="gray",
                                  interpolation="nearest")
                    ax_cam.imshow(cam_cart, cmap="jet", alpha=0.55,
                                  interpolation="nearest")
                    ax_cam.set_title("GradCAM", fontsize=9)
                else:
                    ax_cam.text(0.5, 0.5, "(pending)", ha="center",
                                va="center", fontsize=9)
                ax_cam.set_xticks([]); ax_cam.set_yticks([])
        fig.tight_layout()
        canvas = FigureCanvasTkAgg(fig, master=holder)
        canvas.draw()
        canvas.get_tk_widget().pack(side="top", fill="both", expand=True)
        self._gs_canvas = canvas
        # Hover-q readout (px / q / d-spacing) on each grain-avg image.
        # Grain averages are at raw-detector resolution, so the recip
        # calibration applies directly (q_per_disp_px = rp).
        if rp > 0:
            try:
                from gui_app._ui import attach_hover_q
                for ax_pat, H, W in pat_axes:
                    attach_hover_q(canvas, ax_pat,
                                   center=(H / 2.0, W / 2.0),
                                   q_per_disp_px=rp, units="nm⁻¹")
            except Exception:
                pass
        # Reliably surface the window: a CTkScrollableFrame sometimes
        # leaves the embedded matplotlib canvas blank until the toplevel
        # is mapped/raised, which reads as "stuck / didn't open".
        try:
            win.deiconify(); win.lift()
            win.update_idletasks()
        except Exception:
            pass

    def _compute_grain_stack_gradcam(self):
        """Compute GradCAM (+ avg cartesian) for every grain in the
        stack in a worker thread, then redraw with the GradCAM column."""
        stack = getattr(self, "_grain_stack", []) or []
        if not stack:
            return
        ckpt = self._best_ckpt()
        if ckpt is None:
            messagebox.showinfo("GradCAM", "No checkpoint available.")
            return
        try:
            self._gs_cam_btn.configure(state="disabled", text="computing…")
        except Exception:
            pass

        def _worker():
            for rec in stack:
                if rec.get("cam") is not None:
                    continue
                try:
                    avg_cart, cam_cart, _ig = (
                        self._compute_gradcam_and_ig_from_raw(
                            ckpt, int(rec["cls"]), rec["grain_avg"]))
                    rec["cam"] = (avg_cart, cam_cart)
                except Exception as e:
                    print(f"[stack-gradcam] grain p{rec['cls']} "
                          f"failed: {e!r}", flush=True)

            def _done():
                try:
                    self._gs_cam_btn.configure(state="normal",
                                               text="Compute GradCAM")
                except Exception:
                    pass
                self._redraw_grain_stack()
            self.after(0, _done)
        threading.Thread(target=_worker, daemon=True).start()

    def _clear_grain_stack(self):
        self._grain_stack = []
        win = getattr(self, "_grain_stack_win", None)
        if win is not None and bool(win.winfo_exists()):
            win.destroy()
        self._grain_stack_win = None
        self._set_status("grain stack cleared")

    def _compute_grain_average(self, y, x, apply_filters=True):
        """Return the (raw-detector) average diffraction pattern for the
        connected-component grain at scan pixel (y, x).

        Used by both the right-click grain popup and the batch ACOM
        section.  Returns dict::

            {grain_mask, grain_pix, cls, grain_avg, mean_conf, n_pix}

        or None if the pixel isn't in any grain (out-of-range / invalid
        click).
        """
        from scipy.ndimage import label
        cfg = SAMPLES[self.sample]
        ds = LoadPRZ(cfg["path"], resize=192, vmax=cfg["vmax"])
        Ny, Nx = self._scan_shape
        assigns = self._inf["assigns"]
        assigns_grid = assigns.reshape(Ny, Nx)
        if not (0 <= y < Ny and 0 <= x < Nx):
            return None
        cls = int(assigns_grid[y, x])
        mask = (assigns_grid == cls)
        lab, _n = label(mask)
        grain_id = int(lab[y, x])
        if grain_id == 0:
            return None
        grain_mask = (lab == grain_id)
        grain_pix = np.where(grain_mask.flatten())[0]
        n_pix = int(grain_pix.size)
        patterns = np.stack([ds.get_raw(int(i)) for i in grain_pix],
                              axis=0).astype(np.float32)
        grain_avg = patterns.mean(axis=0)
        if apply_filters:
            try:
                from data import apply_sample_filters
                v = float(cfg.get("vmax", 2.0))
                tmp = np.clip(grain_avg / max(v, 1e-6), 0.0, 1.0)
                tmp = apply_sample_filters(tmp, cfg)
                grain_avg = tmp * v
            except Exception:
                pass
        soft = self._inf["soft_probs"]
        mean_conf = float(soft[grain_pix, cls].mean())
        return dict(grain_mask=grain_mask, grain_pix=grain_pix,
                     cls=cls, grain_avg=grain_avg,
                     mean_conf=mean_conf, n_pix=n_pix)

    def _show_grain_popup(self, y, x):
        """Two-panel grain popup.

        Left panel  : the class map with this grain's pixels painted black
                       (so the grain location is unmistakable).
        Right panel : the average diffraction pattern across the grain
                       pixels, with the same vmax/log controls as the
                       single-pattern popup.

        Grain = connected component (4-connectivity) of same-class
        pixels containing the click.
        """
        from matplotlib.colors import ListedColormap
        import matplotlib.pyplot as plt
        cfg = SAMPLES[self.sample]
        Ny, Nx = self._scan_shape
        assigns = self._inf["assigns"]
        assigns_grid = assigns.reshape(Ny, Nx)
        K = int(self._inf["soft_probs"].shape[1])

        gi = self._compute_grain_average(y, x, apply_filters=True)
        if gi is None:
            self._set_status("grain lookup failed: pixel not in any grain")
            return
        grain_mask = gi["grain_mask"]
        grain_pix = gi["grain_pix"]
        cls = gi["cls"]
        grain_avg = gi["grain_avg"]
        mean_conf = gi["mean_conf"]
        n_pix = gi["n_pix"]
        train_vmax = float(cfg["vmax"])

        # ---- popup ----
        win = tk.Toplevel(self)
        win.title(f"grain @ (y={y}, x={x})  →  class p{cls}  "
                    f"({n_pix} pixels)")
        win.geometry("1720x560")

        # Display-control row at the top.
        ctrl = ctk.CTkFrame(win, fg_color="transparent")
        ctrl.pack(side="top", fill="x", padx=6, pady=4)
        log_var = ctk.BooleanVar(value=False)
        vmax_var = ctk.DoubleVar(value=train_vmax)
        ctk.CTkLabel(ctrl, text="display vmax:").pack(side="left", padx=(2, 2))
        vmax_entry = ctk.CTkEntry(ctrl, textvariable=vmax_var, width=70)
        vmax_entry.pack(side="left", padx=2)
        ctk.CTkButton(ctrl, text="reset",
                       width=60,
                       command=lambda: (vmax_var.set(train_vmax), _redraw())
                       ).pack(side="left", padx=2)
        ctk.CTkCheckBox(ctrl, text="log stretch",
                          variable=log_var,
                          command=lambda: _redraw()
                          ).pack(side="left", padx=8)
        gradcam_btn = ctk.CTkButton(ctrl, text="Compute GradCAM + IG",
                                       width=180,
                                       fg_color=("#2D7A2D", "#1F7A1F"),
                                       command=lambda: _kickoff_gradcam())
        gradcam_btn.pack(side="left", padx=8)
        # Send the grain average to the dedicated ACOM tab as the
        # current test source.  The user can then run detection +
        # calibration + match interactively over there (no black-box
        # batch).
        ctk.CTkButton(ctrl, text="Send grain to ACOM tab",
                       width=180,
                       fg_color=("#4D6FB0", "#3A5380"),
                       command=lambda: self._send_to_acom_tab(
                           grain_avg, y, x, cls, n_pix, mean_conf)
                       ).pack(side="left", padx=4)

        f = Figure(figsize=(17.0, 4.6), dpi=110, facecolor="white")
        ax_map = f.add_subplot(1, 4, 1)
        ax_pat = f.add_subplot(1, 4, 2)
        ax_cam = f.add_subplot(1, 4, 3)
        ax_ig  = f.add_subplot(1, 4, 4)

        # ---- LEFT: class map with grain in black ----
        cmap = plt.get_cmap("tab10")
        palette = ListedColormap([cmap(i) for i in range(K)])
        ax_map.imshow(assigns_grid, cmap=palette,
                       vmin=-0.5, vmax=K - 0.5,
                       interpolation="nearest")
        # Solid-black overlay where grain_mask is True.
        overlay = np.zeros((Ny, Nx, 4), dtype=np.float32)
        overlay[grain_mask, 3] = 1.0          # opaque, RGB stays 0 → black
        ax_map.imshow(overlay, interpolation="nearest")
        ax_map.set_title(f"class map  (grain p{cls}, {n_pix} px in black)",
                          fontsize=10)
        ax_map.set_axis_off()
        nm_per_px = self._real_per_px()
        if nm_per_px > 0:
            from gui_app._calib_utils import add_real_scalebar
            add_real_scalebar(ax_map, nm_per_px, length_nm=100,
                                color="white")

        # ---- MIDDLE: grain-average diffraction pattern ----
        rp = self._recip_per_px()
        # ---- RIGHT 2 panels: GradCAM + IG (computed together) ----
        gc_state = {"avg_cart": None, "cam_cart": None,
                      "ig_cart":  None, "computing": False}

        def _attr_underlay_and_mask(avg):
            """Common machinery for GradCAM/IG overlay panels.  Returns
            (display_avg, beam_mask) with the same beam-mask radius the
            model trained with."""
            try:
                from data import SAMPLES as _S
                bm_r_192 = int((_S.get(self.sample) or {}).get(
                    "center_mask_radius", 0))
            except Exception:
                bm_r_192 = 0
            H = avg.shape[0]
            if bm_r_192 > 0:
                yy, xx = np.ogrid[:H, :H]; cy = cx = H / 2.0
                bm = ((yy - cy) ** 2 + (xx - cx) ** 2) > bm_r_192 ** 2
            else:
                bm = np.ones((H, H), dtype=bool)
            return avg * bm, bm

        def _draw_attr_axis(ax, attr, label):
            """Draw a single attribution overlay panel."""
            ax.clear()
            if gc_state["computing"]:
                ax.text(0.5, 0.5,
                          f"computing {label} …  (~10–20 s)",
                          ha="center", va="center", fontsize=10,
                          transform=ax.transAxes)
                ax.set_axis_off()
                return
            if gc_state["avg_cart"] is None or attr is None:
                ax.text(0.5, 0.5,
                          "(click 'Compute GradCAM + IG' above)",
                          ha="center", va="center", fontsize=10,
                          color="#666", transform=ax.transAxes)
                ax.set_axis_off()
                return
            avg = gc_state["avg_cart"]
            disp, bm = _attr_underlay_and_mask(avg)
            cn = (attr - attr.min()) * bm
            mx = float(cn.max())
            cn = (cn / mx) if mx > 0 else cn
            ax.imshow(disp, cmap="gray", aspect="equal",
                        interpolation="nearest")
            ax.imshow(cn, cmap="jet", alpha=0.55, aspect="equal",
                        interpolation="nearest")
            ax.set_title(f"{label}  (target = p{cls})", fontsize=10)
            ax.set_xticks([]); ax.set_yticks([])
            try:
                if rp > 0:
                    from gui_app._calib_utils import (
                        q_per_polar_bin, get_raw_detector_size,
                        add_recip_scalebar)
                    qpx = q_per_polar_bin(
                        rp, get_raw_detector_size(self.sample))
                    add_recip_scalebar(ax, q_per_disp_px=qpx,
                                          length_q=0.2)
            except Exception:
                pass

        def _draw_cam():
            _draw_attr_axis(ax_cam, gc_state.get("cam_cart"), "GradCAM")
            _draw_attr_axis(ax_ig,  gc_state.get("ig_cart"),
                                "Integrated Gradients")

        def _kickoff_gradcam():
            if gc_state["computing"]:
                return
            ckpt = self._best_ckpt()
            if ckpt is None:
                messagebox.showinfo("GradCAM + IG",
                    "No checkpoint available."); return
            gc_state["computing"] = True
            gradcam_btn.configure(state="disabled")
            _draw_cam()
            c.draw_idle()
            import threading
            def _worker():
                try:
                    avg_cart, cam_cart, ig_cart = (
                        self._compute_gradcam_and_ig_from_raw(
                            ckpt, cls, grain_avg))
                    def _done():
                        gc_state["avg_cart"] = avg_cart
                        gc_state["cam_cart"] = cam_cart
                        gc_state["ig_cart"]  = ig_cart
                        gc_state["computing"] = False
                        gradcam_btn.configure(state="normal")
                        _draw_cam()
                        c.draw_idle()
                    self.after(0, _done)
                except Exception as e:
                    err = repr(e)
                    def _fail():
                        gc_state["computing"] = False
                        gradcam_btn.configure(state="normal")
                        _draw_cam()
                        c.draw_idle()
                        messagebox.showerror("GradCAM + IG", err)
                    self.after(0, _fail)
            threading.Thread(target=_worker, daemon=True).start()

        def _redraw():
            ax_pat.clear()
            try:
                vm = max(float(vmax_var.get()), 1e-6)
            except Exception:
                vm = train_vmax
            img = np.clip(grain_avg / vm, 0.0, 1.0)
            if log_var.get():
                img = np.log1p(img * 50)
            ax_pat.imshow(img, cmap="inferno",
                           aspect="equal", interpolation="nearest")
            if rp > 0:
                ax_pat.set_xlabel(f"k_x ({rp:.3g} nm⁻¹/px)")
                ax_pat.set_ylabel(f"k_y ({rp:.3g} nm⁻¹/px)")
                from gui_app._calib_utils import add_recip_scalebar
                # Grain avg is computed at raw-detector resolution.
                add_recip_scalebar(ax_pat, q_per_disp_px=rp, length_q=0.2)
            stretch_tag = "  log1p×50" if log_var.get() else ""
            tag_vmax = ("vmax=train"
                         if abs(float(vmax_var.get()) - train_vmax) < 1e-9
                         else f"vmax={float(vmax_var.get()):.3g}")
            ax_pat.set_title(
                f"grain avg ({n_pix} px, ⟨p_cls⟩={mean_conf:.3f})  "
                f"[{tag_vmax}{stretch_tag}]",
                fontsize=10)
            _draw_cam()
            c.draw_idle()

        c = FigureCanvasTkAgg(f, master=win)
        c.get_tk_widget().pack(fill="both", expand=True)
        tb = NavigationToolbar2Tk(c, win, pack_toolbar=False)
        tb.update(); tb.pack(side="bottom", fill="x")
        vmax_entry.bind("<Return>", lambda _e: _redraw())
        _redraw()
        # Hover-q on the grain pattern (raw-detector resolution).
        try:
            from gui_app._ui import attach_hover_q
            if rp > 0:
                H, W = grain_avg.shape
                attach_hover_q(c, ax_pat,
                                  center=(H / 2.0, W / 2.0),
                                  q_per_disp_px=rp, units="nm⁻¹")
        except Exception: pass

    def _diffraction_popup_window(self, title, subtitle, raw_2d, train_vmax):
        """Shared popup implementation for pattern / grain views.

        Defaults to vmax-normalized model-input view; checkbox for
        log1p stretch; entry for display-vmax override.
        """
        win = tk.Toplevel(self)
        win.title(title)
        win.geometry("520x600")

        # Display-control row at the top.
        ctrl = ctk.CTkFrame(win, fg_color="transparent")
        ctrl.pack(side="top", fill="x", padx=6, pady=4)
        log_var = ctk.BooleanVar(value=False)            # default: raw input
        vmax_var = ctk.DoubleVar(value=train_vmax)       # default: training vmax
        ctk.CTkLabel(ctrl, text="display vmax:").pack(side="left", padx=(2, 2))
        vmax_entry = ctk.CTkEntry(ctrl, textvariable=vmax_var, width=70)
        vmax_entry.pack(side="left", padx=2)
        ctk.CTkButton(ctrl, text="reset",
                       width=60,
                       command=lambda: (vmax_var.set(train_vmax), _redraw())
                       ).pack(side="left", padx=2)
        ctk.CTkCheckBox(ctrl, text="log stretch",
                          variable=log_var,
                          command=lambda: _redraw()
                          ).pack(side="left", padx=8)

        f = Figure(figsize=(5.0, 5.0), dpi=110, facecolor="white")
        ax = f.add_subplot(111)
        rp = self._recip_per_px()

        def _redraw():
            ax.clear()
            try:
                vm = max(float(vmax_var.get()), 1e-6)
            except Exception:
                vm = train_vmax
            img = np.clip(raw_2d / vm, 0.0, 1.0)
            if log_var.get():
                img = np.log1p(img * 50)
            ax.imshow(img, cmap="inferno",
                       aspect="equal", interpolation="nearest")
            if rp > 0:
                ax.set_xlabel(f"k_x ({rp:.3g} nm⁻¹/px)")
                ax.set_ylabel(f"k_y ({rp:.3g} nm⁻¹/px)")
                # Pattern popup shows the raw pattern (no resize/crop), so
                # 1 disp px = 1 raw-detector-px.
                from gui_app._calib_utils import add_recip_scalebar
                add_recip_scalebar(ax, q_per_disp_px=rp, length_q=0.2)
            stretch_tag = "  log1p×50" if log_var.get() else ""
            tag_vmax = ("vmax=train"
                         if abs(float(vmax_var.get()) - train_vmax) < 1e-9
                         else f"vmax={float(vmax_var.get()):.3g}")
            ax.set_title(f"{subtitle}  [{tag_vmax}{stretch_tag}]",
                          fontsize=10)
            c.draw_idle()

        c = FigureCanvasTkAgg(f, master=win)
        c.get_tk_widget().pack(fill="both", expand=True)
        tb = NavigationToolbar2Tk(c, win, pack_toolbar=False)
        tb.update(); tb.pack(side="bottom", fill="x")
        # Apply current vmax-entry value when user presses Return.
        vmax_entry.bind("<Return>", lambda _e: _redraw())
        _redraw()
        # Hover-q on the diffraction pattern (raw-detector resolution).
        try:
            from gui_app._ui import attach_hover_q
            if rp > 0:
                H, W = raw_2d.shape
                attach_hover_q(c, ax,
                                  center=(H / 2.0, W / 2.0),
                                  q_per_disp_px=rp, units="nm⁻¹")
        except Exception: pass

    def _refresh_class_dropdown_lazy(self):
        if self._inf is not None and self._class_menu.cget("values") in (["—"], []):
            self._refresh_class_dropdown()
            return True
        return False

    # raw class averages
    def _compute_class_averages(self, top_n=200, weighted=True):
        """Per-class average diffraction pattern.

        top_n   : number of highest-confidence members to average.
                   Pass None (or <=0) to use ALL members of each class
                   — the honest, unfiltered class average.
        weighted: weight each member by its soft-prob.  When using
                   all members, set weighted=False for a plain mean.
        """
        import torch
        import torch.nn.functional as F
        from torchvision.transforms import v2 as T
        from torchvision.transforms import InterpolationMode
        cfg = SAMPLES[self.sample]
        ds = LoadPRZ(cfg["path"], resize=192, vmax=cfg["vmax"])
        K = int(self._inf["soft_probs"].shape[1])
        soft = self._inf["soft_probs"]; ass = self._inf["assigns"]
        H = 192
        use_all = (top_n is None) or (top_n <= 0)
        cart_pre = T.Compose([
            T.CenterCrop(140),
            T.Resize(H, interpolation=InterpolationMode.BILINEAR, antialias=True),
        ])
        avgs = []
        for c in range(K):
            idx = np.where(ass == c)[0]
            if idx.size == 0:
                avgs.append(np.zeros((H, H), np.float32)); continue
            s = soft[idx, c]
            if use_all:
                sel = idx                       # everything in the class
                w = (s.astype(np.float32) if weighted
                       else np.ones(len(idx), np.float32))
            else:
                order = np.argsort(-s)[:min(top_n, len(idx))]
                sel = idx[order]
                w = (s[order].astype(np.float32) if weighted
                       else np.ones(len(sel), np.float32))
            patterns = np.stack([ds.get_raw(int(i)) for i in sel],
                                   0).astype(np.float32)
            wavg = (patterns * w[:, None, None]).sum(0) / (w.sum() + 1e-12)
            wn = np.clip(wavg / max(float(cfg["vmax"]), 1e-6), 0.0, 1.0)
            x = torch.from_numpy(wn).unsqueeze(0).unsqueeze(0).float()
            x = F.interpolate(x, size=(H, H), mode="bilinear",
                                align_corners=False)
            # Apply SAMPLES-stored blur σ + log stretch so the class
            # average matches what the model actually saw at 192.
            try:
                from data import apply_sample_filters
                arr = x.squeeze().cpu().numpy()
                arr = apply_sample_filters(arr, cfg)
                x = torch.from_numpy(arr).unsqueeze(0).unsqueeze(0).float()
            except Exception:
                pass
            avgs.append(cart_pre(x)[0, 0].cpu().numpy())
        return avgs

    def _disp_no_mask(self, avg):
        ref = avg.flatten()
        if not ref.size or ref.max() <= 0:
            return avg
        lo, hi = np.percentile(ref, 2), np.percentile(ref, 99.5)
        return np.log1p(np.clip(avg, lo, hi) - lo)

    # ============= Merge classes =====================================
    def _on_merge_method_change(self):
        """Show only the row relevant to the chosen merge method, and
        seed the threshold field with a sensible default for the metric
        the user just picked."""
        m = self._merge_method_var.get()
        try:
            if m == "manual":
                self._merge_thresh_row.pack_forget()
                self._merge_manual_row.pack(fill="x", padx=10, pady=2)
            else:
                self._merge_manual_row.pack_forget()
                self._merge_thresh_row.pack(fill="x", padx=10, pady=2)
                # Update default + hint per metric.  Only overwrite the
                # user's current value when it was left at the previous
                # method's default (so manual tweaks aren't clobbered).
                if m.startswith("SSIM"):
                    default = 0.70
                    hint = ("SSIM ∈ [-1, 1].  Class-avg patterns typically "
                             "land in 0.40–0.85.  Try 0.70 to merge only "
                             "near-identical avgs; 0.50 to be aggressive. "
                             "Hit 'Preview pairs' to see actual values "
                             "before applying.")
                else:  # cosine on centroids
                    default = 0.95
                    hint = ("cosine ∈ [-1, 1] on L2-normalised embedding "
                             "centroids.  Tight clusters often sit at "
                             "0.97–0.99.  Try 0.95 for safe merges; "
                             "0.90 to be aggressive. Hit 'Preview pairs' "
                             "to see actual values before applying.")
                # Replace only when current value is the other method's
                # default (or an obviously-untouched value).
                try:
                    cur = float(self._merge_thresh_var.get())
                except Exception:
                    cur = 0.85
                if abs(cur - 0.85) < 1e-6 or abs(cur - 0.70) < 1e-6 \
                       or abs(cur - 0.95) < 1e-6:
                    self._merge_thresh_var.set(default)
                self._merge_thresh_hint.configure(text=hint)
        except Exception:
            pass

    def _snapshot_inference(self):
        """Save a one-time copy of the original assigns / soft_probs so
        the user can hit Reset to undo merges."""
        if getattr(self, "_inf_orig", None) is not None:
            return
        if self._inf is None:
            return
        self._inf_orig = {
            "soft_probs": np.asarray(self._inf["soft_probs"]).copy(),
            "assigns":    np.asarray(self._inf["assigns"]).copy(),
            "embeds":     (np.asarray(self._inf["embeds"]).copy()
                            if self._inf.get("embeds") is not None
                            else None),
        }

    def _reset_merge(self):
        """Restore the snapshotted (pre-merge) inference and drop the
        cached merged state so the View toggle goes back to neutral."""
        if getattr(self, "_inf_orig", None) is None:
            self._merge_status.configure(text="nothing to reset.")
            return
        self._inf["soft_probs"] = self._inf_orig["soft_probs"].copy()
        self._inf["assigns"]    = self._inf_orig["assigns"].copy()
        if self._inf_orig.get("embeds") is not None:
            self._inf["embeds"] = self._inf_orig["embeds"].copy()
        self._inf_merged = None
        try: self._merge_view_var.set("merged")
        except Exception: pass
        self._reset_caches_post_merge()
        self._merge_status.configure(
            text=f"reset → K={int(self._inf['soft_probs'].shape[1])}.")
        try: self._render_classmap()
        except Exception: pass

    def _reset_caches_post_merge(self):
        # Class averages / BF / HAADF depend on assignments; invalidate
        # them so re-renders pick up the new K.
        try: self._class_avgs_cache = None
        except Exception: pass
        # Drop class dropdown so it refreshes.
        try:
            self._class_menu.configure(values=["—"])
            self._class_var.set("—")
        except Exception:
            pass

    def _apply_merge(self):
        """Dispatch to the chosen merge method, update self._inf in
        place, and re-render the class map."""
        if not self._ensure_inference():
            return
        self._snapshot_inference()
        method = self._merge_method_var.get()
        soft = np.asarray(self._inf["soft_probs"]).copy()
        ass = np.asarray(self._inf["assigns"]).copy()
        K = int(soft.shape[1])
        try:
            if method == "manual":
                pairs = self._merge_pairs_manual(K)
            elif method.startswith("SSIM"):
                thr = float(self._merge_thresh_var.get())
                pairs = self._merge_pairs_ssim(thr)
            else:  # cosine threshold
                thr = float(self._merge_thresh_var.get())
                pairs = self._merge_pairs_cosine(thr)
        except Exception as e:
            messagebox.showerror("Merge", repr(e)); return
        if not pairs:
            self._merge_status.configure(
                text="no merges proposed.")
            return
        # Union-find collapse so chains of merges resolve cleanly.
        parent = list(range(K))
        def _find(i):
            while parent[i] != i:
                parent[i] = parent[parent[i]]
                i = parent[i]
            return i
        def _union(a, b):
            ra, rb = _find(a), _find(b)
            if ra == rb: return
            # keep the smaller id as root (more-populous after re-rank).
            if ra > rb: ra, rb = rb, ra
            parent[rb] = ra
        for (a, b) in pairs:
            _union(int(a), int(b))
        groups = {}
        for i in range(K):
            groups.setdefault(_find(i), []).append(i)
        # Build new K' columns by summing source columns.
        old_to_new = {}
        new_K = 0
        new_cols = []
        for root in sorted(groups.keys()):
            members = groups[root]
            new_cols.append(soft[:, members].sum(axis=1))
            for m in members:
                old_to_new[m] = new_K
            new_K += 1
        new_soft = np.stack(new_cols, axis=1)
        new_ass = np.array([old_to_new[int(c)] for c in ass],
                              dtype=np.int32)
        # Renormalize soft (sums per row stay 1 modulo float noise).
        s = new_soft.sum(axis=1, keepdims=True).clip(min=1e-12)
        new_soft = new_soft / s
        # Re-rank by frequency (matches the eval_all convention).
        counts = np.bincount(new_ass, minlength=new_K)
        rank = np.argsort(-counts)
        remap = np.zeros(new_K, dtype=np.int32)
        for i, old in enumerate(rank):
            remap[old] = i
        new_ass = remap[new_ass]
        new_soft = new_soft[:, rank]

        self._inf["soft_probs"] = new_soft
        self._inf["assigns"]    = new_ass
        # Cache the merged state so the View toggle can swap back to it
        # after switching to "original".
        self._inf_merged = {
            "soft_probs": new_soft.copy(),
            "assigns":    new_ass.copy(),
            "embeds":     (np.asarray(self._inf["embeds"]).copy()
                              if self._inf.get("embeds") is not None
                              else None),
        }
        self._reset_caches_post_merge()

        # Report what happened.
        bits = []
        for root in sorted(groups.keys()):
            members = groups[root]
            if len(members) > 1:
                bits.append("{" + ",".join(map(str, members)) + "}")
        self._merge_status.configure(
            text=f"{method}: K {K} → {new_K}.  merged: "
                  + (" ".join(bits) or "(none)")
                  + "  |  view = merged")
        try: self._merge_view_var.set("merged")
        except Exception: pass
        try: self._render_classmap()
        except Exception: pass

    def _switch_merge_view(self):
        """Swap self._inf between the snapshotted original and the
        cached merged state, or render side-by-side."""
        if (getattr(self, "_inf_orig", None) is None
                or getattr(self, "_inf_merged", None) is None):
            self._merge_status.configure(
                text="apply a merge first.")
            return
        view = self._merge_view_var.get()
        if view == "original":
            src = self._inf_orig
            tag = "view = original"
        elif view.startswith("compare"):
            self._render_classmap_compare()
            self._merge_status.configure(
                text=(self._merge_status.cget("text") or "")
                     .split("  |  ")[0]
                     + "  |  view = compare")
            return
        else:
            src = self._inf_merged
            tag = "view = merged"
        self._inf["soft_probs"] = src["soft_probs"].copy()
        self._inf["assigns"]    = src["assigns"].copy()
        if src.get("embeds") is not None:
            self._inf["embeds"] = src["embeds"].copy()
        self._reset_caches_post_merge()
        old = (self._merge_status.cget("text") or "").split(
            "  |  ")[0]
        self._merge_status.configure(text=old + "  |  " + tag)
        try: self._render_classmap()
        except Exception: pass

    def _render_classmap_compare(self):
        """Side-by-side class maps: left = original (snapshot), right
        = merged. Both rendered with their own palettes so each map
        keeps the eval_all frequency-sort convention internally; class
        ids are not aligned across panels."""
        if (getattr(self, "_inf_orig", None) is None
                or getattr(self, "_inf_merged", None) is None):
            messagebox.showinfo("merge view",
                "apply a merge first."); return
        if self._scan_shape is None:
            messagebox.showinfo("merge view",
                "no scan_shape — load a run first."); return
        Ny, Nx = self._scan_shape
        a_orig = np.asarray(self._inf_orig["assigns"])
        a_merg = np.asarray(self._inf_merged["assigns"])
        K_orig = int(self._inf_orig["soft_probs"].shape[1])
        K_merg = int(self._inf_merged["soft_probs"].shape[1])
        fig = self._new_fig()
        ax1 = fig.add_subplot(1, 2, 1)
        ax2 = fig.add_subplot(1, 2, 2)
        for ax, ass, K_act, title in (
                (ax1, a_orig, K_orig,
                 f"Original  (K = {K_orig})"),
                (ax2, a_merg, K_merg,
                 f"Merged  (K = {K_merg})")):
            cmap = _adaptive_cmap(K_act)
            from matplotlib.colors import (BoundaryNorm,
                                              ListedColormap)
            palette = ListedColormap([cmap(i) for i in range(K_act)])
            norm = BoundaryNorm(np.arange(K_act + 1) - 0.5, K_act)
            grid = ass.reshape(Ny, Nx) if ass.size == Ny * Nx \
                else ass.reshape(1, -1)
            im = ax.imshow(grid, cmap=palette, norm=norm,
                              aspect="equal",
                              interpolation="nearest")
            ax.set_title(title, fontsize=11)
            ax.set_axis_off()
            cb = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04,
                                  ticks=list(range(K_act)))
            cb.set_label("class id", fontsize=8)
            try:
                nm_per_px = self._real_per_px()
            except Exception:
                nm_per_px = 0.0
            if nm_per_px > 0:
                from gui_app._calib_utils import add_real_scalebar
                add_real_scalebar(ax, nm_per_px, length_nm=100,
                                    color="white")
        fig.suptitle(f"{self.sample} — class map: original vs merged",
                       fontsize=12)
        self._redraw()

    def _merge_pairs_manual(self, K: int) -> list[tuple[int, int]]:
        x_str = self._merge_x_var.get().strip()
        y_str = self._merge_y_var.get().strip()
        if not x_str or not y_str:
            raise ValueError("Provide X and Y class ids.")
        x = int(x_str); y = int(y_str)
        if not (0 <= x < K and 0 <= y < K):
            raise ValueError(f"class ids must be in [0, {K - 1}].")
        if x == y:
            raise ValueError("X and Y are the same class.")
        return [(x, y)]

    def _preview_merge_pairs(self):
        """Compute the full pairwise similarity matrix for the chosen
        metric and pop a modal that lists every (i, j) pair and its
        score, sorted descending.  Pairs at-or-above the current
        threshold are highlighted so the user can dial the threshold
        before hitting Apply.

        Cheap to recompute on demand — the cosine path is O(K² · D) and
        the SSIM path is O(K² · 192²); for K ≤ 30 this is sub-second.
        """
        if not self._ensure_inference():
            return
        method = self._merge_method_var.get()
        try:
            thr = float(self._merge_thresh_var.get())
        except Exception:
            thr = 0.0
        soft = np.asarray(self._inf["soft_probs"])
        K = int(soft.shape[1])

        try:
            if method.startswith("SSIM"):
                from skimage.metrics import structural_similarity as ssim
                avgs = self._compute_class_averages()
                norm = []
                for a in avgs:
                    mn = float(np.min(a)); mx = float(np.max(a))
                    if mx - mn < 1e-9:
                        norm.append(np.zeros_like(a, dtype=np.float32))
                    else:
                        norm.append(((a - mn) / (mx - mn)
                                       ).astype(np.float32))
                S = np.full((K, K), -np.inf, dtype=np.float32)
                for i in range(K):
                    S[i, i] = 1.0
                    for j in range(i + 1, K):
                        try:
                            s = float(ssim(norm[i], norm[j],
                                              data_range=1.0))
                        except Exception:
                            s = float("nan")
                        S[i, j] = s; S[j, i] = s
                metric_name = "SSIM"
            else:  # cosine
                embeds = self._inf.get("embeds")
                if embeds is None:
                    messagebox.showerror(
                        "Preview",
                        "no embeddings cached — re-run inference first.")
                    return
                ass = np.asarray(self._inf["assigns"])
                cen = np.zeros((K, embeds.shape[1]), dtype=np.float32)
                for c in range(K):
                    m = ass == c
                    if m.any():
                        cen[c] = embeds[m].mean(axis=0)
                n = np.linalg.norm(cen, axis=1, keepdims=True).clip(
                    min=1e-9)
                cn = cen / n
                S = (cn @ cn.T).astype(np.float32)
                metric_name = "cosine"
        except Exception as e:
            messagebox.showerror("Preview", repr(e)); return

        # Build the ranked pair list (upper triangle only, i < j).
        rows = []
        for i in range(K):
            for j in range(i + 1, K):
                s = float(S[i, j])
                rows.append((s, i, j))
        rows.sort(key=lambda r: -r[0])

        # Modal listing
        win = ctk.CTkToplevel(self)
        win.title(f"Merge preview — {metric_name} (K={K})")
        win.geometry("520x520")
        try:
            win.transient(self.winfo_toplevel())
        except Exception:
            pass
        head = ctk.CTkLabel(win, justify="left", wraplength=480,
            font=("Segoe UI", 10),
            text=(f"Pairwise {metric_name} similarity between the K={K} "
                  f"current classes, sorted high → low.  Rows in green "
                  f"are at-or-above the current threshold "
                  f"({thr:g}) and would merge if you hit Apply."))
        head.pack(padx=10, pady=(8, 4), anchor="w")

        list_frame = ctk.CTkFrame(win)
        list_frame.pack(fill="both", expand=True, padx=10, pady=4)
        import tkinter as _tk
        lb = _tk.Listbox(list_frame, font=("Consolas", 10),
                           activestyle="dotbox",
                           selectmode=_tk.SINGLE)
        lb.pack(side="left", fill="both", expand=True, padx=(4, 0),
                 pady=4)
        sb = _tk.Scrollbar(list_frame, command=lb.yview)
        sb.pack(side="right", fill="y", padx=(0, 4), pady=4)
        lb.config(yscrollcommand=sb.set)

        for k, (s, i, j) in enumerate(rows):
            line = f"  {metric_name}={s:+.4f}    p{i:>2d}  ↔  p{j:<2d}"
            lb.insert(_tk.END, line)
            if s >= thr:
                lb.itemconfig(k, bg="#1F4D1F", fg="white")

        # Threshold editor at the bottom for live re-colouring.
        bot = ctk.CTkFrame(win, fg_color="transparent")
        bot.pack(fill="x", padx=10, pady=(0, 8))
        ctk.CTkLabel(bot, text="threshold:").pack(side="left")
        thr_local = ctk.DoubleVar(value=thr)
        e = ctk.CTkEntry(bot, textvariable=thr_local, width=70)
        e.pack(side="left", padx=4)

        def _recolor():
            try:
                t = float(thr_local.get())
            except Exception:
                return
            for k, (s, _i, _j) in enumerate(rows):
                if s >= t:
                    lb.itemconfig(k, bg="#1F4D1F", fg="white")
                else:
                    lb.itemconfig(k, bg="white", fg="black")

        e.bind("<Return>", lambda _e: _recolor())
        ctk.CTkButton(bot, text="Re-shade", width=90,
                        command=_recolor).pack(side="left", padx=2)

        def _adopt_and_close():
            try:
                self._merge_thresh_var.set(float(thr_local.get()))
            except Exception:
                pass
            win.destroy()

        ctk.CTkButton(bot, text="Use this threshold",
                        fg_color=("#2D7A2D", "#1F7A1F"), width=160,
                        command=_adopt_and_close
                        ).pack(side="right", padx=2)
        ctk.CTkButton(bot, text="Close", width=80,
                        command=win.destroy).pack(side="right", padx=2)

    def _merge_pairs_cosine(self, thr: float
                              ) -> list[tuple[int, int]]:
        """Pairs of classes whose centroid cosine > thr."""
        embeds = self._inf.get("embeds")
        if embeds is None:
            raise RuntimeError(
                "no embeddings cached — re-run inference first.")
        ass = np.asarray(self._inf["assigns"])
        K = int(self._inf["soft_probs"].shape[1])
        cen = np.zeros((K, embeds.shape[1]), dtype=np.float32)
        for c in range(K):
            m = ass == c
            if m.any():
                cen[c] = embeds[m].mean(axis=0)
        # Cosine similarity matrix.
        n = np.linalg.norm(cen, axis=1, keepdims=True).clip(min=1e-9)
        cn = cen / n
        S = cn @ cn.T
        pairs = []
        for i in range(K):
            for j in range(i + 1, K):
                if S[i, j] >= thr:
                    pairs.append((i, j))
        return pairs

    def _merge_pairs_ssim(self, thr: float
                            ) -> list[tuple[int, int]]:
        """Pairs of classes whose class-average SSIM > thr."""
        try:
            from skimage.metrics import structural_similarity as ssim
        except Exception as e:
            raise RuntimeError(
                f"scikit-image not available: {e!r}")
        avgs = self._compute_class_averages()
        K = len(avgs)
        # Normalize to [0, 1] per image so SSIM is comparable.
        norm = []
        for a in avgs:
            mn = float(np.min(a)); mx = float(np.max(a))
            if mx - mn < 1e-9:
                norm.append(np.zeros_like(a, dtype=np.float32))
            else:
                norm.append(((a - mn) / (mx - mn)).astype(np.float32))
        pairs = []
        for i in range(K):
            for j in range(i + 1, K):
                try:
                    s = float(ssim(norm[i], norm[j], data_range=1.0))
                except Exception:
                    s = 0.0
                if s >= thr:
                    pairs.append((i, j))
        return pairs

    def _render_class_distribution(self):
        """Two-panel class-occupancy plot:
          left  — soft occupancy: mean softmax probability per prototype
                   over the full scan (`p_bar`). Mirrors the notebook's
                   "Class Occupancy" bar chart and is what `effK =
                   exp(H(p_bar))` measures.
          right — hard occupancy: fraction of patterns whose argmax
                   lands on each prototype (= bincount / N). The two
                   columns rarely match exactly — gap = how confident
                   the model is at each prototype.
        """
        if not self._ensure_inference(): return
        soft = np.asarray(self._inf["soft_probs"])     # (N, K)
        assigns = np.asarray(self._inf["assigns"])     # (N,)
        N, K = soft.shape
        p_bar = soft.mean(axis=0)                       # (K,)
        # effK and active-class count (same definition the train log uses)
        pb = np.clip(p_bar, 1e-12, 1.0)
        effK = float(np.exp(-(pb * np.log(pb)).sum()))
        hard_counts = np.bincount(assigns, minlength=K)
        hard_frac = hard_counts / max(N, 1)
        active = int((hard_counts > 0).sum())

        cmap = _adaptive_cmap(K)
        colors = [cmap(c) for c in range(K)]

        fig = self._new_fig()
        ax_s = fig.add_subplot(1, 2, 1)
        x = np.arange(K)
        ax_s.bar(x, p_bar, color=colors, edgecolor="black",
                   linewidth=0.6)
        ax_s.set_xlabel("prototype id")
        ax_s.set_ylabel("mean softmax probability  (p_bar)")
        ax_s.set_title(
            f"Soft occupancy   (eff K = {effK:.2f})", fontsize=10)
        ax_s.set_xticks(x)
        ax_s.grid(axis="y", linestyle="--", alpha=0.4)
        # Annotate value on each bar.
        for k, v in enumerate(p_bar):
            ax_s.text(k, v + p_bar.max() * 0.02,
                        f"{v:.3f}", ha="center", va="bottom",
                        fontsize=8)
        ax_s.set_ylim(0, float(p_bar.max()) * 1.18)

        ax_h = fig.add_subplot(1, 2, 2)
        ax_h.bar(x, hard_frac, color=colors, edgecolor="black",
                   linewidth=0.6)
        ax_h.set_xlabel("prototype id")
        ax_h.set_ylabel("fraction of patterns  (argmax)")
        ax_h.set_title(
            f"Hard occupancy   (active = {active}/{K},  "
            f"N = {N})", fontsize=10)
        ax_h.set_xticks(x)
        ax_h.grid(axis="y", linestyle="--", alpha=0.4)
        for k, v in enumerate(hard_frac):
            ax_h.text(k, v + hard_frac.max() * 0.02,
                        f"{hard_counts[k]}",
                        ha="center", va="bottom",
                        fontsize=8)
        ax_h.set_ylim(0, float(hard_frac.max()) * 1.18)

        fig.suptitle(f"{self.sample} — class occupancy "
                       f"(K = {K}, N = {N})", fontsize=11)
        self._redraw()
        self._set_status(
            f"class occupancy: effK={effK:.2f}, active={active}/{K}.")

    def _render_class_avgs(self):
        if not self._ensure_inference(): return
        self._refresh_class_dropdown_lazy()
        K = int(self._inf["soft_probs"].shape[1])
        self._set_status("computing class averages…")
        avgs = self._compute_class_averages()
        counts = np.bincount(self._inf["assigns"], minlength=K).tolist()
        fig = self._new_fig()
        cols = min(K, 4)
        rows = (K + cols - 1) // cols
        # Reciprocal-space pixel scale shared across all panels.
        rp = self._recip_per_px()
        qpx = 0.0
        if rp > 0:
            from gui_app._calib_utils import (q_per_polar_bin,
                                                 get_raw_detector_size,
                                                 add_recip_scalebar)
            qpx = q_per_polar_bin(rp, get_raw_detector_size(self.sample))
        for c in range(K):
            r, cc = divmod(c, cols)
            ax = fig.add_subplot(rows, cols, c + 1)
            ax.imshow(self._disp_no_mask(avgs[c]), cmap="inferno",
                       aspect="equal", interpolation="nearest")
            ax.set_title(f"p{c}  N={counts[c]}", fontsize=10)
            ax.set_xticks([]); ax.set_yticks([])
            if qpx > 0:
                add_recip_scalebar(ax, q_per_disp_px=qpx, length_q=0.2)
        fig.suptitle(f"{self.sample} — class averages (no central mask)",
                      fontsize=11)
        self._redraw()
        self._set_status("class averages rendered.")

    # ---- Shared 1D-radial profile compute (used by basic + advanced) ----
    def _compute_class_radials(self):
        """Return dict with keys profiles (K, n_r), xs, xlab, K, counts,
        cc, ps, pmc, q_per_bin (or 0 if uncalibrated). Mirrors the
        compute_radial_profile.py paper pipeline; reads cc/ps/pmc from
        run_summary.json so they exactly match the trained run."""
        import torch
        from torchvision.transforms import v2 as T
        from torchvision.transforms import InterpolationMode
        from dino_sr_ablation import PolarTransform, PolarMaskLeft

        avgs = self._compute_class_averages()
        K = int(self._inf["soft_probs"].shape[1])
        counts = np.bincount(self._inf["assigns"], minlength=K).tolist()
        cc, ps, pmc = 140, 192, 45
        rs = os.path.join(self.outdir or "", "run_summary.json")
        if os.path.exists(rs):
            try:
                js = json.load(open(rs))
                cfg = js.get("cfg", {})
                cc  = int(cfg.get("center_crop_size", cc))
                ps  = int(cfg.get("polar_size", ps))
                pmc = int(cfg.get("polar_mask_cols", pmc))
            except Exception:
                pass
        pre = T.Compose([
            T.CenterCrop(cc),
            T.Resize(ps, interpolation=InterpolationMode.BILINEAR,
                       antialias=True),
            PolarTransform(output_size=ps),
            PolarMaskLeft(k_cols=pmc),
        ])
        profiles = []
        with torch.no_grad():
            for c in range(K):
                x = (torch.from_numpy(avgs[c]).unsqueeze(0).unsqueeze(0)
                       .float())
                pol = pre(x)
                I_r = pol.sum(dim=(0, 1, 2)).cpu().numpy()
                profiles.append(np.maximum(I_r, 1e-6))
        n_r = profiles[0].shape[0]

        # X-axis calibration (1/nm). Convert to 1/Å downstream.
        rp = self._recip_per_px()
        raw_size = None
        if rp > 0 and self._cube_path and os.path.exists(self._cube_path):
            try:
                from gui_app.pre_panel import _open_lazy as _ol
                _arr = _ol(self._cube_path)
                raw_size = int(_arr.shape[-1])
            except Exception:
                raw_size = None
        if rp > 0:
            scale = (cc / ps) * ((raw_size or 192) / 192.0)
            q_per_bin = rp * scale         # 1/nm per polar bin
            xs = np.arange(n_r) * q_per_bin
            xlab = (f"q  (nm⁻¹)   "
                    f"[1 polar-bin = {scale:.3f} raw-px]")
        else:
            xs = np.arange(n_r).astype(float)
            xlab = "polar r-bin  (set top-bar reciprocal res for nm⁻¹)"
            q_per_bin = 0.0
        return dict(profiles=np.asarray(profiles), xs=xs, xlab=xlab,
                     K=K, counts=counts,
                     cc=cc, ps=ps, pmc=pmc, q_per_bin=q_per_bin)

    @staticmethod
    def _subtract_poly_baseline(y, deg=3):
        """Robust polynomial baseline subtraction. Fits the lower 70% of
        the profile, then subtracts. Negative residuals clipped to 0
        when y is on a log-scale axis."""
        n = len(y)
        x = np.arange(n)
        thresh = np.percentile(y, 70)
        mask = y < thresh
        if mask.sum() < deg + 2:
            mask = np.ones_like(y, dtype=bool)
        try:
            p = np.polyfit(x[mask], y[mask], deg)
            bg = np.polyval(p, x)
        except Exception:
            return y, np.zeros_like(y)
        return y - bg, bg

    def _render_radial_1d(self):
        """Per-class 1D radial profile I(r) of the class-average pattern.

        Pipeline mirrors `compute_radial_profile.py` (the paper's canonical
        1D radial that feeds cluster1d):
            CenterCrop(cc) → Resize(ps) → PolarTransform(ps)
                             → PolarMaskLeft(pmc) → sum over θ
        cc, ps, pmc default to 140/192/45 (paper master) and are read
        from <outdir>/run_summary.json when present so they exactly
        match the trained run.
        """
        if not self._ensure_inference(): return
        self._refresh_class_dropdown_lazy()
        self._set_status("computing per-class radial profiles "
                          "(paper compute_radial_profile.py pipeline)…")

        import torch
        from torchvision.transforms import v2 as T
        from torchvision.transforms import InterpolationMode
        from dino_sr_ablation import PolarTransform, PolarMaskLeft

        avgs = self._compute_class_averages()
        K = int(self._inf["soft_probs"].shape[1])
        counts = np.bincount(self._inf["assigns"], minlength=K).tolist()

        # Pull the run's actual pipeline params from run_summary.json so
        # the plot uses the same crop / mask as the trained model. Fall
        # back to paper-master defaults if missing.
        cc, ps, pmc = 140, 192, 45
        rs = os.path.join(self.outdir or "", "run_summary.json")
        if os.path.exists(rs):
            try:
                js = json.load(open(rs))
                cfg = js.get("cfg", {})
                cc  = int(cfg.get("center_crop_size", cc))
                ps  = int(cfg.get("polar_size", ps))
                pmc = int(cfg.get("polar_mask_cols", pmc))
            except Exception:
                pass

        pre = T.Compose([
            T.CenterCrop(cc),
            T.Resize(ps, interpolation=InterpolationMode.BILINEAR,
                       antialias=True),
            PolarTransform(output_size=ps),
            PolarMaskLeft(k_cols=pmc),
        ])

        # Apply once per class to its top-200 confidence-weighted average.
        profiles = []
        with torch.no_grad():
            for c in range(K):
                x = (torch.from_numpy(avgs[c]).unsqueeze(0).unsqueeze(0)
                       .float())                          # (1, 1, 192, 192)
                pol = pre(x)                              # (1, 1, θ, r)
                I_r = pol.sum(dim=(0, 1, 2)).cpu().numpy() # (r,)
                profiles.append(np.maximum(I_r, 1e-6))
        n_r = profiles[0].shape[0]

        # X-axis. The user's `recip_res` is in nm⁻¹/raw-detector-px.
        # After LoadPRZ does cv2.resize(raw → 192) and the transform does
        # CenterCrop(cc) → Resize(ps), one polar r-bin spans
        #     (cc / ps) * (raw_size / 192) raw-px.
        # We probe raw_size from the sample's cube on disk (mmap header
        # only — no full read).
        rp = self._recip_per_px()
        raw_size = None
        if rp > 0 and self._cube_path and os.path.exists(self._cube_path):
            try:
                from gui_app.pre_panel import _open_lazy as _ol
                _arr = _ol(self._cube_path)
                raw_size = int(_arr.shape[-1])
            except Exception:
                raw_size = None
        if rp > 0:
            scale = (cc / ps) * ((raw_size or 192) / 192.0)
            q_per_bin = rp * scale
            xs = np.arange(n_r) * q_per_bin
            xlab = (f"q  (nm⁻¹)   "
                    f"[1 polar-bin = {scale:.3f} raw-px"
                    + (f", raw frame = {raw_size}²"
                        if raw_size else ", raw size unknown") + "]")
        else:
            xs = np.arange(n_r).astype(float)
            xlab = "polar r-bin  (set top-bar reciprocal res for nm⁻¹)"

        cmap = _adaptive_cmap(K)
        colors = [cmap(c) for c in range(K)]

        fig = self._new_fig()
        ax = fig.add_subplot(111)
        lines = []
        for c in range(K):
            line, = ax.plot(xs, profiles[c], color=colors[c], lw=1.4,
                              label=f"p{c}  N={counts[c]}",
                              picker=5)
            lines.append(line)
        # mark the polar-mask edge so the user can see what the model
        # zeros out
        if pmc > 0:
            edge = (xs[pmc - 1] if pmc - 1 < len(xs) else xs[-1])
            ax.axvline(edge, color="#888", lw=0.8, ls=":")
            ax.text(edge, ax.get_ylim()[1] if False else profiles[0].max(),
                     f" mask (k<{pmc})", color="#666", fontsize=8,
                     va="top", ha="left")
        ax.set_yscale("log")
        ax.set_xlabel(xlab)
        ax.set_ylabel("∫ I(θ, r) dθ   (log scale)")
        ax.set_title(
            f"{self.sample} — per-class 1D radial  "
            f"[CenterCrop({cc}) → Resize({ps}) → Polar({ps}) → "
            f"MaskLeft({pmc})]\n"
            f"matches compute_radial_profile.py (paper pipeline)",
            fontsize=10)
        ax.grid(alpha=0.3, which="both")
        leg = ax.legend(fontsize=8, ncol=2, loc="upper right")

        # legend-pick: click an entry to toggle that curve on/off
        legmap = {}
        for legline, line in zip(leg.get_lines(), lines):
            legline.set_picker(True); legline.set_pickradius(8)
            legmap[legline] = line

        def _on_pick(event):
            if event.artist in legmap:
                line = legmap[event.artist]
                vis = not line.get_visible()
                line.set_visible(vis)
                event.artist.set_alpha(1.0 if vis else 0.25)
                self._canvas.draw_idle()
        self._pick_cids = [self._canvas.mpl_connect("pick_event", _on_pick)]

        # crosshair cursor for reading off (q, I)
        self._enable_crosshair_on(ax)
        self._redraw()
        self._set_status("1D radial profiles rendered "
                          "(click legend entry to toggle, hover for crosshair).")

    # ---- 1D advanced popup: per-class show/hide, baseline, CIF ----
    def _open_radial_1d_popup(self):
        """Toplevel popup with class checkboxes, polynomial-baseline
        subtraction, and CIF g-line overlay over the per-class 1D
        radial profiles."""
        if not self._ensure_inference():
            return
        try:
            data = self._compute_class_radials()
        except Exception as e:
            messagebox.showerror("1D advanced", f"compute failed:\n{e!r}")
            return
        K        = data["K"]
        profiles = data["profiles"]
        xs_nm    = data["xs"]            # 1/nm  (or bin index if uncal.)
        xlab     = data["xlab"]
        counts   = data["counts"]
        pmc      = data["pmc"]
        q_per_bin_nm = data["q_per_bin"]  # 1/nm per bin (0 if uncal.)
        # Display in nm⁻¹ (matches the top-bar reciprocal calibration
        # and the rest of the GUI). CIF g-magnitudes are in Å⁻¹ and
        # get converted to nm⁻¹ on overlay (×10).
        if q_per_bin_nm > 0:
            xs = xs_nm
            xlab_disp = "q  (nm⁻¹)"
            q_max_inv_nm = float(xs.max())
            calibrated = True
        else:
            xs = xs_nm
            xlab_disp = xlab
            q_max_inv_nm = 0.0
            calibrated = False

        win = tk.Toplevel(self)
        win.title(f"{self.sample}  —  1D radial advanced")
        win.geometry("1180x720")

        # ---- top control rows ----
        ctrl_top = ctk.CTkFrame(win)
        ctrl_top.pack(side="top", fill="x", padx=4, pady=2)
        ctk.CTkLabel(ctrl_top, text="show:",
                       font=("Segoe UI", 10, "bold")
                       ).pack(side="left", padx=(4, 2))
        show_vars = []
        cmap = _adaptive_cmap(K)
        for c in range(K):
            v = ctk.BooleanVar(value=True)
            show_vars.append(v)
            cb = ctk.CTkCheckBox(ctrl_top, text=f"p{c}",
                                    variable=v, width=50,
                                    command=lambda: _redraw())
            cb.pack(side="left", padx=1)
        ctk.CTkButton(ctrl_top, text="all",
                       width=40,
                       command=lambda: (
                           [v.set(True) for v in show_vars], _redraw())
                       ).pack(side="left", padx=4)
        ctk.CTkButton(ctrl_top, text="none",
                       width=40,
                       command=lambda: (
                           [v.set(False) for v in show_vars], _redraw())
                       ).pack(side="left", padx=2)

        ctrl_mid = ctk.CTkFrame(win)
        ctrl_mid.pack(side="top", fill="x", padx=4, pady=2)
        log_y_var  = ctk.BooleanVar(value=True)
        sub_bg_var = ctk.BooleanVar(value=False)
        deg_var    = ctk.IntVar(value=3)
        ctk.CTkCheckBox(ctrl_mid, text="log y",
                          variable=log_y_var,
                          command=lambda: _redraw()
                          ).pack(side="left", padx=4)
        ctk.CTkCheckBox(ctrl_mid, text="subtract polynomial baseline",
                          variable=sub_bg_var,
                          command=lambda: _redraw()
                          ).pack(side="left", padx=4)
        ctk.CTkLabel(ctrl_mid, text="poly deg:").pack(side="left",
                                                          padx=(4, 0))
        deg_entry = ctk.CTkEntry(ctrl_mid, textvariable=deg_var, width=40)
        deg_entry.pack(side="left", padx=2)

        ctrl_cif = ctk.CTkFrame(win)
        ctrl_cif.pack(side="top", fill="x", padx=4, pady=2)
        ctk.CTkLabel(ctrl_cif, text="CIF:",
                       font=("Segoe UI", 10, "bold")
                       ).pack(side="left", padx=(4, 2))
        cif_var   = ctk.StringVar(value="")
        kmax_var  = ctk.DoubleVar(value=2.0)
        show_cif  = ctk.BooleanVar(value=True)
        ctk.CTkEntry(ctrl_cif, textvariable=cif_var,
                       width=300).pack(side="left", padx=2)
        ctk.CTkButton(ctrl_cif, text="Browse",
                       width=70,
                       command=lambda: _pick_cif()).pack(side="left",
                                                            padx=2)
        ctk.CTkLabel(ctrl_cif, text="k_max (1/Å):"
                       ).pack(side="left", padx=(8, 2))
        ctk.CTkEntry(ctrl_cif, textvariable=kmax_var, width=60
                       ).pack(side="left", padx=2)
        ctk.CTkButton(ctrl_cif, text="Load + overlay",
                       command=lambda: _load_cif()).pack(side="left",
                                                            padx=4)
        ctk.CTkCheckBox(ctrl_cif, text="show CIF g-lines",
                          variable=show_cif,
                          command=lambda: _redraw()
                          ).pack(side="left", padx=4)
        ctk.CTkButton(ctrl_cif, text="Save snapshot",
                       command=lambda: _save_snap()
                       ).pack(side="right", padx=4)

        # status under controls
        status_lbl = ctk.CTkLabel(win, text="", font=("Consolas", 9),
                                     anchor="w")
        status_lbl.pack(side="top", fill="x", padx=8)

        # ---- canvas ----
        fig = Figure(figsize=(11.0, 5.0), dpi=110, facecolor="white")
        ax = fig.add_subplot(111)
        canv = FigureCanvasTkAgg(fig, master=win)
        canv.get_tk_widget().pack(fill="both", expand=True)
        NavigationToolbar2Tk(canv, win)

        # ---- shared CIF state ----
        cif_state = dict(g_mag=None, sf_int=None, label="")

        def _pick_cif():
            p = filedialog.askopenfilename(
                filetypes=[("CIF", "*.cif"), ("All", "*.*")])
            if p:
                cif_var.set(p)

        def _load_cif():
            p = cif_var.get().strip()
            if not p or not os.path.exists(p):
                messagebox.showerror("CIF", f"Not found: {p!r}"); return
            try:
                from py4DSTEM.process.diffraction import Crystal
            except Exception as e:
                messagebox.showerror("py4DSTEM",
                    f"py4DSTEM unavailable:\n{e!r}"); return
            try:
                try:
                    cr = Crystal.from_CIF(p)
                except AttributeError:
                    cr = Crystal(filepath=p)
                cr.calculate_structure_factors(float(kmax_var.get()))
                g_vec = np.asarray(getattr(cr, "g_vec_all", None))
                g_mag = np.linalg.norm(g_vec, axis=0)
                sf_int = getattr(cr, "struct_factors_int", None)
                if sf_int is None:
                    sf = getattr(cr, "struct_factors", None)
                    sf_int = np.abs(sf) ** 2 if sf is not None else None
                cif_state["g_mag"] = np.asarray(g_mag)
                cif_state["sf_int"] = (np.asarray(sf_int)
                                          if sf_int is not None else None)
                cif_state["label"] = os.path.basename(p)
                status_lbl.configure(text=(
                    f"loaded {cif_state['label']}: "
                    f"{cif_state['g_mag'].size} reflections "
                    f"@ k_max = {kmax_var.get()} Å⁻¹"))
            except Exception as e:
                messagebox.showerror("CIF", repr(e)); return
            _redraw()

        def _redraw():
            ax.clear()
            # Curves per class (with optional baseline subtraction)
            log_y = log_y_var.get()
            do_bg = sub_bg_var.get()
            try:
                deg = max(1, int(deg_var.get()))
            except Exception:
                deg = 3
            any_visible = False
            for c in range(K):
                if not show_vars[c].get():
                    continue
                any_visible = True
                y = profiles[c].astype(np.float64).copy()
                if do_bg:
                    y, _bg = self._subtract_poly_baseline(y, deg=deg)
                # log y can't show negatives; clip if needed
                if log_y:
                    y = np.maximum(y, 1e-6)
                ax.plot(xs, y, color=cmap(c), lw=1.4,
                          label=f"p{c}  N={counts[c]}")
            # Polar-mask edge marker
            if pmc > 0 and any_visible:
                edge = (xs[pmc - 1] if pmc - 1 < len(xs) else xs[-1])
                ax.axvline(edge, color="#888", lw=0.8, ls=":")
            # CIF overlay: vertical lines at predicted q. py4DSTEM gives
            # g-magnitudes in Å⁻¹; we display in nm⁻¹, so ×10.
            if (show_cif.get()
                    and cif_state["g_mag"] is not None
                    and calibrated):
                gm_nm = cif_state["g_mag"] * 10.0      # Å⁻¹ → nm⁻¹
                si = cif_state["sf_int"]
                if si is None or si.size != gm_nm.size:
                    si = np.ones_like(gm_nm)
                mx = float(si.max()) if si.size else 1.0
                # Get current y-range to scale the vertical bars.
                y0, y1 = ax.get_ylim()
                if log_y:
                    y0 = max(y0, 1e-6)
                for k in range(gm_nm.size):
                    if gm_nm[k] > q_max_inv_nm:
                        continue
                    if si[k] < mx * 1e-3:   # filter very weak refls
                        continue
                    alpha = 0.25 + 0.6 * (si[k] / mx)
                    ax.axvline(gm_nm[k], color="orange",
                                 alpha=float(alpha), lw=0.9)
                ax.plot([], [], color="orange", lw=1.0,
                          label=f"CIF: {cif_state['label']}")
            ax.set_xlabel(xlab_disp)
            ax.set_ylabel("∫ I(θ, r) dθ"
                            + ("  (baseline subtracted)" if do_bg else ""))
            if log_y:
                ax.set_yscale("log")
            else:
                ax.set_yscale("linear")
            ax.set_title(
                f"{self.sample} — per-class 1D radial   "
                + ("[poly-bg subtracted]" if do_bg else ""),
                fontsize=10)
            ax.grid(alpha=0.3, which="both")
            try:
                ax.legend(fontsize=8, ncol=2, loc="upper right")
            except Exception:
                pass
            fig.tight_layout()
            canv.draw_idle()

        def _save_snap():
            if self.outdir is None:
                messagebox.showinfo("save", "no run linked"); return
            out_dir = os.path.join(self.outdir, "radial_1d_advanced")
            os.makedirs(out_dir, exist_ok=True)
            tag_bits = []
            shown = [c for c in range(K) if show_vars[c].get()]
            tag_bits.append("p" + "_".join(str(c) for c in shown))
            if sub_bg_var.get(): tag_bits.append(f"poly{deg_var.get()}")
            if (cif_state["g_mag"] is not None
                    and show_cif.get()):
                tag_bits.append("cif")
            tag = "_".join(tag_bits) or "snap"
            png = os.path.join(out_dir, f"radial1d_{tag}.png")
            try:
                fig.savefig(png, dpi=140)
            except Exception as e:
                messagebox.showerror("save", repr(e)); return
            status_lbl.configure(text=f"saved → {png}")

        deg_entry.bind("<Return>", lambda _e: _redraw())
        _redraw()

    def _render_radial_1d_saxs(self):
        """SAXS-treated 1D residuals — exactly what cluster1d_loss compares.

        Mirrors `compute_radial_profile.py` end to end:
            CenterCrop → Resize → Polar → MaskLeft → sum-θ → I(r)
            [Q_LOW:Q_HIGH] → /Σ → log → polynomial-baseline-subtract
            → mean-center  →  residual

        Two stacked axes: top shows log(I/ΣI) + the fitted polynomial
        baseline (dotted), bottom shows the residual that the cosine
        loss actually sees.  Per-class centroid pairwise cosine printed
        to the status bar — that's the signal-strength diagnostic
        (healthy = wide spread, dead = all positive and clustered)."""
        if not self._ensure_inference(): return
        self._refresh_class_dropdown_lazy()
        self._set_status("computing per-class SAXS residuals "
                          "(loss-side 1D)…")

        import torch
        from torchvision.transforms import v2 as T
        from torchvision.transforms import InterpolationMode
        from dino_sr_ablation import PolarTransform, PolarMaskLeft
        from compute_radial_profile import (Q_LOW, Q_HIGH, POLY_ORDER,
                                                EPS)

        avgs = self._compute_class_averages()
        K = int(self._inf["soft_probs"].shape[1])
        counts = np.bincount(self._inf["assigns"], minlength=K).tolist()

        # Per-run pipeline params (fall back to paper-master defaults)
        cc, ps, pmc = 140, 192, 45
        rs = os.path.join(self.outdir or "", "run_summary.json")
        if os.path.exists(rs):
            try:
                js = json.load(open(rs))
                cfg = js.get("cfg", {})
                cc  = int(cfg.get("center_crop_size", cc))
                ps  = int(cfg.get("polar_size", ps))
                pmc = int(cfg.get("polar_mask_cols", pmc))
            except Exception:
                pass

        pre = T.Compose([
            T.CenterCrop(cc),
            T.Resize(ps, interpolation=InterpolationMode.BILINEAR,
                       antialias=True),
            PolarTransform(output_size=ps),
            PolarMaskLeft(k_cols=pmc),
        ])

        # Step 1 — I(r) per class average
        Irs = []
        with torch.no_grad():
            for c in range(K):
                x = (torch.from_numpy(avgs[c]).unsqueeze(0).unsqueeze(0)
                       .float())
                pol = pre(x)
                Ir = pol.sum(dim=(0, 1, 2)).cpu().numpy()
                Irs.append(np.maximum(Ir, EPS))

        # Step 2 — SAXS treatment, identical to compute_radial_profile.py
        n_keep = Q_HIGH - Q_LOW
        q = np.arange(n_keep, dtype=np.float64)
        residuals, raw_logs, baselines = [], [], []
        for c in range(K):
            Ir_keep = Irs[c][Q_LOW:Q_HIGH].astype(np.float64)
            s = max(float(Ir_keep.sum()), EPS)
            log_I = np.log(Ir_keep / s + EPS)
            raw_logs.append(log_I)
            # First pass polynomial fit
            coef = np.polyfit(q, log_I, POLY_ORDER)
            bg = np.polyval(coef, q)
            r = log_I - bg
            med = np.median(r); mad = np.median(np.abs(r - med)) + 1e-12
            peak_mask = (r - med) > 2.0 * mad
            if peak_mask.sum() < n_keep - POLY_ORDER - 1:
                non_peak = ~peak_mask
                coef = np.polyfit(q[non_peak], log_I[non_peak], POLY_ORDER)
                bg = np.polyval(coef, q)
            r2 = log_I - bg
            residuals.append((r2 - r2.mean()).astype(np.float64))
            baselines.append(bg)

        # X-axis (same scaling as the main 1D plot)
        rp = self._recip_per_px()
        raw_size = None
        if rp > 0 and self._cube_path and os.path.exists(self._cube_path):
            try:
                from gui_app.pre_panel import _open_lazy as _ol
                _arr = _ol(self._cube_path)
                raw_size = int(_arr.shape[-1])
            except Exception:
                raw_size = None
        if rp > 0:
            scale = (cc / ps) * ((raw_size or 192) / 192.0)
            q_per_bin = rp * scale
            xs = np.arange(Q_LOW, Q_HIGH) * q_per_bin
            xlab = (f"q  (nm⁻¹)   "
                    f"[1 polar-bin = {scale:.3f} raw-px"
                    + (f", raw frame = {raw_size}²"
                       if raw_size else ", raw size unknown") + "]")
        else:
            xs = np.arange(Q_LOW, Q_HIGH).astype(float)
            xlab = "polar r-bin"

        cmap = _adaptive_cmap(K)
        colors = [cmap(c) for c in range(K)]

        fig = self._new_fig()
        ax_top = fig.add_subplot(2, 1, 1)
        ax_bot = fig.add_subplot(2, 1, 2)

        # One row of 3 artists per class so legend click can toggle the
        # whole class (top log line + dotted baseline + bottom residual)
        # together.
        cls_artists = []
        for c in range(K):
            top_line, = ax_top.plot(xs, raw_logs[c], color=colors[c],
                                       lw=1.0, alpha=0.85,
                                       label=f"p{c} N={counts[c]}")
            base_line, = ax_top.plot(xs, baselines[c], color=colors[c],
                                        lw=0.8, ls=":", alpha=0.6)
            bot_line, = ax_bot.plot(xs, residuals[c], color=colors[c],
                                       lw=1.5,
                                       label=f"p{c} N={counts[c]}")
            cls_artists.append((top_line, base_line, bot_line))

        ax_top.set_ylabel("log(I(r) / ΣI)")
        ax_top.set_title(
            f"{self.sample} — SAXS treatment (matches "
            f"compute_radial_profile.py)\n"
            f"top: log-norm I(r) (solid) + polynomial baseline order "
            f"{POLY_ORDER} (dotted), bins {Q_LOW}..{Q_HIGH}",
            fontsize=10)
        ax_top.grid(alpha=0.3)
        leg = ax_top.legend(fontsize=8, ncol=2, loc="best")

        # Map each legend entry to the 3 artists for that class.  Click
        # toggles all 3 + dims the legend entry so the user can see at
        # a glance which classes are off.
        legmap = {}
        for legline, artists in zip(leg.get_lines(), cls_artists):
            legline.set_picker(True); legline.set_pickradius(8)
            legmap[legline] = artists

        def _on_pick(event):
            arts = legmap.get(event.artist)
            if not arts:
                return
            # Use first artist's visibility as the source of truth.
            vis = not arts[0].get_visible()
            for a in arts:
                a.set_visible(vis)
            event.artist.set_alpha(1.0 if vis else 0.25)
            self._canvas.draw_idle()
        self._pick_cids = [
            self._canvas.mpl_connect("pick_event", _on_pick)]

        ax_bot.axhline(0, color="#888", lw=0.5)
        ax_bot.set_xlabel(xlab)
        ax_bot.set_ylabel("residual  (mean-centered)")
        ax_bot.set_title("bottom: residuals — exactly what "
                          "cluster1d_loss compares via cosine",
                          fontsize=10)
        ax_bot.grid(alpha=0.3)

        # Diagnostic: pairwise cosine between per-class residual centroids
        R = np.stack(residuals, 0)
        n = np.linalg.norm(R, axis=1) + 1e-8
        cosM = (R @ R.T) / np.outer(n, n)
        off = cosM[~np.eye(K, dtype=bool)]
        # Health rule of thumb (described in chat): healthy clustering wants
        # wide spread (off-diag cosine spans both signs and isn't bunched
        # near +0.3-0.4).
        if off.size:
            health = (
                "OK" if (off.max() < 0.7 and off.min() < 0.0)
                else "WARN: residuals look bunched — likely SAXS treatment "
                      "is removing the discriminative signal"
            )
        else:
            health = "n/a"
        msg = (
            f"SAXS residuals rendered (click legend entry to toggle a "
            f"class on/off).  per-class centroid pairwise cosine "
            f"(K={K}, off-diag): "
            f"min={off.min() if off.size else 0:+.2f}, "
            f"max={off.max() if off.size else 0:+.2f}, "
            f"mean={off.mean() if off.size else 0:+.2f}. "
            f"[{health}]")
        self._enable_crosshair_on(ax_bot)
        self._redraw()
        self._set_status(msg)

    def _render_centroid_matrix(self):
        if not self._ensure_inference(): return
        embeds = self._inf["embeds"]; ass = self._inf["assigns"]
        K = int(self._inf["soft_probs"].shape[1])
        cents = np.zeros((K, embeds.shape[1]), dtype=np.float32)
        for c in range(K):
            m = (ass == c)
            if m.any():
                cents[c] = embeds[m].mean(0)
        cents = cents / (np.linalg.norm(cents, axis=1, keepdims=True) + 1e-12)
        cos = cents @ cents.T
        fig = self._new_fig()
        ax = fig.add_subplot(111)
        im = ax.imshow(cos, cmap="RdBu_r", vmin=-1, vmax=1,
                        interpolation="nearest")
        for i in range(K):
            for j in range(K):
                ax.text(j, i, f"{cos[i, j]:.2f}", ha="center", va="center",
                         fontsize=9, color="black" if abs(cos[i, j]) < 0.6
                                                    else "white")
        ax.set_title(f"{self.sample} — class-centroid cosine (K={K})",
                      fontsize=11)
        ax.set_xticks(range(K)); ax.set_yticks(range(K))
        fig.colorbar(im, ax=ax, fraction=0.04, pad=0.02)
        self._redraw()
        self._set_status("centroid matrix rendered.")

    # ---- UMAP ----
    def _render_umap(self):
        """2-D UMAP (or PCA fallback) of the per-pattern embeddings,
        coloured by argmax class. Uses the same `umap_or_pca` helper
        that `eval_all` uses for paper figures, so the GUI render
        matches the paper one byte-for-byte."""
        if not self._ensure_inference(): return
        embeds = self._inf.get("embeds")
        if embeds is None:
            messagebox.showinfo("UMAP",
                "This run's inference.npz has no 'embeds' array."); return
        ass = self._inf["assigns"]
        K = int(self._inf["soft_probs"].shape[1])
        N = embeds.shape[0]
        self._set_status(f"computing UMAP on {N} embeddings…")

        try:
            from contrastive_eval import umap_or_pca
            coords, method = umap_or_pca(embeds)
        except Exception as e:
            messagebox.showerror("UMAP failed", repr(e)); return

        cmap = _adaptive_cmap(K)
        colors = np.array([cmap(int(c))[:3] for c in ass],
                            dtype=np.float32)

        fig = self._new_fig()
        ax = fig.add_subplot(111)
        # Plot per-class so the legend is clean
        for c in range(K):
            m = (ass == c)
            if not m.any(): continue
            ax.scatter(coords[m, 0], coords[m, 1],
                        s=4, alpha=0.6,
                        color=cmap(c), edgecolors="none",
                        label=f"p{c} (N={int(m.sum())})")
        ax.set_aspect("equal", adjustable="datalim")
        ax.set_xticks([]); ax.set_yticks([])
        ax.set_title(
            f"{self.sample} — {method} of student embeddings  "
            f"(K={K}, N={N})", fontsize=11)
        ax.legend(loc="best", fontsize=8, ncol=2,
                   markerscale=2.5, framealpha=0.85)
        self._redraw()
        self._set_status(f"{method} rendered  "
                          f"(N={N}, dim={embeds.shape[1]} → 2).")

    # ---- BF / HAADF ----
    def _active_cube_path(self):
        """Cube to analyze: the run-linked cube if present, else the cube
        loaded in the Dataset tab (top of the app)."""
        if self._cube_path:
            return self._cube_path
        pre = getattr(self.app, "pre", None)
        try:
            p = pre.get_loaded_path() if pre is not None else None
        except Exception:
            p = None
        return p or None

    def _ensure_BF_HAADF(self):
        if self._BF is not None and self._HA is not None:
            return True
        cube_path = self._active_cube_path()
        if not cube_path:
            messagebox.showinfo("No dataset",
                "Load a cube in the Dataset tab (top) first.")
            return False
        self._set_status("computing BF + HAADF (~10–30 s) …")
        try:
            cube = _open_lazy(cube_path)
            Ny, Nx, H, W = cube.shape
            r_BF = 0.06 * H
            r_in = 0.18 * H; r_out = 0.45 * H
            bm = _radial_mask(H, W, 0, r_BF)
            ham = _radial_mask(H, W, r_in, r_out)
            BF = np.zeros((Ny, Nx), dtype=np.float64)
            HA = np.zeros((Ny, Nx), dtype=np.float64)
            for y in range(Ny):
                block = np.asarray(cube[y]).astype(np.float32)
                BF[y] = (block * bm).sum(axis=(1, 2))
                HA[y] = (block * ham).sum(axis=(1, 2))
            self._BF = BF; self._HA = HA
            return True
        except Exception as e:
            messagebox.showerror("BF/HAADF failed", str(e))
            return False

    def _render_bf(self):
        if not self._ensure_BF_HAADF(): return
        img = self._BF
        fig = self._new_fig()
        ax = fig.add_subplot(111)
        lo, hi = np.percentile(img, 1), np.percentile(img, 99)
        im = ax.imshow(np.clip(img, lo, hi), cmap="gray",
                        aspect="equal", interpolation="nearest")
        ax.set_title(f"{self.sample} — virtual BF "
                      f"(r ≤ {0.06 * 192:.0f} px proportional)",
                      fontsize=11)
        ax.set_xticks([]); ax.set_yticks([])
        fig.colorbar(im, ax=ax, fraction=0.04, pad=0.02)
        self._redraw()
        self._set_status("virtual BF rendered.")

    def _render_haadf(self):
        if not self._ensure_BF_HAADF(): return
        img = self._HA
        fig = self._new_fig()
        ax = fig.add_subplot(111)
        lo, hi = np.percentile(img, 1), np.percentile(img, 99)
        im = ax.imshow(np.clip(img, lo, hi), cmap="gray",
                        aspect="equal", interpolation="nearest")
        ax.set_title(f"{self.sample} — virtual HAADF "
                      f"(annulus 0.18–0.45 of frame)",
                      fontsize=11)
        ax.set_xticks([]); ax.set_yticks([])
        fig.colorbar(im, ax=ax, fraction=0.04, pad=0.02)
        self._redraw()
        self._set_status("virtual HAADF rendered.")

    # ---- arbitrary annular virtual detector (user-chosen r, dr) ----
    def _open_annular_popup(self):
        """Popup to build a virtual image from an arbitrary annular detector.

        The user picks an inner radius r and ring width dr (detector pixels).
        The left panel shows a diffraction pattern — either a single scan
        frame (changeable) or the DP-max over the whole scan — with the
        annulus (r .. r+dr) drawn on it, so the chosen detector is visible.
        'Compute map' integrates that annulus at every probe position to give
        the virtual annular image on the right.
        """
        if not self._active_cube_path():
            messagebox.showinfo("Annular detector",
                "Load a cube in the Dataset tab (top) first.")
            return
        from matplotlib.figure import Figure
        from matplotlib.backends.backend_tkagg import (
            FigureCanvasTkAgg, NavigationToolbar2Tk)
        from matplotlib.patches import Circle

        st = {"cube": None, "shape": None, "dpmax": None}

        win = ctk.CTkToplevel(self)
        win.title(f"Annular virtual detector — {self.sample or ''}")
        win.geometry("1000x640")
        try:
            win.after(200, lambda: (win.lift(), win.focus_force()))
        except Exception:
            pass

        ctrl = ctk.CTkFrame(win)
        ctrl.pack(side="top", fill="x", padx=8, pady=6)
        ctk.CTkLabel(ctrl, text="r (px):").pack(side="left", padx=(8, 2))
        r_var = ctk.StringVar(value="")
        ctk.CTkEntry(ctrl, textvariable=r_var, width=70).pack(side="left")
        ctk.CTkLabel(ctrl, text="dr (px):").pack(side="left", padx=(10, 2))
        dr_var = ctk.StringVar(value="")
        ctk.CTkEntry(ctrl, textvariable=dr_var, width=70).pack(side="left")
        ctk.CTkLabel(ctrl, text="background:").pack(side="left", padx=(12, 2))
        bg_var = ctk.StringVar(value="single frame")
        ctk.CTkOptionMenu(ctrl, values=["single frame", "DP-max"],
                          variable=bg_var, width=130,
                          command=lambda _v: _update_view()).pack(side="left")
        ctk.CTkLabel(ctrl, text="frame idx:").pack(side="left", padx=(10, 2))
        idx_var = ctk.StringVar(value="0")
        ctk.CTkEntry(ctrl, textvariable=idx_var, width=80).pack(side="left")
        ctk.CTkButton(ctrl, text="Update view", width=100,
                      command=lambda: _update_view()).pack(side="left", padx=(12, 2))
        ctk.CTkButton(ctrl, text="Compute map", width=110,
                      fg_color=("#2D7A2D", "#1F7A1F"),
                      command=lambda: _compute_map()).pack(side="left", padx=2)
        # Frame slider (single-frame background) — drag to scrub frames.
        frame_row = ctk.CTkFrame(win, fg_color="transparent")
        frame_row.pack(side="top", fill="x", padx=8, pady=(2, 0))
        ctk.CTkLabel(frame_row, text="frame:").pack(side="left", padx=(2, 6))
        frame_slider = ctk.CTkSlider(
            frame_row, from_=0, to=1, number_of_steps=1,
            command=lambda v: (idx_var.set(str(int(float(v)))),
                               _update_view() if bg_var.get() == "single frame"
                               else None))
        frame_slider.pack(side="left", fill="x", expand=True, padx=4)
        frame_slider.set(0)
        status = ctk.CTkLabel(win, text="", anchor="w")
        status.pack(side="bottom", fill="x", padx=8, pady=(0, 4))

        fig = Figure(figsize=(9.6, 5.0))
        ax_dp = fig.add_subplot(121)
        ax_map = fig.add_subplot(122)
        for a in (ax_dp, ax_map):
            a.set_xticks([]); a.set_yticks([])
        canvas = FigureCanvasTkAgg(fig, master=win)
        canvas.get_tk_widget().pack(side="top", fill="both", expand=True)
        NavigationToolbar2Tk(canvas, win)

        def _load_cube():
            if st["cube"] is None:
                status.configure(text="opening cube…"); win.update_idletasks()
                st["cube"] = _open_lazy(self._active_cube_path())
                st["shape"] = tuple(st["cube"].shape)
                H = st["shape"][2]
                if not r_var.get():
                    r_var.set(f"{0.25 * H:.0f}")
                if not dr_var.get():
                    dr_var.set(f"{0.06 * H:.0f}")
                N = int(st["shape"][0]) * int(st["shape"][1])
                try:
                    frame_slider.configure(to=max(1, N - 1),
                                           number_of_steps=max(1, N - 1))
                except Exception:
                    pass
            return st["cube"]

        def _vals():
            try: r = max(0.0, float(r_var.get()))
            except Exception: r = 0.0
            try: dr = max(0.0, float(dr_var.get()))
            except Exception: dr = 0.0
            return r, dr

        def _get_bg():
            cube = _load_cube()
            Ny, Nx, H, W = st["shape"]
            if bg_var.get() == "DP-max":
                if st["dpmax"] is None:
                    status.configure(
                        text="computing DP-max (reads all frames, ~10–30 s)…")
                    win.update_idletasks()
                    acc = np.zeros((H, W), np.float32)
                    for y in range(Ny):
                        blk = np.asarray(cube[y]).astype(np.float32)
                        acc = np.maximum(acc, blk.max(axis=0))
                    st["dpmax"] = acc
                return st["dpmax"], "DP-max"
            try: idx = int(float(idx_var.get()))
            except Exception: idx = 0
            idx = max(0, min(Ny * Nx - 1, idx))
            y, x = divmod(idx, Nx)
            try:
                frame = np.asarray(cube[y, x]).astype(np.float32)
            except Exception:
                frame = np.asarray(cube[y])[x].astype(np.float32)
            return frame, f"frame idx={idx}  (y={y}, x={x})"

        def _update_view():
            try:
                bg, label = _get_bg()
            except Exception as e:
                messagebox.showerror("Annular detector", str(e)); return
            H, W = bg.shape
            r, dr = _vals()
            ax_dp.clear()
            ax_dp.imshow(np.log1p(np.clip(bg, 0, None)), cmap="inferno",
                         aspect="equal", interpolation="nearest")
            cy, cx = H / 2.0, W / 2.0
            ax_dp.add_patch(Circle((cx, cy), r, fill=False, ec="cyan", lw=1.4))
            ax_dp.add_patch(Circle((cx, cy), r + dr, fill=False,
                                   ec="yellow", lw=1.4))
            ax_dp.set_title(f"{label}\nannulus  r={r:.0f} → {r + dr:.0f} px",
                            fontsize=9)
            ax_dp.set_xticks([]); ax_dp.set_yticks([])
            canvas.draw_idle()
            status.configure(
                text="view updated — adjust r/dr, then 'Compute map'.")

        def _compute_map():
            try:
                cube = _load_cube()
                Ny, Nx, H, W = st["shape"]
                r, dr = _vals()
                mask = _radial_mask(H, W, r, r + dr)
                n = int(mask.sum())
                if n == 0:
                    messagebox.showinfo("Annular detector",
                        "Annulus contains no pixels — increase dr or r.")
                    return
                status.configure(
                    text=f"computing annular map (r={r:.0f}, dr={dr:.0f}, "
                         f"{n} det px, reads all frames)…")
                win.update_idletasks()
                amap = np.zeros((Ny, Nx), np.float64)
                for y in range(Ny):
                    blk = np.asarray(cube[y]).astype(np.float32)
                    amap[y] = (blk * mask).sum(axis=(1, 2))
                ax_map.clear()
                lo, hi = np.percentile(amap, 1), np.percentile(amap, 99)
                im = ax_map.imshow(np.clip(amap, lo, hi), cmap="gray",
                                   aspect="equal", interpolation="nearest")
                ax_map.set_title(f"annular VDF   r={r:.0f}, dr={dr:.0f} px",
                                 fontsize=10)
                ax_map.set_xticks([]); ax_map.set_yticks([])
                fig.colorbar(im, ax=ax_map, fraction=0.046, pad=0.04)
                _update_view()
                canvas.draw_idle()
                status.configure(
                    text=f"done — r={r:.0f}, dr={dr:.0f} px, {n} detector px. "
                         f"Save via the toolbar (disk icon).")
            except Exception as e:
                messagebox.showerror("Annular detector", str(e))

        _update_view()

    def _render_overlay(self, on_haadf=True):
        if not self._ensure_inference(): return
        if not self._ensure_BF_HAADF(): return
        sel = [c for c, v in self._cls_check_vars.items() if v.get()]
        if not sel:
            messagebox.showinfo("overlay",
                "Tick at least one class in the multi-class panel.")
            return
        Ny, Nx = self._scan_shape
        ass = self._inf["assigns"].reshape(Ny, Nx)
        K = int(self._inf["soft_probs"].shape[1])
        bg = self._HA if on_haadf else self._BF
        lo, hi = np.percentile(bg, 1), np.percentile(bg, 99)
        bg_disp = np.clip(bg, lo, hi)
        # Each selected class gets its own colour from the same adaptive
        # cmap as the class map / legend strip, so the colour code is
        # consistent across the tab.
        cmap = _adaptive_cmap(K)
        overlay = np.zeros((Ny, Nx, 4), dtype=np.float32)
        legend_handles = []
        from matplotlib.patches import Patch
        per_cls_counts = []
        for cls_id in sorted(sel):
            r, g, b = cmap(cls_id)[:3]
            m = (ass == cls_id)
            overlay[m, 0] = r
            overlay[m, 1] = g
            overlay[m, 2] = b
            overlay[m, 3] = 0.60
            per_cls_counts.append((cls_id, int(m.sum())))
            legend_handles.append(
                Patch(facecolor=(r, g, b, 0.85), edgecolor="white",
                      label=f"p{cls_id}  N={int(m.sum())}"))
        fig = self._new_fig()
        ax = fig.add_subplot(111)
        ax.imshow(bg_disp, cmap="gray", aspect="equal",
                   interpolation="nearest")
        ax.imshow(overlay, aspect="equal", interpolation="nearest")
        which = "HAADF" if on_haadf else "BF"
        total_sel = sum(n for _, n in per_cls_counts)
        ax.set_title(
            f"{self.sample} — multi-class overlay on {which}   "
            f"({total_sel}/{Ny*Nx} = {100*total_sel/(Ny*Nx):.1f}% of pixels)",
            fontsize=10)
        ax.set_xticks([]); ax.set_yticks([])
        ax.legend(handles=legend_handles, loc="upper right",
                   fontsize=8, framealpha=0.85)
        self._redraw()
        self._set_status(
            f"overlay rendered on {which} "
            f"(per-class colors, {len(sel)} classes).")

    # ---- per-class / per-frame GradCAM + IG triptych ----
    def _current_class(self):
        try:
            return int(self._class_var.get())
        except (ValueError, TypeError):
            return None

    def _current_frame_idx(self):
        try:
            return int(self._frame_idx_var.get())
        except (ValueError, TypeError):
            return None

    def _render_triptych(self):
        """Class average / single frame  +  GradCAM  +  Integrated
        Gradients, all in one figure. Source = 'class' uses the top-200
        confidence-weighted class average; 'frame' uses one raw frame."""
        if not self._ensure_inference(): return
        if self._busy: return
        source = self._attr_source.get()
        if source == "class":
            c = self._current_class()
            if c is None:
                messagebox.showinfo("class", "Pick a class id first."); return
            target_idx = None
            tag = f"p{c}  (top-200 weighted avg)"
        else:
            target_idx = self._current_frame_idx()
            if target_idx is None:
                messagebox.showinfo("idx",
                    "Type a scan-flat index, or right-click a pixel on "
                    "the class map."); return
            Ny, Nx = (self._scan_shape if self._scan_shape else (None, None))
            if Ny is not None and not (0 <= target_idx < Ny * Nx):
                messagebox.showerror("idx",
                    f"idx {target_idx} out of range for scan {self._scan_shape}.")
                return
            # default target class for the frame = the assigned class
            ass = self._inf["assigns"]
            c_assigned = int(ass[target_idx])
            c_picked   = self._current_class()
            c = c_assigned if c_picked is None else c_picked
            y, x = divmod(target_idx, Nx) if Nx else (target_idx, 0)
            tag = (f"frame idx={target_idx}  (y={y}, x={x})  "
                   f"target=p{c}  assigned=p{c_assigned}")
        self._busy = True
        self._set_status(f"computing GradCAM + IG  ({tag}) …")
        threading.Thread(target=self._triptych_worker,
                          args=(c, source, target_idx, tag), daemon=True).start()

    def _triptych_worker(self, c, source, target_idx, tag):
        try:
            ckpt = self._best_ckpt()
            if ckpt is None:
                raise RuntimeError("no checkpoint")
            raw_full, avg, cam, ig = self._compute_attribution_pair(
                ckpt, c, source, target_idx)
            self.after(0, lambda: self._draw_triptych(
                c, source, target_idx, tag, avg, cam, ig,
                raw_full=raw_full))
        except Exception as e:
            err = str(e)
            self.after(0, lambda: messagebox.showerror("attribution failed", err))
        finally:
            self._busy = False
            self.after(0, lambda: self._set_status(
                f"triptych rendered  ({tag})."))

    def _read_polar_pipeline_cfg(self) -> tuple[int, int, int, bool]:
        """Return (mask_r, polar_mask_cols, center_crop_size, com)
        for the active run. Tries run_summary.json first (manual
        training runs), then _train_kwargs.json (GUI-launched runs),
        then falls back to paper-recipe defaults.
        """
        mask_r, mask_cols, ccrop, com = 15, 45, 140, True
        try:
            rs = os.path.join(self.outdir or "", "run_summary.json")
            if os.path.exists(rs):
                with open(rs) as f:
                    js = json.load(f)
                cf = js.get("cfg", {})
                mask_r    = int(cf.get("center_mask_radius", mask_r))
                mask_cols = int(cf.get("polar_mask_cols", mask_cols))
                ccrop     = int(cf.get("center_crop_size", ccrop))
                if "com_centering" in cf:
                    com = bool(cf.get("com_centering", com))
        except Exception:
            pass
        try:
            tk = os.path.join(self.outdir or "", "_train_kwargs.json")
            if os.path.exists(tk):
                with open(tk) as f:
                    blob = json.load(f)
                # Try several common layouts the GUI uses.
                src = blob.get("kwargs", blob)
                if isinstance(src, dict):
                    mask_r    = int(src.get("center_mask_radius",
                                                 mask_r))
                    mask_cols = int(src.get("polar_mask_cols",
                                                 mask_cols))
                    ccrop     = int(src.get("center_crop_size",
                                                 ccrop))
                    if "com_centering" in src:
                        com = bool(src.get("com_centering", com))
        except Exception:
            pass
        return mask_r, mask_cols, ccrop, com

    def _resolve_cam_target(self, encoder, name: str):
        """Pick the conv module to hook for GradCAM.

        'last' (or empty) → final stage of the trained encoder
        (works for any n_layers∈{1..4}). 'layerN' → the named ResNet
        stage if it exists in this encoder.

        Encoder children are flat: [conv1, bn1, relu, layer1, ..., layerN].
        """
        children = list(encoder.children())
        if not children:
            raise ValueError("encoder has no children to hook")
        n = (name or "last").strip().lower()
        if n in ("", "last"):
            return children[-1]
        if n.startswith("layer"):
            try:
                ln = int(n.replace("layer", ""))
            except ValueError:
                raise ValueError(f"bad cam layer name {name!r}")
            # layer1..N live at indices 3..3+N-1 in our encoder.
            for idx in range(3, len(children)):
                # Match by class name to be robust against monkey-patching.
                child = children[idx]
                if child is getattr(encoder, "_layer_aliases", {}) \
                          .get(n):
                    return child
            # Fallback: count residual stages from index 3.
            stage_idx = 3 + (ln - 1)
            if 0 <= stage_idx < len(children):
                return children[stage_idx]
            raise ValueError(
                f"requested {name!r} but encoder only has "
                f"{max(0, len(children) - 3)} residual stages "
                f"(was the run trained with n_layers ≥ {ln}?)")
        raise ValueError(f"unknown cam layer {name!r}")

    def _dense_to_proto_ids(self, model, cfg, polar_pre, device):
        """dense class id -> original prototype index for the loaded run.
        GradCAM/IG target the 60-way prototype head, so a dense id (0..K-1)
        must be mapped to the actual prototype index or the attribution points
        at an inactive/wrong prototype. Cached per run dir; recovers + caches
        into inference.npz for runs that predate K_original_ids being saved."""
        from viz_gradcam import (load_original_prototype_ids,
                                  resolve_prototype_ids)
        run_dir = self.outdir
        cache = getattr(self, "_proto_ids_cache", None)
        if cache is None:
            cache = self._proto_ids_cache = {}
        if run_dir in cache:
            return cache[run_dir]
        ids = load_original_prototype_ids(run_dir)
        if ids is None:
            inf = getattr(self, "_inf", None)
            ki = inf.get("K_original_ids") if isinstance(inf, dict) else None
            if ki is not None and len(ki):
                ids = [int(x) for x in np.asarray(ki).ravel().tolist()]
        if ids is None:
            try:
                from data import LoadPRZ
                ds_full = LoadPRZ(cfg["path"], resize=192, vmax=cfg["vmax"])
                ids = resolve_prototype_ids(run_dir, model, ds_full, device,
                                            polar_pre=polar_pre)
            except Exception as e:
                print(f"[gradcam] prototype-id resolve failed: {e!r}")
        cache[run_dir] = ids
        return ids

    def _compute_gradcam_and_ig_from_raw(
            self, ckpt_path: str, c: int,
            raw_pattern_2d: np.ndarray
            ) -> tuple[np.ndarray, np.ndarray, np.ndarray]:
        """GradCAM **and** Integrated Gradients on a custom raw 2D
        pattern (e.g. grain average).  One forward + one polar
        preproc shared between the two attributions.  Returns
        ``(avg_cart_192, gradcam_cart_192, ig_cart_192)``."""
        import torch
        import torch.nn.functional as F
        from scipy.ndimage import gaussian_filter
        from dino_sr_contrastive_model import load_contrastive_checkpoint
        from viz_gradcam import (GradCAM, integrated_gradients,
                                   polar_cam_to_cartesian, dense_target)
        from viz_paper_attribution import (build_polar_preproc,
                                              build_cart_preproc)
        device = torch.device("cuda" if torch.cuda.is_available()
                                  else "cpu")
        cfg = SAMPLES[self.sample]
        mr, mask_cols, ccrop, _com = self._read_polar_pipeline_cfg()
        # Mirror infer_scan: the polar pipeline masks the beam AFTER the
        # warp (PolarMaskLeft) and applies no Cartesian CenterMask, but it
        # DOES honour COM centering.
        polar_pre = build_polar_preproc(polar_size=192,
                                          polar_mask_cols=mask_cols,
                                          center_crop_size=ccrop,
                                          center_mask_radius=0,
                                          com_centering=bool(_com))
        cart_pre  = build_cart_preproc(polar_size=192,
                                          center_crop_size=ccrop,
                                          center_mask_radius=mr)
        # Run the raw pattern through the DATASET's own preprocessing
        # (resize → rescale_like_vmax → ellipticity → blur → log-stretch)
        # so the attribution is computed on exactly the image the model was
        # fed.  Previously this clipped by vmax BEFORE resizing and skipped
        # blur/log/ellipticity entirely.
        x_full = self._model_frame_from_raw(raw_pattern_2d, cfg, device)
        x_cart  = cart_pre(x_full)
        x_polar = polar_pre(x_full)

        model, _, _, _ = load_contrastive_checkpoint(ckpt_path,
                                                         device=device)
        for p in model.student_encoder.parameters():   p.requires_grad_(True)
        for p in model.student_projector.parameters(): p.requires_grad_(True)
        for p in model.prototypes.parameters():        p.requires_grad_(True)
        model.eval()
        cam_layer_name = "last"
        try:
            cam_layer_name = self._cam_layer_var.get()
        except Exception:
            pass
        last_mod = self._resolve_cam_target(model.student_encoder,
                                                cam_layer_name)
        cam_tool = GradCAM(model, last_mod)
        c_t = dense_target(self._dense_to_proto_ids(model, cfg, polar_pre, device), c)
        with torch.enable_grad():
            xp_cam = x_polar.detach().requires_grad_(True)
            cam_p  = cam_tool(xp_cam, target_class=c_t)
        cam_cart = polar_cam_to_cartesian(cam_p).detach().cpu().numpy()
        cam_cart = gaussian_filter(np.abs(cam_cart), sigma=2.0)
        # IG (uses a fresh input tensor so the GradCAM autograd graph
        # doesn't interfere with the IG accumulation).
        with torch.enable_grad():
            ig_p = integrated_gradients(model, x_polar.detach(),
                                          target_class=c_t, n_steps=50)
        ig_cart = polar_cam_to_cartesian(ig_p).detach().cpu().numpy()
        ig_cart = gaussian_filter(np.abs(ig_cart), sigma=2.0)
        avg_cart = x_cart[0, 0].detach().cpu().numpy()
        return avg_cart, cam_cart, ig_cart

    def _compute_gradcam_from_raw(self, ckpt_path: str, c: int,
                                    raw_pattern_2d: np.ndarray
                                    ) -> tuple[np.ndarray, np.ndarray]:
        """GradCAM on a custom raw 2D pattern (e.g. grain average).
        Returns (avg_cart_192, gradcam_cart_192)."""
        import torch
        import torch.nn.functional as F
        from scipy.ndimage import gaussian_filter
        from dino_sr_contrastive_model import load_contrastive_checkpoint
        from viz_gradcam import GradCAM, polar_cam_to_cartesian, dense_target
        from viz_paper_attribution import (build_polar_preproc,
                                              build_cart_preproc)
        device = torch.device("cuda" if torch.cuda.is_available()
                                  else "cpu")
        cfg = SAMPLES[self.sample]
        mr, mask_cols, ccrop, _com = self._read_polar_pipeline_cfg()
        # Mirror infer_scan: the polar pipeline masks the beam AFTER the
        # warp (PolarMaskLeft) and applies no Cartesian CenterMask, but it
        # DOES honour COM centering.
        polar_pre = build_polar_preproc(polar_size=192,
                                          polar_mask_cols=mask_cols,
                                          center_crop_size=ccrop,
                                          center_mask_radius=0,
                                          com_centering=bool(_com))
        cart_pre  = build_cart_preproc(polar_size=192,
                                          center_crop_size=ccrop,
                                          center_mask_radius=mr)

        # Run the raw pattern through the DATASET's own preprocessing
        # (resize → rescale_like_vmax → ellipticity → blur → log-stretch)
        # so the attribution is computed on exactly the image the model was
        # fed.  Previously this clipped by vmax BEFORE resizing and skipped
        # blur/log/ellipticity entirely.
        x_full = self._model_frame_from_raw(raw_pattern_2d, cfg, device)
        x_cart  = cart_pre(x_full)
        x_polar = polar_pre(x_full)

        model, _, _, _ = load_contrastive_checkpoint(ckpt_path,
                                                         device=device)
        for p in model.student_encoder.parameters():   p.requires_grad_(True)
        for p in model.student_projector.parameters(): p.requires_grad_(True)
        for p in model.prototypes.parameters():        p.requires_grad_(True)
        model.eval()
        cam_layer_name = "last"
        try:
            cam_layer_name = self._cam_layer_var.get()
        except Exception:
            pass
        last_mod = self._resolve_cam_target(model.student_encoder,
                                                cam_layer_name)
        cam_tool = GradCAM(model, last_mod)
        c_t = dense_target(self._dense_to_proto_ids(model, cfg, polar_pre, device), c)
        with torch.enable_grad():
            xp_cam = x_polar.detach().requires_grad_(True)
            cam_p  = cam_tool(xp_cam, target_class=c_t)
        cam_cart = polar_cam_to_cartesian(cam_p).detach().cpu().numpy()
        cam_cart = gaussian_filter(np.abs(cam_cart), sigma=2.0)
        avg_cart = x_cart[0, 0].detach().cpu().numpy()
        return avg_cart, cam_cart

    def _model_frame_from_raw(self, raw2d, cfg, device):
        """Raw detector pattern -> the exact (1,1,192,192) tensor the model
        is fed, by delegating to LoadPRZ.preprocess_raw (resize → vmax
        rescale → ellipticity → blur → log-stretch).  Falls back to a plain
        vmax clip + resize if the cube can't be opened."""
        import torch
        import torch.nn.functional as F
        arr = np.asarray(raw2d, dtype=np.float32)
        try:
            ds = LoadPRZ(cfg["path"], resize=192, vmax=cfg["vmax"])
            img = ds.preprocess_raw(arr)
            return (torch.from_numpy(np.ascontiguousarray(img))
                      .unsqueeze(0).unsqueeze(0).to(device).float())
        except Exception as e:
            print(f"[posthoc] preprocess_raw unavailable ({e!r}); "
                  f"falling back to vmax-clip", flush=True)
            wn = np.clip(arr / max(float(cfg.get("vmax", 2.0)), 1e-6), 0.0, 1.0)
            x = (torch.from_numpy(wn).unsqueeze(0).unsqueeze(0)
                   .to(device).float())
            return F.interpolate(x, size=(192, 192), mode="bilinear",
                                   align_corners=False)

    def _compute_attribution_pair(self, ckpt_path, c, source, target_idx):
        """Run BOTH GradCAM and IG on the same input. Returns
        (raw_full_HxW, avg_cart_192, gradcam_cart_192, ig_cart_192).

        ``raw_full`` is the normalised (vmax-clipped, blur/log applied)
        pattern at the original detector resolution — so callers can
        display the CAM/IG overlay on the same field-of-view as the rest
        of the GUI rather than the cart-cropped 192×192 model input.
        """
        import torch
        import torch.nn.functional as F
        from scipy.ndimage import gaussian_filter
        from dino_sr_contrastive_model import load_contrastive_checkpoint
        from viz_gradcam import (GradCAM, integrated_gradients,
                                   polar_cam_to_cartesian, dense_target)
        from viz_paper_attribution import (build_polar_preproc,
                                             build_cart_preproc)

        device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        cfg = SAMPLES[self.sample]
        ds = LoadPRZ(cfg["path"], resize=192, vmax=cfg["vmax"])
        mask_r, mask_cols, ccrop, _com = \
            self._read_polar_pipeline_cfg()
        # The model's polar input carries NO Cartesian CenterMask (the beam
        # is masked post-polar by PolarMaskLeft) but DOES honour COM
        # centering -- mirror infer_scan exactly.  cart_pre is display-only,
        # so it keeps the beam mask for a readable underlay.
        polar_pre = build_polar_preproc(polar_size=192,
                                          polar_mask_cols=mask_cols,
                                          center_crop_size=ccrop,
                                          center_mask_radius=0,
                                          com_centering=bool(_com))
        cart_pre  = build_cart_preproc(polar_size=192,
                                          center_crop_size=ccrop,
                                          center_mask_radius=mask_r)

        # ---- pick the input pattern ----
        # Use the DATASET's own frames (ds[i]), not ds.get_raw(): __getitem__
        # applies resize -> rescale_like_vmax -> ellipticity -> blur ->
        # log-stretch, i.e. exactly what the model was fed.  The old code
        # took raw counts, clipped by vmax BEFORE resizing and skipped
        # blur/log/ellipticity, so attributions were computed on a subtly
        # different image than the one that produced the class assignment.
        if source == "class":
            soft = self._inf["soft_probs"]; ass = self._inf["assigns"]
            idx = np.where(ass == c)[0]
            if idx.size == 0:
                raise RuntimeError(f"class {c} is empty")
            s = soft[idx, c]
            order = np.argsort(-s)[:min(200, len(idx))]
            top = idx[order]
            frames = torch.stack([ds[int(i)] for i in top]).float()  # (N,1,192,192)
            w = torch.from_numpy(
                np.ascontiguousarray(s[order], dtype=np.float32))
            x_full = ((frames * w[:, None, None, None]).sum(0)
                        / (float(w.sum()) + 1e-12)).unsqueeze(0)
        else:
            x_full = ds[int(target_idx)].unsqueeze(0).float()
        x_full = x_full.to(device).float()
        x_cart  = cart_pre(x_full)
        x_polar = polar_pre(x_full)

        model, _, _, _ = load_contrastive_checkpoint(ckpt_path, device=device)
        for p in model.student_encoder.parameters():  p.requires_grad_(True)
        for p in model.student_projector.parameters(): p.requires_grad_(True)
        for p in model.prototypes.parameters():        p.requires_grad_(True)
        model.eval()
        cam_layer_name = "last"
        try:
            cam_layer_name = self._cam_layer_var.get()
        except Exception:
            pass
        last_mod = self._resolve_cam_target(model.student_encoder,
                                                cam_layer_name)
        cam_tool = GradCAM(model, last_mod)
        c_t = dense_target(self._dense_to_proto_ids(model, cfg, polar_pre, device), c)

        with torch.enable_grad():
            xp_cam = x_polar.detach().requires_grad_(True)
            cam_p  = cam_tool(xp_cam, target_class=c_t)
        cam_cart = polar_cam_to_cartesian(cam_p).detach().cpu().numpy()
        cam_cart = gaussian_filter(np.abs(cam_cart), sigma=2.0)

        # IG uses its own clean input tensor
        with torch.enable_grad():
            ig_p = integrated_gradients(model, x_polar.detach(),
                                          target_class=c_t, n_steps=50)
        ig_cart = polar_cam_to_cartesian(ig_p).detach().cpu().numpy()
        ig_cart = gaussian_filter(np.abs(ig_cart), sigma=2.0)

        avg_cart = x_cart[0, 0].detach().cpu().numpy()
        # Full (un-cropped) model-view frame, i.e. the 192x192 image before
        # CenterCrop/Resize.  Kept for callers that want the wider field of
        # view; _draw_triptych deliberately ignores it and shows everything
        # at the cropped scale the attribution was computed on.
        raw_full = x_full[0, 0].detach().cpu().numpy().astype(np.float32)
        return raw_full, avg_cart, cam_cart, ig_cart

    def _draw_triptych(self, c, source, target_idx, tag, avg, cam, ig,
                          raw_full=None):
        """Render a 1×3 panel: [class avg / frame] [GradCAM overlay]
        [IG overlay].

        All three panels are at the 192×192 cart-cropped scale — i.e.
        the EXACT input the model saw (CenterCrop(ccrop) → Resize(192)
        → optional CenterMask).  GradCAM/IG are computed in the polar
        frame and back-projected onto this same 192-cart grid, so this
        is the only spatially-faithful view.  ``raw_full`` is ignored
        for the display but kept in the signature for compatibility.
        """
        # Display beam-mask radius — matches training CenterMask.
        try:
            from data import SAMPLES as _S
            bm_r_192 = int((_S.get(self.sample) or {}).get(
                "center_mask_radius", 0))
        except Exception:
            bm_r_192 = 0
        H = avg.shape[0]
        if bm_r_192 > 0:
            yy, xx = np.ogrid[:H, :H]; cy = cx = H / 2.0
            bm = ((yy - cy) ** 2 + (xx - cx) ** 2) > bm_r_192 ** 2
        else:
            bm = np.ones((H, H), dtype=bool)
        disp = avg * bm

        def _norm(a):
            a = a - a.min(); mx = a.max()
            return (a / mx if mx > 0 else a) * bm

        cam_n = _norm(cam); ig_n = _norm(ig)
        fig = self._new_fig()
        axes = [fig.add_subplot(1, 3, i + 1) for i in range(3)]
        axes[0].imshow(disp, cmap="inferno", aspect="equal",
                          interpolation="nearest")
        axes[0].set_title(
            ("class avg" if source == "class" else "frame")
            + " (model view)", fontsize=10)
        axes[1].imshow(disp, cmap="gray", aspect="equal",
                          interpolation="nearest")
        axes[1].imshow(cam_n, cmap="jet", alpha=0.55,
                          aspect="equal", interpolation="nearest")
        axes[1].set_title("GradCAM", fontsize=10)
        axes[2].imshow(disp, cmap="gray", aspect="equal",
                          interpolation="nearest")
        axes[2].imshow(ig_n, cmap="jet", alpha=0.55,
                          aspect="equal", interpolation="nearest")
        axes[2].set_title("Integrated Gradients", fontsize=10)
        for a in axes:
            a.set_xticks([]); a.set_yticks([])

        # Reciprocal scale bar — cart pipeline scales 1 raw px to
        # (192/ccrop) display px, so use q_per_polar_bin to back this out.
        rp = self._recip_per_px()
        if rp > 0:
            from gui_app._calib_utils import (q_per_polar_bin,
                                                 get_raw_detector_size,
                                                 add_recip_scalebar)
            qpx = q_per_polar_bin(
                rp, get_raw_detector_size(self.sample))
            for a in axes:
                add_recip_scalebar(a, q_per_disp_px=qpx, length_q=0.2)

        cam_lab = ""
        try:
            cl = self._cam_layer_var.get()
            cam_lab = f"   [CAM: {cl}]" if cl else ""
        except Exception:
            pass
        fig.suptitle(f"{self.sample}  —  {tag}{cam_lab}", fontsize=11)
        self._redraw()

    # ---- legacy aliases (kept for any callers) ----
    def _render_one_class_avg(self):
        self._attr_source.set("class"); self._render_triptych()

    def _render_gradcam(self):
        self._attr_source.set("class"); self._render_triptych()

    def _render_ig(self):
        self._attr_source.set("class"); self._render_triptych()

    # ===================================================================
    # Phase B — active labelling + fine-tune
    # ===================================================================
    def _open_active_labeler(self):
        """Open the pair-labelling dialog in 'active' mode with the
        proposer the user picked from the dropdown."""
        if not self._ensure_inference(): return
        from gui_app.pair_labeler import PairLabelerWindow
        from gui_app.pair_proposers import (
            cross_class_proposer, intra_class_proposer,
            mixed_intra_inter_proposer, scan_edge_proposer,
            low_margin_proposer)

        soft   = self._inf["soft_probs"]
        ass    = self._inf["assigns"]
        K      = int(soft.shape[1])
        Ny, Nx = self._scan_shape

        # Parse the optional class filters from the small entries.
        # Blank → None → no restriction (current default behaviour).
        def _parse_int(s):
            s = (s or "").strip()
            if not s: return None
            try: v = int(s)
            except ValueError: return None
            if not (0 <= v < K):
                raise ValueError(f"class id {v} out of range for K={K}")
            return v
        try:
            r_intra = _parse_int(self._ft_intra_class.get())
            ia      = _parse_int(self._ft_inter_a.get())
            ib      = _parse_int(self._ft_inter_b.get())
        except ValueError as e:
            messagebox.showerror("class filter", str(e))
            return
        r_inter = (ia, ib) if (ia is not None and ib is not None) else None
        if (ia is None) != (ib is None):
            messagebox.showinfo("inter filter",
                "Set BOTH inter entries to a class id, or leave both "
                "blank. (One filled and one blank is treated as no "
                "filter.)")
            r_inter = None

        kind = self._ft_proposer_var.get()
        try:
            if kind == "cross+intra":
                proposer = mixed_intra_inter_proposer(
                    soft, K,
                    restrict_intra_class=r_intra,
                    restrict_inter_pair=r_inter)
                mode = "active"
            elif kind == "cross_class":
                proposer = cross_class_proposer(
                    soft, K, restrict_pair=r_inter)
                mode = "active"
            elif kind == "intra_class":
                proposer = intra_class_proposer(
                    soft, K, restrict_class=r_intra)
                mode = "active"
            elif kind == "scan_edge":
                proposer = scan_edge_proposer(ass, (Ny, Nx))
                mode = "active"
            elif kind == "low_margin":
                proposer = low_margin_proposer(soft)
                mode = "active"
            elif kind == "random":
                proposer = None
                mode = "random"
            else:
                messagebox.showerror("source", f"unknown source: {kind}")
                return
        except Exception as e:
            messagebox.showerror("proposer build failed", repr(e))
            return

        # Resolve cube path + load mmap so the labeller can fetch raw
        # patterns. Reuse the loaded path stored on the panel.
        cube_path = self._cube_path
        if not cube_path or not os.path.exists(cube_path):
            messagebox.showerror("cube",
                "Could not find this run's cube file. Make sure the "
                "Pre-processing tab has the same sample loaded.")
            return
        try:
            from gui_app.pre_panel import _open_lazy
            cube = _open_lazy(cube_path)
        except Exception as e:
            messagebox.showerror("cube open failed", repr(e)); return

        cfg = SAMPLES.get(self.sample, {})
        vmax = float(cfg.get("vmax", 2.0))

        if proposer is not None and getattr(proposer, "total", 0) == 0:
            messagebox.showinfo("active labeller",
                f"The '{kind}' proposer found 0 candidate pairs for "
                f"this run.  (e.g. scan_edge needs >0 pixels assigned "
                f"to differing classes.)  Try a different source.")
            return

        sample_key = self.sample or os.path.basename(cube_path)

        def _after_close(counts):
            self._set_status(
                f"active labelling done.  total pairs on disk: "
                f"{counts.get('total', 0)}.")
            self._refresh_ft_status()

        try:
            win = PairLabelerWindow(
                self,
                cube_path=cube_path,
                sample_key=sample_key,
                cube=cube, vmax=vmax, cmap="inferno",
                mode=mode,
                pair_proposer=proposer,
                on_close=_after_close,
            )
            win.transient(self); win.focus_set()
        except Exception as e:
            messagebox.showerror("active labeller failed", repr(e))

    def _refresh_ft_status(self):
        """Refresh the small status line in the Fine-tune section."""
        if not getattr(self, "_ft_status", None): return
        try:
            cube_path = self._cube_path
            if not cube_path:
                self._ft_status.configure(text="(no cube linked)")
                return
            from gui_app.pair_labels import (
                load_pair_labels, label_count, label_path_for_cube)
            d = load_pair_labels(cube_path)
            c = label_count(d)
            sp = label_path_for_cube(cube_path)
            if c["total"] == 0:
                self._ft_status.configure(
                    text=f"no labels yet @ {os.path.basename(sp)}")
            else:
                self._ft_status.configure(
                    text=f"{c['total']} pairs available  "
                          f"({c['same']} same / {c['diff']} diff)\n"
                          f"@ {os.path.basename(sp)}")
        except Exception as e:
            self._ft_status.configure(text=f"(label-info error: {e})")

    def _start_finetune(self):
        """Spawn a fine-tune subprocess that warm-starts from this
        run's `best.pth`, freezes the encoder, restricts the loader to
        a `subsample_n` random subset, and adds high `λ_pair` so the
        labelled pairs drive the gradient."""
        if self.outdir is None or self.sample is None:
            messagebox.showinfo("fine-tune",
                "Load a run first."); return

        # Find the parent checkpoint
        ckpt = self._best_ckpt()
        if ckpt is None:
            messagebox.showerror("checkpoint",
                f"No best.pth or ckpt_ep*.pth in {self.outdir}")
            return

        # Find the parent run_summary.json so we can read the trained cfg
        parent_summary_path = os.path.join(self.outdir,
                                              "run_summary.json")
        if not os.path.exists(parent_summary_path):
            messagebox.showerror("run summary",
                f"No run_summary.json in {self.outdir} — cannot "
                f"reconstruct the training kwargs."); return
        try:
            parent = json.load(open(parent_summary_path,
                                       encoding="utf-8"))
        except Exception as e:
            messagebox.showerror("run summary unreadable", repr(e))
            return
        pcfg = parent.get("cfg", {})

        # Pair labels
        cube_path = self._cube_path
        from gui_app.pair_labels import (
            label_path_for_cube, load_pair_labels, label_count)
        pair_path = label_path_for_cube(cube_path) if cube_path else None
        if not pair_path or not os.path.exists(pair_path):
            messagebox.showerror("no labels",
                "Click 'Active-label pairs…' first to collect pair "
                "labels for this sample."); return
        labels_obj = load_pair_labels(cube_path)
        n_lbl = label_count(labels_obj)["total"]
        if n_lbl == 0:
            messagebox.showerror("no labels",
                "The labels file exists but has 0 pairs."); return

        # Compute pre / post breakdown using the parent run's
        # timestamp as the cutoff. Used both in the confirmation
        # modal AND (if the checkbox is on) as the
        # `pair_labels_min_timestamp` knob in the spec.
        parent_ts = parent.get("timestamp")
        only_post = bool(self._ft_only_post.get())
        n_pre = n_post = 0
        if parent_ts:
            for p in labels_obj.get("pairs", []):
                t = str(p.get("t", ""))
                if t > str(parent_ts):
                    n_post += 1
                else:
                    n_pre += 1
        else:
            n_pre = n_lbl
        n_use = n_post if only_post else n_lbl
        if only_post and n_post == 0:
            messagebox.showinfo("no post-parent labels",
                f"The 'only labels added since parent run' option is "
                f"checked, but every label in the sidecar predates the "
                f"parent run's timestamp ({parent_ts}). "
                f"Either uncheck the option or label a few new pairs "
                f"first."); return

        # Build the fine-tune outdir under the parent run
        from datetime import datetime as _dt
        stamp = _dt.now().strftime("%Y%m%d_%H%M%S")
        ft_outdir = os.path.join(
            self.outdir, f"finetune_{n_lbl}pairs_{stamp}")
        os.makedirs(ft_outdir, exist_ok=True)

        # Confirm
        epochs   = max(1, int(self._ft_epochs.get()))
        lam_pair = float(self._ft_lambda_pair.get())
        lam_c1d  = float(self._ft_lambda_c1d.get())
        sub_n    = max(0, int(self._ft_subsample_n.get()))
        freeze   = bool(self._ft_freeze_enc.get())
        # Build labels-policy line for the confirmation modal.
        if parent_ts:
            label_line = (
                f"Labels : {n_lbl} pairs total  "
                f"(pre-parent: {n_pre}, since-parent: {n_post})\n"
                f"         policy: "
                + ("ONLY post-parent  →  using "
                    f"{n_post}" if only_post
                   else f"CUMULATIVE  →  using {n_lbl}")
                + f"\n"
            )
        else:
            label_line = (f"Labels : {n_lbl} pairs  "
                          f"(parent has no timestamp)\n")
        if not messagebox.askyesno("Fine-tune?",
            f"Sample : {self.sample}\n"
            f"Parent : {os.path.basename(self.outdir)}\n"
            f"Ckpt   : {os.path.basename(ckpt)}\n\n"
            + label_line +
            f"Epochs : {epochs}\n"
            f"λ_pair : {lam_pair}\n"
            f"λ_1d   : {lam_c1d}\n"
            f"subset : {sub_n}\n"
            f"freeze : {freeze}\n\n"
            f"Output → {ft_outdir}\n\n"
            f"Start fine-tune?"):
            return

        # Build the spec — copy the parent cfg and override the
        # fine-tune knobs so the run reproduces the same data
        # pipeline.
        ft_kwargs = self._build_finetune_kwargs(
            pcfg=pcfg, ckpt=ckpt, pair_path=pair_path,
            epochs=epochs, lam_pair=lam_pair,
            lam_c1d=lam_c1d, sub_n=sub_n, freeze=freeze,
            pair_min_timestamp=(parent_ts if only_post else None))

        # Spawn via the same TrainingJob plumbing the Train tab uses.
        try:
            from gui_app.runner import TrainingJob, read_training_log
        except Exception as e:
            messagebox.showerror("runner import failed", repr(e))
            return

        self._ft_job = TrainingJob(self.sample, ft_outdir, ft_kwargs)
        self._ft_job.start()
        self._ft_status.configure(
            text=f"fine-tune RUNNING…  outdir = "
                  f"{os.path.relpath(ft_outdir, self.outdir)}")
        self._poll_finetune()

    def _build_finetune_kwargs(self, *, pcfg, ckpt, pair_path,
                                  epochs, lam_pair, lam_c1d,
                                  sub_n, freeze,
                                  pair_min_timestamp=None):
        """Construct run_config kwargs for a fine-tune run, inheriting
        the parent's pipeline + augmentation + COM choices so the
        forward pass is identical to what the parent saw."""
        return dict(
            epochs=int(epochs),
            seed=int(pcfg.get("seed", 42)),
            batch_size=int(pcfg.get("batch_size", 128)),
            lr=float(pcfg.get("lr", 3e-4)) * 0.5,    # gentler for FT
            weight_decay=float(pcfg.get("weight_decay", 1e-6)),
            num_prototypes=int(pcfg.get("num_prototypes", 6)),
            t0=float(pcfg.get("T0", 0.04)),
            tfin=float(pcfg.get("Tfin", 0.07)),
            warmup_epochs=int(pcfg.get("warmup_epochs", 20)),
            ramp_epochs=int(pcfg.get("ramp_epochs", 10)),
            entropy_gate=False,
            projection_dim=int(pcfg.get("projection_dim", 128)),
            projection_hidden=int(pcfg.get("projection_hidden", 256)),
            theta_shift_range=None,
            theta_shift_range_student=int(
                pcfg.get("theta_shift_range_student", 192)),
            theta_shift_range_teacher=int(
                pcfg.get("theta_shift_range_teacher", 16)),
            center_mask_radius=int(
                pcfg.get("center_mask_radius", 15)),
            center_crop_size=int(pcfg.get("center_crop_size", 140)),
            vmax=None,
            polar_size=int(pcfg.get("polar_size", 192)),
            polar_mask_cols=int(pcfg.get("polar_mask_cols", 45)),
            pipeline=str(pcfg.get("pipeline", "polar")),
            centroid_lambda=0.0,
            centroid_margin=float(pcfg.get("centroid_margin", 0.3)),
            conf_weight_gamma=float(
                pcfg.get("conf_weight_gamma", 0.0)),
            entropy_gate_override=None,
            lam_spatial=0.0,
            architecture=str(pcfg.get("architecture", "resnet")),
            n_layers=int(pcfg.get("n_layers", 1)),
            w_ent=float(pcfg.get("w_ent", 0.0)),
            com_centering=bool(pcfg.get("com_centering", True)),
            com_search_radius_factor=float(
                pcfg.get("com_search_radius_factor", 2.0)),
            aug_disable=list(
                pcfg.get("aug_disable",
                          ["hflip", "vflip", "colorjitter"])),
            supcon_radials_path=pcfg.get("supcon_radials_path"),
            supcon_thresholds_path=pcfg.get(
                "supcon_thresholds_path"),
            supcon_lambda=0.0,
            supcon_temperature=float(
                pcfg.get("supcon_temperature", 0.3)),
            contrastive_lambda_override=0.0,
            proto_repel_lambda=0.0,
            proto_repel_threshold=0.5,
            cluster1d_lambda=float(lam_c1d),
            cluster1d_margin=float(
                pcfg.get("cluster1d_margin", 0.4)),
            cluster1d_min_cluster_mass=1.0,
            cluster1d_warmup_frac=0.0,
            cluster1d_ramp_frac=0.0,
            pair_labels_path=pair_path,
            pair_labels_min_timestamp=pair_min_timestamp,
            lambda_pair=float(lam_pair),
            pair_entropy_reg=0.0,
            pair_per_batch=32,
            init_from_checkpoint=ckpt,
            freeze_encoder=bool(freeze),
            subsample_n=(int(sub_n) if sub_n else None),
            save_every=max(1, epochs),
        )

    def _poll_finetune(self):
        """Poll the fine-tune job every second and refresh status. On
        completion, re-link the panel to the fine-tune outdir so the
        class map etc. show the new model."""
        job = getattr(self, "_ft_job", None)
        if job is None: return
        try:
            from gui_app.runner import read_training_log
            rows = read_training_log(job.csv_path)
            n = len(rows)
        except Exception:
            n = 0
        if job.is_running():
            elapsed = job.elapsed()
            tot = job.kwargs.get("epochs", 0)
            self._ft_status.configure(
                text=f"fine-tune RUNNING…  ep {n}/{tot}  "
                      f"({elapsed:.0f}s)\noutdir = "
                      f"{os.path.relpath(job.outdir, self.outdir)}")
            self.after(1000, self._poll_finetune)
            return
        # Finished
        status = job.status().upper()
        if status == "DONE":
            self._ft_status.configure(
                text=f"fine-tune DONE  ({job.elapsed():.0f}s).  "
                      f"re-linking class map to fine-tune outdir.")
            # Re-link the Post-hoc panel to the fine-tune output so
            # subsequent renders use the fine-tuned model.
            try:
                self.link_run(job.outdir, self.sample)
                self._render_classmap()
            except Exception as e:
                messagebox.showerror("re-link failed", repr(e))
        else:
            err = job.error() or "(no error message)"
            self._ft_status.configure(
                text=f"fine-tune {status} after "
                      f"{job.elapsed():.0f}s")
            if status == "FAILED":
                messagebox.showerror("Fine-tune failed", err[:1000])
        self._ft_job = None
