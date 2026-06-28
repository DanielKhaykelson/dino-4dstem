"""nmf_panel.py -- NMF + clustering tab.

Independent of the trained DINO model.  Used as a baseline /
benchmark pipeline:

    polar-pre-process the cube  →  NMF (n_comp components)
                                →  clustering (K-means / Aglo / HDBSCAN)
                                →  per-pattern class label, reshape to
                                   the scan grid for a class map.

NMF "variant" dropdown decides what the input matrix to NMF looks like:

    Polar              (θ, r), masked beam, default
    Polar + θ-shift    augmented with θ-rolled copies → rotation invariant
    Polar + log        log1p(I) before NMF (handles dynamic range)
    Polar + sparse     l1-regularised H (sparser activations)
    Cartesian (flat)   raw (cropped + resized) 2D patterns
    1D radial          ∫_θ I(θ, r) — strict 1D baseline (matches K-means
                        on radial).

Three clustering methods:

    K-means       (default — Euclidean over W)
    Aglo          (Euclidean / Cosine / polar-xcorr distance,
                    polar-xcorr only available when input is polar)
    HDBSCAN       (no K, uses min_cluster_size; reports the cluster count
                    it discovered)

Auto modes:

    n_comp:  knee of reconstruction error vs n_comp
    K:       best silhouette over K∈[2..12] (HDBSCAN ignores)

"Report All": sweeps every NMF variant × every clustering method and
saves a folder of publication-ready PNGs + a PPTX.
"""
from __future__ import annotations
import os, json, time, threading, webbrowser

import numpy as np
import customtkinter as ctk
import tkinter as tk
from tkinter import filedialog, messagebox

import matplotlib
matplotlib.use("TkAgg")
import matplotlib.pyplot as plt
from matplotlib.figure import Figure
from matplotlib.backends.backend_tkagg import (FigureCanvasTkAgg,
                                                 NavigationToolbar2Tk)
from matplotlib.colors import ListedColormap


# ---------------------------------------------------------------------------
# Each variant carries a verified citation (ref + url). DOIs checked live
# against OpenAlex/Crossref (2026-06). NOTE: the earlier "Wang 2024/2025"
# attributions were INCORRECT — the real polar-NMF-for-4D-STEM references are
# Uesugi et al. 2020 (Ultramicroscopy, foundational) and Kimoto et al. 2025
# (Sci. Rep., polar + domain constraints / log handling).
NMF_VARIANTS = {
    "Polar  (Uesugi et al. 2020)":
        dict(input="polar", log=False, sparse=False, theta_shift=False,
             ref="Uesugi, Koshiya, Kikkawa, Nagai et al., "
                 "Ultramicroscopy 218, 113168 (2020)",
             url="https://doi.org/10.1016/j.ultramic.2020.113168"),
    "Polar + θ-shift  (Krajnak & Etheridge 2020)":
        dict(input="polar", log=False, sparse=False, theta_shift=True,
             ref="Krajňák & Etheridge, PNAS 117(45), 27805–27810 (2020)",
             url="https://doi.org/10.1073/pnas.2006975117"),
    "Polar + log  (Kimoto et al. 2025)":
        dict(input="polar", log=True,  sparse=False, theta_shift=False,
             ref="Kimoto, Uesugi, Harano, Kikkawa, "
                 "Scientific Reports 15 (2025)",
             url="https://doi.org/10.1038/s41598-025-23541-7"),
    "Polar + sparse  (Hoyer 2004)":
        dict(input="polar", log=False, sparse=True,  theta_shift=False,
             ref="Hoyer, J. Mach. Learn. Res. 5, 1457–1469 (2004)",
             url="https://www.jmlr.org/papers/v5/hoyer04a.html"),
    "Cartesian flat  (Spurgeon et al. 2020)":
        dict(input="cart",  log=False, sparse=False, theta_shift=False,
             ref="Spurgeon et al., Nature Materials 20, 274–279 (2020)",
             url="https://doi.org/10.1038/s41563-020-00833-z"),
    "1D radial  (baseline)":
        dict(input="radial", log=False, sparse=False, theta_shift=False,
             ref="Baseline (azimuthal average); cf. py4DSTEM — Savitzky "
                 "et al., Microsc. Microanal. 27, 712–743 (2021)",
             url="https://doi.org/10.1017/s1431927621000477"),
}


def _safe_name(s: str) -> str:
    """Slug-safe version of a variant key (parentheses, spaces, etc.)."""
    out = []
    for ch in s:
        if ch.isalnum():
            out.append(ch)
        elif ch in (" ", "-", "+"):
            out.append("_")
    return "".join(out).strip("_") or "variant"
CLUSTER_METHODS = ["K-means", "Aglo", "HDBSCAN", "FCM"]
AGLO_DISTANCES  = ["euclidean", "cosine", "polar-xcorr"]

CENTER_CROP = 140
POLAR_SIZE  = 192
POLAR_MASK_COLS = 45


# ---------------------------------------------------------------------------
def _section(parent, title):
    f = ctk.CTkFrame(parent, fg_color="transparent")
    f.pack(fill="x", padx=2, pady=(8, 0))
    ctk.CTkLabel(f, text=title, font=("Segoe UI", 11, "bold")).pack(anchor="w")
    sep = ctk.CTkFrame(parent, height=2, fg_color=("#cccccc", "#444444"))
    sep.pack(fill="x", padx=2, pady=(0, 4))
    return parent


# ---------------------------------------------------------------------------
# Build the (N, D) input matrix for NMF.
# ---------------------------------------------------------------------------
def build_nmf_input(sample_key: str, variant_cfg: dict,
                     progress_cb=None, batch: int = 128,
                     vmax_override: float | None = None
                     ) -> tuple[np.ndarray, np.ndarray | None,
                                  tuple, dict]:
    """Returns (X, X_aug, comp_shape, info).

    X       : (N, D) non-negative matrix to transform downstream.
    X_aug   : (N·k, D) augmented matrix used for FITTING (same shape D),
              or None for non-augmented variants. Only rows of X get a
              W; the augmented rows are discarded after fit.
    comp_shape : the (H, W) reshape for each NMF component for display
                  (theta, r) for polar, (H, W) for cart, or (1, n_bins)
                  for 1D radial.
    info    : dict with extras (mask, etc.).
    """
    import torch
    from data import SAMPLES, LoadPRZ
    from torchvision.transforms import v2 as T
    from torchvision.transforms import InterpolationMode
    from dino_sr_ablation import PolarTransform, PolarMaskLeft
    cfg = SAMPLES[sample_key]
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    vmax = float(vmax_override) if vmax_override and vmax_override > 0 \
            else float(cfg["vmax"])
    ds = LoadPRZ(cfg["path"], resize=POLAR_SIZE, vmax=vmax)
    N = len(ds)
    inp = variant_cfg["input"]
    if inp == "polar" or inp == "radial":
        pre = T.Compose([
            T.CenterCrop(CENTER_CROP),
            T.Resize(POLAR_SIZE,
                       interpolation=InterpolationMode.BILINEAR,
                       antialias=True),
            PolarTransform(output_size=POLAR_SIZE),
            PolarMaskLeft(k_cols=POLAR_MASK_COLS),
        ])
    elif inp == "cart":
        pre = T.Compose([
            T.CenterCrop(CENTER_CROP),
            T.Resize(POLAR_SIZE,
                       interpolation=InterpolationMode.BILINEAR,
                       antialias=True),
        ])
    else:
        raise ValueError(f"unknown input mode {inp!r}")

    chunks = []
    with torch.no_grad():
        for i in range(0, N, batch):
            j = min(i + batch, N)
            xs = torch.stack([ds[k] for k in range(i, j)]).to(device).float()
            xs = pre(xs).squeeze(1)             # (b, ?, ?)
            if inp == "radial":
                xs = xs.sum(dim=-2)              # (b, r) — sum over theta
            chunks.append(xs.cpu().numpy())
            if progress_cb is not None and (i // batch) % 4 == 0:
                progress_cb(j, N, "polar pipeline")

    if inp == "radial":
        X_full = np.concatenate(chunks, axis=0)              # (N, r)
        comp_shape = (1, X_full.shape[1])
    else:
        X_full = np.concatenate(chunks, axis=0)              # (N, h, w)
        comp_shape = X_full.shape[1:]
    if variant_cfg["log"]:
        X_full = np.log1p(np.clip(X_full, 0, None))
    X_full = np.clip(X_full, 0, None).astype(np.float32)

    # Augment for θ-shift variant: roll along theta axis.
    X_aug = None
    if variant_cfg["theta_shift"] and inp == "polar":
        n_shifts = 4
        shifts = np.linspace(0, POLAR_SIZE, n_shifts, endpoint=False
                                ).astype(int)
        rolls = []
        for s in shifts:
            rolls.append(np.roll(X_full, int(s), axis=1))
        X_aug = np.concatenate(rolls, axis=0).reshape(
            n_shifts * X_full.shape[0], -1).astype(np.float32)

    if inp == "radial":
        X = X_full
    else:
        X = X_full.reshape(X_full.shape[0], -1)
    info = dict(N=int(X.shape[0]), D=int(X.shape[1]),
                  variant=variant_cfg, vmax=vmax)
    return X, X_aug, comp_shape, info


# ---------------------------------------------------------------------------
# NMF + clustering helpers.
# ---------------------------------------------------------------------------
def fit_nmf(X: np.ndarray, X_aug, n_components: int,
             sparse: bool = False, max_iter: int = 200):
    """Fit (optionally on X_aug) and transform X. Returns (W, H, err)."""
    from sklearn.decomposition import NMF
    kw = dict(n_components=int(n_components), init="nndsvda",
                solver="cd", max_iter=int(max_iter), tol=1e-4,
                random_state=42)
    if sparse:
        kw["alpha_H"] = 0.1
        kw["l1_ratio"] = 1.0
    nmf = NMF(**kw)
    if X_aug is not None:
        nmf.fit(X_aug)
        W = nmf.transform(X)
        H = nmf.components_
        err = float(nmf.reconstruction_err_)
    else:
        W = nmf.fit_transform(X)
        H = nmf.components_
        err = float(nmf.reconstruction_err_)
    return W.astype(np.float32), H.astype(np.float32), err


def auto_n_comp(X, X_aug, range_=(2, 14), sparse=False, progress_cb=None
                ) -> tuple[int, np.ndarray]:
    """Knee of reconstruction-error curve."""
    ns, errs = [], []
    for n in range(range_[0], range_[1]):
        _, _, e = fit_nmf(X, X_aug, n, sparse=sparse, max_iter=120)
        ns.append(n); errs.append(e)
        if progress_cb is not None:
            progress_cb(n - range_[0] + 1, range_[1] - range_[0],
                          "scan n_comp")
    x = np.asarray(ns, float); y = np.asarray(errs, float)
    xn = (x - x.min()) / (x.max() - x.min() + 1e-12)
    yn = (y - y.min()) / (y.max() - y.min() + 1e-12)
    line = yn[0] + (yn[-1] - yn[0]) * xn
    knee = int(x[np.argmax(line - yn)])
    return knee, np.asarray(errs, float)


def fuzzy_cmeans(X: np.ndarray, k: int, m: float = 2.0,
                  max_iter: int = 200, tol: float = 1e-5,
                  seed: int = 42
                  ) -> tuple[np.ndarray, np.ndarray]:
    """Fuzzy c-means clustering on rows of X (Bezdek 1981).

    Returns (hard_labels: (N,), memberships: (N, k)). The hard labels
    are argmax over memberships; soft memberships are kept in case the
    caller wants entropy / overlap diagnostics.
    """
    rng = np.random.default_rng(seed)
    X = X.astype(np.float32)
    n = X.shape[0]
    if n == 0:
        return np.zeros(0, np.int32), np.zeros((0, k), np.float32)
    U = rng.uniform(size=(n, int(k))).astype(np.float32)
    U /= U.sum(axis=1, keepdims=True).clip(min=1e-12)
    exponent = -2.0 / (float(m) - 1.0)
    for _ in range(int(max_iter)):
        Um = U ** float(m)
        # Cluster centres = weighted means.
        C = (Um.T @ X) / Um.sum(axis=0)[:, None].clip(min=1e-12)
        # Distances (n, k) via |x - c|² = |x|² + |c|² - 2 x·c.
        x_sq = (X * X).sum(axis=1, keepdims=True)
        c_sq = (C * C).sum(axis=1)
        D2 = x_sq + c_sq[None, :] - 2.0 * (X @ C.T)
        D2 = np.clip(D2, 1e-12, None)
        D = np.sqrt(D2)
        # u_ik = D_ik^(-2/(m-1)) / Σ_j D_ij^(-2/(m-1))
        Dp = D ** exponent
        U_new = Dp / Dp.sum(axis=1, keepdims=True).clip(min=1e-12)
        if np.linalg.norm(U_new - U) < tol:
            U = U_new
            break
        U = U_new
    labels = U.argmax(axis=1).astype(np.int32)
    return labels, U.astype(np.float32)


def cluster_W(W, method, k=None, distance="euclidean",
                min_cluster_size=30, polar_full=None,
                fcm_m: float = 2.0):
    """Cluster the NMF weight matrix W.  Returns int labels (N,)."""
    from sklearn.cluster import KMeans, AgglomerativeClustering
    method = method.lower().replace("-", "")
    if method == "kmeans":
        km = KMeans(n_clusters=int(k), n_init=10, random_state=42)
        return km.fit_predict(W).astype(np.int32)
    if method == "fcm":
        labels, _ = fuzzy_cmeans(W, int(k), m=float(fcm_m))
        return labels
    if method == "aglo":
        if distance == "polar-xcorr":
            if polar_full is None:
                raise ValueError("polar-xcorr needs polar_full")
            D = _polar_xcorr_distance_matrix(polar_full)
            ag = AgglomerativeClustering(n_clusters=int(k),
                                            metric="precomputed",
                                            linkage="average")
            return ag.fit_predict(D).astype(np.int32)
        ag = AgglomerativeClustering(n_clusters=int(k), metric=distance,
                                        linkage="average")
        return ag.fit_predict(W).astype(np.int32)
    if method == "hdbscan":
        try:
            import hdbscan
            h = hdbscan.HDBSCAN(min_cluster_size=int(min_cluster_size))
            labels = h.fit_predict(W)
        except Exception:
            try:
                from sklearn.cluster import HDBSCAN as SKHDBSCAN
                h = SKHDBSCAN(min_cluster_size=int(min_cluster_size))
                labels = h.fit_predict(W)
            except Exception as e:
                raise RuntimeError(
                    f"HDBSCAN not available (pip install hdbscan): {e!r}")
        # Re-label noise (-1) to a real id at the end so downstream
        # bincount / cmap doesn't crash.
        labels = labels.copy()
        if (labels < 0).any():
            labels[labels < 0] = labels.max() + 1
        return labels.astype(np.int32)
    raise ValueError(f"unknown method {method!r}")


def _polar_xcorr_distance_matrix(polar_full: np.ndarray) -> np.ndarray:
    """Wang 2024-style angular cross-correlation distance.

    For each pair (i, j), distance = 1 - max_θ' ⟨I_i(θ), I_j(θ + θ')⟩
    over the angular axis. Implemented via FFT for speed:
        max-θ-corr(i, j) = max_τ Re[ifft(F_i · conj(F_j)) ](τ)
        — both summed over the radial axis.

    polar_full shape: (N, theta, r)
    Returns D shape (N, N), values in [0, 1] (cosine-style after norm).
    """
    N, T, R = polar_full.shape
    # Build per-pattern angular-FFT, summed over radial axis (one
    # FFT per (i, r) is too memory-heavy; sum r first to keep an
    # interpretable signal, then FFT once per pattern).
    sig = polar_full.sum(axis=2)                # (N, T)  — I(θ) integrated over r
    # Normalize per pattern so xcorr is cosine-like.
    nrm = np.linalg.norm(sig, axis=1, keepdims=True).clip(min=1e-9)
    sig_n = (sig / nrm).astype(np.float32)
    F = np.fft.rfft(sig_n, axis=1)              # (N, T//2+1) complex
    D = np.zeros((N, N), dtype=np.float32)
    for i in range(N):
        cross = F[i:i+1] * np.conj(F)            # broadcasting (N, T//2+1)
        corr = np.fft.irfft(cross, n=T, axis=1)  # (N, T)
        max_corr = np.maximum(corr.max(axis=1), 0.0)
        D[i] = 1.0 - max_corr                    # 0 = perfect match
    return D


def auto_K(W, range_=(2, 12), progress_cb=None) -> tuple[int, np.ndarray]:
    """Best K by silhouette score using K-means."""
    from sklearn.cluster import KMeans
    from sklearn.metrics import silhouette_score
    Ks, S = [], []
    for k in range(range_[0], range_[1]):
        km = KMeans(n_clusters=int(k), n_init=10, random_state=42)
        lbls = km.fit_predict(W)
        try:
            s = float(silhouette_score(W, lbls))
        except Exception:
            s = float("nan")
        Ks.append(k); S.append(s)
        if progress_cb is not None:
            progress_cb(k - range_[0] + 1, range_[1] - range_[0],
                          "scan K")
    Sa = np.asarray(S, float)
    Sa = np.where(np.isnan(Sa), -np.inf, Sa)
    return int(Ks[int(np.argmax(Sa))]), np.asarray(S, float)


# ---------------------------------------------------------------------------
class NMFPanel(ctk.CTkFrame):

    def __init__(self, master, app=None):
        super().__init__(master)
        self.app = app
        self.outdir = None
        self.sample = None
        self._scan_shape = None
        self._last = None        # dict(W, H, comp_shape, labels{...})
        self._compute_running = False
        self._compute_progress = ""
        self._stop_requested = False
        self._lock = threading.Lock()
        self._thread = None
        self._build()

    # ----- run linkage --------------------------------------------------
    def link_run(self, outdir, sample):
        self.outdir = outdir
        self.sample = sample
        try:
            self._vars["sample"].set(sample)
        except Exception:
            pass
        self._refresh_scan_shape()
        if self._vars["use_sample_vmax"].get():
            self._snap_vmax_to_sample()
        self._info_lbl.configure(
            text=f"sample: {sample}   (run dir = "
                  f"{os.path.basename(outdir) if outdir else '—'})  "
                  f"vmax = {self._vars['vmax'].get():g}")
        self._last = None
        self._render_idle()

    def on_runtime_sample_added(self, key):
        # A cube was loaded in the Data tab -> make it THIS panel's dataset.
        try:
            self._vars["sample"].set(key)
            self._on_sample_change()
        except Exception:
            pass

    def _sync_from_pre(self):
        """Adopt the cube currently loaded in the Data tab (app.pre)."""
        pre = getattr(self.app, "pre", None)
        try:
            k = pre.get_sample_key() if pre is not None else None
        except Exception:
            k = None
        if k and k != self.sample:
            self._vars["sample"].set(k)
            self._on_sample_change()
        return self.sample

    def _refresh_scan_shape(self):
        try:
            from data import SAMPLES
            cfg = SAMPLES.get(self.sample)
            self._scan_shape = (cfg.get("scan_shape")
                                  or cfg.get("scan_size")
                                  or None) if cfg else None
        except Exception:
            self._scan_shape = None

    # ----- variant citation --------------------------------------------
    def _on_variant_change(self):
        """Update the clickable citation label for the selected variant."""
        cfg = NMF_VARIANTS.get(self._vars["variant"].get(), {})
        self._cite_url = cfg.get("url", "")
        ref = cfg.get("ref", "")
        try:
            self._cite_lbl.configure(
                text=("↗ " + ref) if ref else "")
        except Exception:
            pass

    def _open_variant_ref(self):
        url = getattr(self, "_cite_url", "")
        if url:
            try:
                webbrowser.open(url)
            except Exception:
                pass

    # ----- UI ----------------------------------------------------------
    def _build(self):
        self._vars = {
            "sample":    ctk.StringVar(value=""),
            "vmax":      ctk.DoubleVar(value=2.0),    # DINO-style scaling
            "use_sample_vmax": ctk.BooleanVar(value=True),
            "variant":   ctk.StringVar(value=next(iter(NMF_VARIANTS))),
            "n_comp":    ctk.IntVar(value=8),
            "auto_n":    ctk.BooleanVar(value=True),
            "use_kmeans": ctk.BooleanVar(value=True),
            "use_aglo":   ctk.BooleanVar(value=False),
            "use_hdbscan": ctk.BooleanVar(value=False),
            "use_fcm":   ctk.BooleanVar(value=False),
            "fcm_m":     ctk.DoubleVar(value=2.0),
            "aglo_dist":  ctk.StringVar(value="euclidean"),
            "K":          ctk.IntVar(value=6),
            "auto_K":     ctk.BooleanVar(value=True),
            "min_cluster_size": ctk.IntVar(value=30),
            "view_mode":  ctk.StringVar(value="class average"),
            "frame_idx":  ctk.IntVar(value=0),
        }

        # ---- top bar ----
        top = ctk.CTkFrame(self)
        top.pack(side="top", fill="x", padx=6, pady=6)
        try:
            from data import SAMPLES
            sample_values = sorted(SAMPLES.keys())
        except Exception:
            sample_values = [""]
        # Dataset follows the Data tab (top of the app) — no per-panel sample
        # dropdown / loader (avoids defaulting to the wrong sample).
        ctk.CTkLabel(top, text="dataset:").pack(side="left", padx=(8, 4))
        self._ds_lbl = ctk.CTkLabel(top, text="(load a cube in the Data tab)",
                                     font=("Consolas", 10, "bold"))
        self._ds_lbl.pack(side="left", padx=4)
        self._info_lbl = ctk.CTkLabel(top, text="(no sample yet)",
                                        font=("Consolas", 10))
        self._info_lbl.pack(side="left", padx=8)
        self._run_btn = ctk.CTkButton(top, text="Run",
                                         width=100,
                                         fg_color=("#2D7A2D", "#1F7A1F"),
                                         command=self._kickoff_run)
        self._run_btn.pack(side="right", padx=4)
        self._stop_btn = ctk.CTkButton(top, text="Stop",
                                          width=80,
                                          fg_color=("#a23030", "#7a1f1f"),
                                          state="disabled",
                                          command=self._on_stop_clicked)
        self._stop_btn.pack(side="right", padx=4)
        self._report_btn = ctk.CTkButton(top,
            text="Report All  →  PPTX + PNGs",
            width=220,
            fg_color=("#2D7A2D", "#1F7A1F"),
            command=self._kickoff_report_all)
        self._report_btn.pack(side="right", padx=4)

        # ---- body ----
        body = ctk.CTkFrame(self)
        body.pack(side="top", fill="both", expand=True, padx=6, pady=4)
        sb = ctk.CTkScrollableFrame(body, width=320)
        sb.pack(side="left", fill="y")

        _section(sb, "Pre-process scaling (DINO-style)")
        v_row = ctk.CTkFrame(sb, fg_color="transparent")
        v_row.pack(fill="x", padx=4, pady=2)
        ctk.CTkLabel(v_row, text="vmax:", width=60).pack(side="left")
        ctk.CTkEntry(v_row, textvariable=self._vars["vmax"],
                       width=70).pack(side="left", padx=2)
        ctk.CTkCheckBox(v_row, text="use sample default",
                          variable=self._vars["use_sample_vmax"],
                          command=self._on_use_sample_vmax
                          ).pack(side="left", padx=8)

        _section(sb, "NMF variant")
        ctk.CTkOptionMenu(sb, variable=self._vars["variant"],
                            values=list(NMF_VARIANTS.keys()),
                            width=240,
                            command=lambda _v: self._on_variant_change()
                            ).pack(anchor="w", padx=8, pady=2)
        # clickable citation for the selected variant (opens DOI in browser)
        self._cite_lbl = ctk.CTkLabel(
            sb, text="", width=240, anchor="w", justify="left",
            wraplength=240, font=("Segoe UI", 10),
            text_color=("#1565C0", "#5BA4F0"), cursor="hand2")
        self._cite_lbl.pack(anchor="w", padx=8, pady=(0, 2))
        self._cite_lbl.bind(
            "<Button-1>",
            lambda _e: self._open_variant_ref())
        self._on_variant_change()      # initialise label for default variant

        _section(sb, "n_comp")
        n_row = ctk.CTkFrame(sb, fg_color="transparent")
        n_row.pack(fill="x", padx=4, pady=2)
        ctk.CTkLabel(n_row, text="n_comp:", width=80).pack(side="left")
        ctk.CTkEntry(n_row, textvariable=self._vars["n_comp"],
                       width=60).pack(side="left", padx=4)
        ctk.CTkCheckBox(n_row, text="auto (knee)",
                          variable=self._vars["auto_n"]
                          ).pack(side="left", padx=8)

        _section(sb, "Clustering methods")
        ctk.CTkCheckBox(sb, text="K-means  (default)",
                          variable=self._vars["use_kmeans"]
                          ).pack(anchor="w", padx=8, pady=1)
        a_row = ctk.CTkFrame(sb, fg_color="transparent")
        a_row.pack(fill="x", padx=4, pady=1)
        ctk.CTkCheckBox(a_row, text="Aglo",
                          variable=self._vars["use_aglo"]
                          ).pack(side="left", padx=8)
        ctk.CTkLabel(a_row, text="distance:").pack(side="left",
                                                       padx=(4, 2))
        ctk.CTkOptionMenu(a_row, variable=self._vars["aglo_dist"],
                            values=AGLO_DISTANCES, width=110
                            ).pack(side="left", padx=2)
        h_row = ctk.CTkFrame(sb, fg_color="transparent")
        h_row.pack(fill="x", padx=4, pady=1)
        ctk.CTkCheckBox(h_row, text="HDBSCAN",
                          variable=self._vars["use_hdbscan"]
                          ).pack(side="left", padx=8)
        ctk.CTkLabel(h_row, text="min size:").pack(side="left",
                                                       padx=(4, 2))
        ctk.CTkEntry(h_row, textvariable=self._vars["min_cluster_size"],
                       width=60).pack(side="left", padx=2)
        f_row = ctk.CTkFrame(sb, fg_color="transparent")
        f_row.pack(fill="x", padx=4, pady=1)
        ctk.CTkCheckBox(f_row, text="FCM  (Bezdek 1981)",
                          variable=self._vars["use_fcm"]
                          ).pack(side="left", padx=8)
        ctk.CTkLabel(f_row, text="fuzziness m:").pack(side="left",
                                                          padx=(4, 2))
        ctk.CTkEntry(f_row, textvariable=self._vars["fcm_m"],
                       width=60).pack(side="left", padx=2)

        _section(sb, "K  (ignored by HDBSCAN)")
        k_row = ctk.CTkFrame(sb, fg_color="transparent")
        k_row.pack(fill="x", padx=4, pady=2)
        ctk.CTkLabel(k_row, text="K:", width=80).pack(side="left")
        ctk.CTkEntry(k_row, textvariable=self._vars["K"], width=60
                       ).pack(side="left", padx=4)
        ctk.CTkCheckBox(k_row, text="auto (silhouette)",
                          variable=self._vars["auto_K"]
                          ).pack(side="left", padx=8)
        # Re-cluster the EXISTING NMF (no re-fit) — fast way to change K
        # or add clustering methods without recomputing the decomposition.
        self._cluster_btn = ctk.CTkButton(
            sb, text="Cluster  (re-cluster current NMF)",
            fg_color=("#4D6FB0", "#3A5380"),
            command=self._kickoff_recluster)
        self._cluster_btn.pack(fill="x", padx=8, pady=(2, 6))

        _section(sb, "View / inspect")
        ctk.CTkOptionMenu(sb, variable=self._vars["view_mode"],
            values=["class average", "single frame", "domain + average"],
            width=240,
            command=lambda _v: self._render_last()
            ).pack(anchor="w", padx=8, pady=2)
        f_row = ctk.CTkFrame(sb, fg_color="transparent")
        f_row.pack(fill="x", padx=4, pady=2)
        ctk.CTkLabel(f_row, text="frame idx:", width=80).pack(side="left")
        ctk.CTkEntry(f_row, textvariable=self._vars["frame_idx"],
                       width=60).pack(side="left", padx=4)

        ctk.CTkButton(sb, text="Open interactive map…",
                       fg_color=("#4D6FB0", "#3A5380"),
                       command=self._open_interactive_map
                       ).pack(fill="x", padx=8, pady=(8, 2))
        ctk.CTkButton(sb, text="Show class averages…",
                       command=self._open_class_averages_popup
                       ).pack(fill="x", padx=8, pady=(2, 2))
        ctk.CTkButton(sb, text="Save snapshot",
                       command=self._save_snapshot
                       ).pack(fill="x", padx=8, pady=(2, 4))

        self._status_lbl = ctk.CTkLabel(sb,
            text="", font=("Consolas", 9), justify="left",
            text_color=("#444", "#aaa"), wraplength=300)
        self._status_lbl.pack(anchor="w", padx=8, pady=(8, 4))

        # canvas
        canv = ctk.CTkFrame(body)
        canv.pack(side="left", fill="both", expand=True, padx=(6, 0))
        self._fig = Figure(figsize=(13, 8))
        self._canvas = FigureCanvasTkAgg(self._fig, master=canv)
        self._canvas.get_tk_widget().pack(fill="both", expand=True)
        NavigationToolbar2Tk(self._canvas, canv)
        # Double-click any inline class map -> open the big interactive
        # viewer at that method (left/right/shift+right click inside).
        self._map_axes = {}     # ax -> method name
        self._canvas.mpl_connect("button_press_event",
                                  self._on_inline_map_click)
        self._render_idle()

        # Populate sample default once.
        # Bind to whatever is loaded in the Data tab (kept in sync via
        # on_runtime_sample_added); no default sample selection.
        self.after(200, self._sync_from_pre)
        # If post-hoc has a run linked, adopt its sample automatically.
        try:
            self.after(150, self._try_auto_link_from_loaded_run)
        except Exception:
            pass

    def _on_sample_change(self):
        s = self._vars["sample"].get()
        self.sample = s if s else None
        try:
            self._ds_lbl.configure(
                text=self.sample or "(load a cube in the Data tab)")
        except Exception:
            pass
        self._refresh_scan_shape()
        # When the user changes sample and 'use sample default' is on,
        # snap vmax to the sample's training vmax.
        if self._vars["use_sample_vmax"].get():
            self._snap_vmax_to_sample()
        self._info_lbl.configure(
            text=f"sample: {self.sample}   "
                  f"scan = {self._scan_shape}   "
                  f"vmax = {self._vars['vmax'].get():g}")

    def _snap_vmax_to_sample(self):
        try:
            from data import SAMPLES
            cfg = SAMPLES.get(self.sample) or {}
            v = float(cfg.get("vmax", 2.0))
            self._vars["vmax"].set(v)
        except Exception:
            pass

    def _on_use_sample_vmax(self):
        if self._vars["use_sample_vmax"].get():
            self._snap_vmax_to_sample()
        self._info_lbl.configure(
            text=f"sample: {self.sample}   "
                  f"scan = {self._scan_shape}   "
                  f"vmax = {self._vars['vmax'].get():g}")

    def _load_cube_from_disk(self):
        """Pick a .prz / .npy / .h5 cube from disk and register it as a
        runtime sample. For 3D HDF5 masters (Eiger / Dectris layout), a
        small dialog asks for the scan_shape (Ny, Nx) since master files
        store data as (N_frames, H, W) without scan-grid info."""
        p = filedialog.askopenfilename(
            title="Pick a cube  (.prz / .npz / .npy / .h5)",
            filetypes=[("Cube files",
                          "*.prz *.npz *.npy *.h5 *.hdf5"),
                          ("PRZ / NPZ", "*.prz *.npz"),
                          ("NPY", "*.npy"),
                          ("HDF5", "*.h5 *.hdf5"),
                          ("All files", "*.*")])
        if not p:
            return
        # Friendly nudge: if the user picked a Dectris _data_NNNNNN.h5
        # fragment, point them at the matching master and offer to swap.
        bn = os.path.basename(p)
        import re as _re
        m = _re.match(r"^(.*)_data_(\d+)\.h5$", bn)
        if m:
            master = os.path.join(os.path.dirname(p),
                                     m.group(1) + "_master.h5")
            if os.path.exists(master):
                if messagebox.askyesno("Dectris fragment",
                    f"You picked a Dectris data fragment:\n  {bn}\n\n"
                    f"The matching master file exists:\n  "
                    f"{os.path.basename(master)}\n\n"
                    f"Load the master instead? (recommended — it stitches "
                    f"all data files and carries scan_shape metadata.)"):
                    p = master
        try:
            v = float(self._vars["vmax"].get())
        except Exception:
            v = 2.0
        scan_shape = None
        # For 3D HDF5 masters we need scan_shape.  Peek the file shape.
        if p.lower().endswith((".h5", ".hdf5")):
            try:
                import h5py
                from data import (_h5_find_data_path,
                                      _h5_infer_scan_shape,
                                      _h5_dectris_external_data)
                with h5py.File(p, "r") as fh:
                    try:
                        dpath, ndim = _h5_find_data_path(fh)
                        s = tuple(fh[dpath].shape)
                    except ValueError:
                        # Dectris master with external links.
                        pairs = _h5_dectris_external_data(fh, p)
                        if not pairs:
                            raise
                        total = 0; H_ = W_ = None
                        for fp, dp in pairs:
                            with h5py.File(fp, "r") as gf:
                                sh = tuple(gf[dp].shape)
                                if len(sh) != 3: continue
                                total += sh[0]
                                H_, W_ = sh[1], sh[2]
                        s = (total, H_, W_); ndim = 3
                    if ndim == 3:
                        scan_shape = _h5_infer_scan_shape(fh, s[0])
            except Exception as e:
                messagebox.showerror("h5 peek failed", repr(e)); return
            if ndim == 3 and scan_shape is None:
                N, H, W = s
                from gui_app._dialogs import ask_scan_shape
                scan_shape = ask_scan_shape(self, N, H, W)
                if scan_shape is None:
                    return
        try:
            from data import register_runtime_sample, SAMPLES
            kwargs = dict(vmax=v)
            if scan_shape is not None:
                kwargs["scan_shape"] = scan_shape
            key = register_runtime_sample(p, **kwargs)
        except Exception as e:
            messagebox.showerror("Load cube", f"failed:\n{e!r}"); return
        try:
            self._sample_menu.configure(values=sorted(SAMPLES.keys()))
        except Exception:
            pass
        self._vars["sample"].set(key)
        self._on_sample_change()


    def _try_auto_link_from_loaded_run(self):
        """If the Post-hoc panel already has a linked run, adopt its
        sample so the user doesn't have to pick again."""
        ph = getattr(self.app, "posthoc", None) if self.app else None
        if ph is None:
            return
        s = getattr(ph, "sample", None)
        out = getattr(ph, "outdir", None)
        if not s:
            return
        try:
            from data import SAMPLES
            if s not in SAMPLES:
                return
        except Exception:
            return
        self.outdir = out
        self.sample = s
        self._vars["sample"].set(s)
        self._refresh_scan_shape()
        if self._vars["use_sample_vmax"].get():
            self._snap_vmax_to_sample()
        self._info_lbl.configure(
            text=f"auto-linked: {s}  "
                  f"({os.path.basename(out) if out else 'no run'})  "
                  f"scan = {self._scan_shape}")

    # ----- compute -----------------------------------------------------
    class _Aborted(Exception):
        """Raised by workers when self._stop_requested becomes True."""
        pass

    def _check_stop(self):
        """Called by workers at safe points to exit cleanly on Stop."""
        if self._stop_requested:
            raise NMFPanel._Aborted("user requested stop")

    def _on_stop_clicked(self):
        if not self._compute_running:
            return
        self._stop_requested = True
        with self._lock:
            self._compute_progress = "stopping (waiting for current step to finish)…"
        try: self._stop_btn.configure(state="disabled")
        except Exception: pass

    def _kickoff_run(self):
        if self._compute_running:
            messagebox.showinfo("NMF", "compute already running"); return
        self._sync_from_pre()
        if not self.sample:
            messagebox.showinfo("NMF",
                "Load a cube in the Data tab (top) first."); return
        self._compute_running = True
        self._stop_requested = False
        self._run_btn.configure(state="disabled")
        self._report_btn.configure(state="disabled")
        try:
            self._cluster_btn.configure(state="disabled")
        except Exception:
            pass
        self._stop_btn.configure(state="normal")
        self._compute_progress = "starting…"
        self._thread = threading.Thread(
            target=self._compute_worker, daemon=True)
        self._thread.start()
        self._poll()

    # ---- re-cluster the existing NMF (no re-fit) ----------------------
    def _do_clustering(self, W, polar_full=None):
        """Run the selected clustering method(s) on NMF loadings W.
        Resolves K (auto-silhouette or fixed).  Returns (labels, K, sil)."""
        labels = {}
        min_cs = int(self._vars["min_cluster_size"].get())
        sil = None
        if self._vars["auto_K"].get():
            with self._lock:
                self._compute_progress = "auto K (silhouette)…"
            K, sil = auto_K(W)
            self._vars["K"].set(K)
        else:
            K = int(self._vars["K"].get())
        self._check_stop()
        if self._vars["use_kmeans"].get():
            with self._lock:
                self._compute_progress = "K-means…"
            labels["K-means"] = cluster_W(W, "K-means", k=K)
            self._check_stop()
        if self._vars["use_aglo"].get():
            d = self._vars["aglo_dist"].get()
            with self._lock:
                self._compute_progress = f"Aglo ({d})…"
            if d == "polar-xcorr" and polar_full is None:
                raise RuntimeError("polar-xcorr needs the polar input "
                                   "(re-run NMF with the polar variant)")
            labels["Aglo"] = cluster_W(W, "Aglo", k=K, distance=d,
                                       polar_full=polar_full)
            self._check_stop()
        if self._vars["use_hdbscan"].get():
            with self._lock:
                self._compute_progress = "HDBSCAN…"
            labels["HDBSCAN"] = cluster_W(W, "HDBSCAN",
                                          min_cluster_size=min_cs)
            self._check_stop()
        if self._vars["use_fcm"].get():
            with self._lock:
                self._compute_progress = "FCM…"
            labels["FCM"] = cluster_W(W, "FCM", k=K,
                                      fcm_m=float(self._vars["fcm_m"].get()))
            self._check_stop()
        return labels, K, sil

    def _kickoff_recluster(self):
        if self._compute_running:
            messagebox.showinfo("NMF", "compute already running"); return
        if self._last is None or self._last.get("W") is None:
            messagebox.showinfo(
                "Cluster", "No NMF to re-cluster yet — click 'Run' first. "
                "Then 'Cluster' re-clusters without re-fitting.")
            return
        self._compute_running = True
        self._stop_requested = False
        self._run_btn.configure(state="disabled")
        self._report_btn.configure(state="disabled")
        try:
            self._cluster_btn.configure(state="disabled")
        except Exception:
            pass
        self._stop_btn.configure(state="normal")
        self._compute_progress = "re-clustering…"
        self._thread = threading.Thread(
            target=self._recluster_worker, daemon=True)
        self._thread.start()
        self._poll()

    def _recluster_worker(self):
        try:
            d = self._last
            labels, K, sil = self._do_clustering(d["W"],
                                                 d.get("polar_full"))
            d["labels"] = labels
            d["K"] = K
            d["sil"] = sil
            self.last_cluster_labels = dict(labels)
            with self._lock:
                self._compute_progress = (
                    f"re-clustered.  K={K}  "
                    f"methods={', '.join(labels) or '(none)'}")
        except NMFPanel._Aborted:
            with self._lock:
                self._compute_progress = "stopped by user."
        except Exception as e:
            with self._lock:
                self._compute_progress = f"re-cluster failed: {e!r}"
            print(f"[nmf] recluster failed: {e!r}", flush=True)
        finally:
            self._compute_running = False

    def _compute_worker(self):
        try:
            t0 = time.perf_counter()
            variant_name = self._vars["variant"].get()
            cfg = dict(NMF_VARIANTS[variant_name])
            def cb(done, total, label):
                with self._lock:
                    self._compute_progress = (
                        f"{label}: {done}/{total} "
                        f"({100*done/max(total,1):.0f}%)")
            vmax = float(self._vars["vmax"].get())
            X, X_aug, comp_shape, info = build_nmf_input(
                self.sample, cfg, progress_cb=cb,
                vmax_override=vmax)
            self._check_stop()
            with self._lock:
                self._compute_progress = (
                    f"matrix shape: {X.shape}  "
                    f"(aug: {None if X_aug is None else X_aug.shape})")
            # n_comp
            auto_n = self._vars["auto_n"].get()
            if auto_n:
                with self._lock:
                    self._compute_progress = "auto n_comp scan…"
                n_comp, errs = auto_n_comp(X, X_aug,
                                              sparse=cfg["sparse"],
                                              progress_cb=cb)
                self._vars["n_comp"].set(n_comp)
            else:
                n_comp = int(self._vars["n_comp"].get())
                errs = None
            self._check_stop()
            with self._lock:
                self._compute_progress = (
                    f"fit NMF (n_comp={n_comp}, sparse="
                    f"{cfg['sparse']}, aug={X_aug is not None})…")
            W, H, err_final = fit_nmf(X, X_aug, n_comp,
                                          sparse=cfg["sparse"],
                                          max_iter=300)
            self._check_stop()
            # Keep polar_full (un-flat) for polar-xcorr + later re-cluster.
            polar_full = None
            if cfg["input"] == "polar":
                polar_full = X.reshape((X.shape[0],) + comp_shape)
            self._check_stop()
            # Clustering (shared with the 'Cluster' re-cluster path; also
            # resolves K / auto-K).
            labels, K, sil = self._do_clustering(W, polar_full)
            self._last = dict(
                variant=variant_name, cfg=cfg,
                X_shape=X.shape, comp_shape=comp_shape,
                W=W, H=H, n_comp=n_comp, K=K,
                labels=labels, errs=errs, sil=sil,
                err_final=err_final, polar_full=polar_full,
            )
            # Expose cluster labels for cross-tab overlays (ACOM).
            self.last_cluster_labels = dict(labels)
            with self._lock:
                self._compute_progress = (
                    f"done. ({time.perf_counter() - t0:.1f}s)  "
                    f"n_comp={n_comp}  K={K}  "
                    f"err={err_final:.3f}")
        except NMFPanel._Aborted:
            with self._lock:
                self._compute_progress = "stopped by user."
        except Exception as e:
            err = repr(e)
            print(f"[nmf] worker failed: {err}", flush=True)
            with self._lock:
                self._compute_progress = f"failed: {err}"
        finally:
            self._compute_running = False

    def _poll(self):
        with self._lock:
            prog = self._compute_progress
        self._status_lbl.configure(text=prog or "(running…)")
        if self._compute_running:
            self.after(500, self._poll)
        else:
            self._run_btn.configure(state="normal")
            self._report_btn.configure(state="normal")
            try:
                self._cluster_btn.configure(state="normal")
            except Exception:
                pass
            try: self._stop_btn.configure(state="disabled")
            except Exception: pass
            self._stop_requested = False
            if self._last is not None:
                self._render_last()

    # ----- rendering ---------------------------------------------------
    def _render_idle(self):
        self._fig.clear()
        ax = self._fig.add_subplot(111)
        ax.text(0.5, 0.5,
                "1. Pick a sample.\n"
                "2. Pick an NMF variant + clustering methods.\n"
                "3. Click Run  (Auto modes: knee for n_comp, "
                "silhouette for K).",
                ha="center", va="center", fontsize=11)
        ax.set_axis_off()
        self._canvas.draw_idle()

    def _render_last(self):
        if self._last is None:
            self._render_idle(); return
        d = self._last
        comp_shape = d["comp_shape"]
        H = d["H"]            # (n_comp, D)
        labels = d["labels"]
        K = d["K"]
        n_comp = d["n_comp"]
        Ny, Nx = (self._scan_shape if self._scan_shape
                    else (1, d["X_shape"][0]))
        errs = d.get("errs")
        sil  = d.get("sil")
        diag = (errs is not None) or (sil is not None)

        self._fig.clear()
        n_methods = max(1, len(labels))
        # Layout: row 0 = NMF components, row 1 (optional) = diagnostics
        # (recon-err vs n_comp + silhouette vs K), last row = class maps.
        if diag:
            gs = self._fig.add_gridspec(
                3, max(n_comp, n_methods, 2),
                height_ratios=[1.0, 0.55, 1.0])
            map_row = 2
        else:
            gs = self._fig.add_gridspec(
                2, max(n_comp, n_methods),
                height_ratios=[1.0, 1.0])
            map_row = 1

        # ---- NMF components ----
        for k in range(n_comp):
            ax = self._fig.add_subplot(gs[0, k])
            comp = H[k].reshape(comp_shape)
            if comp.shape[0] == 1:                   # 1D radial
                ax.plot(comp[0], color="black", lw=1.0)
                ax.set_xticks([]); ax.set_yticks([])
            else:
                ax.imshow(comp, cmap="inferno", aspect="auto")
                ax.set_xticks([]); ax.set_yticks([])
            ax.set_title(f"H[{k}]", fontsize=9)

        # ---- diagnostic curves (only when auto modes were used) ----
        if diag:
            n_cols = max(n_comp, n_methods, 2)
            # Use first half for recon-err, second half for silhouette.
            half = max(1, n_cols // 2)
            if errs is not None:
                ax_e = self._fig.add_subplot(gs[1, 0:half])
                xs = np.arange(2, 2 + len(errs))
                ax_e.plot(xs, errs, marker="o", lw=1.4, color="#1f77b4")
                ax_e.axvline(n_comp, color="orange", ls="--", lw=1.2,
                              label=f"chosen n_comp = {n_comp} (knee)")
                # Annotate the knee
                ie = int(np.clip(n_comp - 2, 0, len(errs) - 1))
                ax_e.scatter([n_comp], [errs[ie]], s=60,
                                edgecolors="orange",
                                facecolors="none", lw=2,
                                zorder=5)
                ax_e.set_xlabel("n_comp")
                ax_e.set_ylabel("reconstruction error")
                ax_e.set_title("auto n_comp: knee scan", fontsize=9)
                ax_e.legend(fontsize=7, loc="upper right")
                ax_e.grid(alpha=0.3)
            if sil is not None:
                ax_s = self._fig.add_subplot(gs[1, half:n_cols])
                xs = np.arange(2, 2 + len(sil))
                ax_s.plot(xs, sil, marker="o", lw=1.4, color="#d62728")
                ax_s.axvline(K, color="orange", ls="--", lw=1.2,
                              label=f"chosen K = {K} (max silhouette)")
                ik = int(np.clip(K - 2, 0, len(sil) - 1))
                ax_s.scatter([K], [sil[ik]], s=60,
                                edgecolors="orange",
                                facecolors="none", lw=2,
                                zorder=5)
                ax_s.set_xlabel("K")
                ax_s.set_ylabel("silhouette score")
                ax_s.set_title("auto K: silhouette scan", fontsize=9)
                ax_s.legend(fontsize=7, loc="upper right")
                ax_s.grid(alpha=0.3)

        # ---- Class maps per clustering method ----
        self._map_axes = {}
        for col, (method, lbl) in enumerate(labels.items()):
            ax = self._fig.add_subplot(gs[map_row, col])
            self._map_axes[ax] = method
            n_classes = int(lbl.max()) + 1
            # Match DINO class-map style: tab10/tab20 ListedColormap.
            base = "tab20" if n_classes > 10 else "tab10"
            cmap = plt.get_cmap(base)
            palette = ListedColormap(
                [cmap(i % cmap.N) for i in range(n_classes)])
            grid = lbl.reshape(Ny, Nx) if lbl.size == Ny * Nx \
                else lbl.reshape(1, -1)
            im = ax.imshow(grid, cmap=palette,
                             vmin=-0.5, vmax=n_classes - 0.5,
                             interpolation="nearest", aspect="equal")
            ax.set_title(f"{method}   K={n_classes}", fontsize=10)
            ax.set_axis_off()
            # Integer colorbar like the DINO panel.
            cb = self._fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04,
                                       ticks=list(range(n_classes)))
            cb.set_label("class id", fontsize=8)
            # Real-space 100 nm scale bar.
            try:
                nm_per_px = float(self.app.real_res.get()) if self.app \
                              else 0.0
            except Exception:
                nm_per_px = 0.0
            if nm_per_px > 0:
                from gui_app._calib_utils import add_real_scalebar
                add_real_scalebar(ax, nm_per_px, length_nm=100,
                                    color="white")

        vmax = self._vars["vmax"].get()
        title = (f"{self.sample}   variant: {d['variant']}   "
                 f"vmax={vmax:g}   n_comp={n_comp}   K={K}   "
                 f"recon-err={d['err_final']:.3f}")
        self._fig.suptitle(title, fontsize=11)
        self._fig.tight_layout()
        self._canvas.draw_idle()
        # Auto-save a copy of the result figure next to the loaded data.
        try:
            import assistant_io
            methods = "+".join(_safe_name(m) for m in d.get("labels", {}))
            assistant_io.gui_autosave(
                self.app, "nmf", self._fig,
                name=f"nmf_K{K}_{methods}" if methods else f"nmf_K{K}",
                summary=title)
        except Exception:
            pass

    # ----- interactive map --------------------------------------------
    def _recip_per_px(self):
        try:
            return float(self.app.recip_res.get()) if self.app else 0.0
        except Exception:
            return 0.0

    def _open_interactive_map(self, method=None):
        from data import SAMPLES
        if self._last is None or not self._last.get("labels"):
            messagebox.showinfo("interactive map",
                "Run NMF + clustering first."); return
        if not self.sample or self.sample not in SAMPLES:
            messagebox.showinfo("interactive map",
                "No dataset loaded for this sample."); return
        if not self._scan_shape:
            messagebox.showinfo("interactive map",
                "Scan shape unknown — cannot map clusters to positions.")
            return
        labels = self._last["labels"]
        if method in labels:        # open with the clicked method first
            labels = {method: labels[method],
                       **{m: l for m, l in labels.items() if m != method}}
        try:
            from gui_app.cluster_interactive import (
                open_interactive_clustermap)
            open_interactive_clustermap(
                self, sample=self.sample, scan_shape=self._scan_shape,
                labels=labels, recip_per_px=self._recip_per_px(),
                title=f"NMF + cluster — {self.sample}")
        except Exception as e:
            messagebox.showerror("interactive map", repr(e))

    def _on_inline_map_click(self, event):
        # Double-click an inline class map -> big interactive viewer.
        if not getattr(event, "dblclick", False):
            return
        method = self._map_axes.get(event.inaxes)
        if method is not None:
            self._open_interactive_map(method)

    # ----- save / export ---------------------------------------------
    def _save_snapshot(self):
        if self._last is None or not self.sample:
            messagebox.showinfo("save", "nothing to save"); return
        out_dir = os.path.join(
            self.outdir or ".", "nmf",
            f"{self.sample}_{_safe_name(self._last['variant'])}")
        os.makedirs(out_dir, exist_ok=True)
        try:
            self._fig.savefig(os.path.join(out_dir, "summary.png"),
                                dpi=140, bbox_inches="tight")
            np.save(os.path.join(out_dir, "W.npy"), self._last["W"])
            np.save(os.path.join(out_dir, "H.npy"), self._last["H"])
            for m, lbl in self._last["labels"].items():
                np.save(os.path.join(out_dir, f"labels_{m}.npy"), lbl)
            with open(os.path.join(out_dir, "summary.json"), "w") as fh:
                json.dump({
                    "sample": self.sample,
                    "variant": self._last["variant"],
                    "n_comp": int(self._last["n_comp"]),
                    "K": int(self._last["K"]),
                    "err_final": float(self._last["err_final"]),
                    "methods": list(self._last["labels"].keys()),
                }, fh, indent=2)
        except Exception as e:
            messagebox.showerror("save", repr(e)); return
        self._status_lbl.configure(text=f"saved → {out_dir}")

    # ----- class-averages popup ----------------------------------------
    def _open_class_averages_popup(self):
        """Per-cluster average diffraction pattern for each clustering
        method in the last run. Up to top-N (=300) raw patterns per
        cluster are averaged; large clusters get a random subsample."""
        if self._last is None or not self.sample:
            messagebox.showinfo(
                "averages", "Run NMF + clustering first."); return
        try:
            from data import SAMPLES, LoadPRZ
        except Exception as e:
            messagebox.showerror("data", repr(e)); return
        cfg = SAMPLES.get(self.sample) or {}
        try:
            v = float(self._vars["vmax"].get())
        except Exception:
            v = float(cfg.get("vmax", 2.0))
        labels_dict = self._last["labels"]
        if not labels_dict:
            messagebox.showinfo("averages",
                "No clustering results to average."); return

        win = tk.Toplevel(self)
        win.title(f"{self.sample} — class averages "
                    f"(vmax={v:g})")
        win.geometry("1200x780")
        ctrl = ctk.CTkFrame(win, fg_color="transparent")
        ctrl.pack(side="top", fill="x", padx=6, pady=4)
        method_var = ctk.StringVar(value=next(iter(labels_dict)))
        ctk.CTkLabel(ctrl, text="method:").pack(side="left", padx=(8, 4))
        ctk.CTkOptionMenu(ctrl, variable=method_var,
                            values=list(labels_dict.keys()),
                            width=180,
                            command=lambda _v: _redraw()
                            ).pack(side="left", padx=4)
        log_var = ctk.BooleanVar(value=True)
        ctk.CTkCheckBox(ctrl, text="log stretch",
                          variable=log_var,
                          command=lambda: _redraw()
                          ).pack(side="left", padx=12)
        ctk.CTkButton(ctrl, text="Save PNG",
                       width=100,
                       command=lambda: _save()
                       ).pack(side="right", padx=4)
        status = ctk.CTkLabel(win, text="",
                                font=("Consolas", 9), anchor="w")
        status.pack(side="top", fill="x", padx=8)

        fig = Figure(figsize=(11, 6.5), dpi=110, facecolor="white")
        canv = FigureCanvasTkAgg(fig, master=win)
        canv.get_tk_widget().pack(side="top", fill="both", expand=True)
        NavigationToolbar2Tk(canv, win)

        # Compute averages on demand and cache per method.
        avg_cache: dict[str, list[np.ndarray]] = {}

        def _compute(method):
            if method in avg_cache:
                return avg_cache[method]
            status.configure(text=f"computing class averages for "
                                       f"{method} …")
            win.update_idletasks()
            ds = LoadPRZ(cfg["path"], resize=POLAR_SIZE, vmax=v)
            lbl = labels_dict[method]
            n_classes = int(lbl.max()) + 1
            rng = np.random.default_rng(42)
            top_n = 300
            avgs = []
            for c in range(n_classes):
                idx = np.where(lbl == c)[0]
                if idx.size == 0:
                    avgs.append(np.zeros((POLAR_SIZE, POLAR_SIZE),
                                            np.float32))
                    continue
                pick = (idx if idx.size <= top_n
                          else rng.choice(idx, top_n, replace=False))
                pats = []
                for i in pick:
                    pats.append(ds[int(i)].squeeze().numpy())
                pats = np.stack(pats, axis=0).astype(np.float32)
                avgs.append(pats.mean(axis=0))
            avg_cache[method] = avgs
            return avgs

        def _redraw():
            method = method_var.get()
            avgs = _compute(method)
            n_classes = len(avgs)
            cols = min(n_classes, 4)
            rows = (n_classes + cols - 1) // cols
            fig.clear()
            for c in range(n_classes):
                ax = fig.add_subplot(rows, cols, c + 1)
                a = avgs[c]
                if log_var.get():
                    a_disp = np.log1p(np.clip(a * 50, 0, None))
                else:
                    a_disp = a
                ax.imshow(a_disp, cmap="inferno",
                           aspect="equal", interpolation="nearest")
                n_pts = int((labels_dict[method] == c).sum())
                ax.set_title(f"p{c}   N={n_pts}", fontsize=10)
                ax.set_xticks([]); ax.set_yticks([])
                # Reciprocal scale bar (raw resolution since ds[i] is the
                # vmax-scaled raw pattern resized to 192).
                try:
                    rp = float(self.app.recip_res.get()) \
                        if self.app else 0.0
                except Exception:
                    rp = 0.0
                if rp > 0:
                    from gui_app._calib_utils import (
                        q_per_polar_bin, get_raw_detector_size,
                        add_recip_scalebar)
                    qpx = q_per_polar_bin(
                        rp, get_raw_detector_size(self.sample))
                    add_recip_scalebar(ax, q_per_disp_px=qpx,
                                          length_q=0.2)
            fig.suptitle(
                f"{self.sample}  —  per-cluster averages "
                f"({method}, vmax={v:g})", fontsize=11)
            fig.tight_layout()
            canv.draw_idle()
            status.configure(
                text=f"{method}: {n_classes} classes (top-{300} per "
                      f"class, vmax={v:g}, log-stretch="
                      f"{log_var.get()})")

        def _save():
            if self.outdir is None:
                d = os.path.join(os.getcwd(), "nmf_class_averages")
            else:
                d = os.path.join(self.outdir, "nmf_class_averages")
            os.makedirs(d, exist_ok=True)
            v_safe = _safe_name(self._last["variant"])
            m_safe = _safe_name(method_var.get())
            p = os.path.join(d, f"{self.sample}_{v_safe}_{m_safe}.png")
            try:
                fig.savefig(p, dpi=140, bbox_inches="tight")
                status.configure(text=f"saved → {p}")
            except Exception as e:
                messagebox.showerror("save", repr(e))

        _redraw()

    # ===================================================================
    # Report All — sweep every variant × every clustering method
    # ===================================================================
    def _kickoff_report_all(self):
        if self._compute_running:
            messagebox.showinfo("Report", "compute already running")
            return
        self._sync_from_pre()
        if not self.sample:
            messagebox.showinfo("Report",
                "Load a cube in the Data tab (top) first."); return
        out_root = filedialog.askdirectory(
            initialdir=self.outdir or os.getcwd(),
            title="Pick output folder for NMF report  "
                   "(a 'nmf_report/' sub-folder will be created)",
            mustexist=False)
        if not out_root:
            return
        out_dir = os.path.join(
            out_root, "nmf_report",
            f"{self.sample}_{time.strftime('%Y%m%d-%H%M%S')}")
        png_dir = os.path.join(out_dir, "png")
        os.makedirs(png_dir, exist_ok=True)
        self._compute_running = True
        self._stop_requested = False
        self._run_btn.configure(state="disabled")
        self._report_btn.configure(state="disabled")
        self._stop_btn.configure(state="normal")
        self._compute_progress = "report all: starting…"
        self._thread = threading.Thread(
            target=self._report_worker,
            args=(out_dir, png_dir), daemon=True)
        self._thread.start()
        self._poll()

    def _report_worker(self, out_dir, png_dir):
        try:
            t0 = time.perf_counter()
            png_paths: list[tuple[str, str]] = []
            # Save a parameters image first.
            ppath = self._render_params_image(png_dir)
            if ppath:
                png_paths.append(("00_parameters", ppath))

            # Sweep all variants × clustering methods.
            for v_idx, variant_name in enumerate(NMF_VARIANTS.keys()):
                self._check_stop()
                with self._lock:
                    self._compute_progress = (
                        f"[{v_idx + 1}/{len(NMF_VARIANTS)}] "
                        f"{variant_name}: build + fit…")
                cfg = dict(NMF_VARIANTS[variant_name])
                X, X_aug, comp_shape, _info = build_nmf_input(
                    self.sample, cfg)
                self._check_stop()
                # Use auto n_comp + auto K for a quick survey.
                n_comp, _e = auto_n_comp(X, X_aug,
                                              sparse=cfg["sparse"])
                self._check_stop()
                W, H, err = fit_nmf(X, X_aug, n_comp,
                                          sparse=cfg["sparse"],
                                          max_iter=200)
                self._check_stop()
                K, _s = auto_K(W)
                self._check_stop()
                polar_full = None
                if cfg["input"] == "polar":
                    polar_full = X.reshape((X.shape[0],) + comp_shape)

                method_labels = {}
                method_labels["K-means"] = cluster_W(W, "K-means", k=K)
                method_labels["Aglo-eu"] = cluster_W(
                    W, "Aglo", k=K, distance="euclidean")
                method_labels["Aglo-cos"] = cluster_W(
                    W, "Aglo", k=K, distance="cosine")
                if cfg["input"] == "polar":
                    try:
                        method_labels["Aglo-xcorr"] = cluster_W(
                            W, "Aglo", k=K, distance="polar-xcorr",
                            polar_full=polar_full)
                    except Exception as e:
                        print(f"[report] aglo-xcorr failed for "
                              f"{variant_name}: {e!r}", flush=True)
                try:
                    method_labels["HDBSCAN"] = cluster_W(
                        W, "HDBSCAN", min_cluster_size=30)
                except Exception as e:
                    print(f"[report] HDBSCAN failed for "
                          f"{variant_name}: {e!r}", flush=True)
                try:
                    method_labels["FCM"] = cluster_W(
                        W, "FCM", k=K, fcm_m=2.0)
                except Exception as e:
                    print(f"[report] FCM failed for "
                          f"{variant_name}: {e!r}", flush=True)

                # Save figure.
                png_name = f"{v_idx + 1:02d}_{_safe_name(variant_name)}"
                p = os.path.join(png_dir, f"{png_name}.png")
                self._render_report_figure(
                    p, variant_name, cfg, comp_shape,
                    H=H, W=W, n_comp=n_comp, K=K,
                    err=err, labels=method_labels)
                png_paths.append((png_name, p))

            # Build PPTX.
            pptx_path = None
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                prs = Presentation()
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
                for name, p in png_paths:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    tb = slide.shapes.add_textbox(
                        Inches(0.3), Inches(0.05),
                        Inches(12.7), Inches(0.4))
                    tb.text_frame.text = name
                    for para in tb.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(16)
                            run.font.bold = True
                    slide.shapes.add_picture(
                        p, Inches(0.3), Inches(0.55),
                        height=Inches(6.85))
                pptx_path = os.path.join(out_dir, "nmf_report.pptx")
                prs.save(pptx_path)
            except ImportError:
                print("[report] python-pptx not installed; "
                      "skipping pptx", flush=True)
            except Exception as e:
                print(f"[report] PPTX failed: {e!r}", flush=True)

            elapsed = time.perf_counter() - t0
            with self._lock:
                bits = [f"{len(png_paths)} pngs"]
                if pptx_path: bits.append("PPTX")
                self._compute_progress = (
                    f"report done in {elapsed:.0f}s.  "
                    + " + ".join(bits) + f"  →  {out_dir}")
        except NMFPanel._Aborted:
            with self._lock:
                self._compute_progress = "report stopped by user."
        except Exception as e:
            err = repr(e)
            print(f"[report] worker failed: {err}", flush=True)
            with self._lock:
                self._compute_progress = f"report failed: {err}"
        finally:
            self._compute_running = False

    def _render_params_image(self, png_dir) -> str | None:
        """Render a small text panel with sample + cube info."""
        try:
            from data import SAMPLES
            cfg = SAMPLES.get(self.sample) or {}
            fig = Figure(figsize=(13, 7), dpi=120, facecolor="white")
            ax = fig.add_subplot(111); ax.set_axis_off()
            lines = [
                f"NMF + clustering benchmark report",
                f"sample: {self.sample}",
                f"scan shape: {self._scan_shape}",
                f"vmax: {cfg.get('vmax')}",
                f"path: {cfg.get('path')}",
                f"timestamp: {time.strftime('%Y-%m-%d %H:%M:%S')}",
                "",
                f"NMF variants:  " + ", ".join(NMF_VARIANTS.keys()),
                f"clustering:    K-means, Aglo (euclidean / cosine / "
                f"polar-xcorr), HDBSCAN",
                f"n_comp:        auto (knee of recon error, n∈[2..14))",
                f"K:             auto (best silhouette, K∈[2..12))",
                f"polar pipeline: CenterCrop({CENTER_CROP}) → "
                f"Resize({POLAR_SIZE}) → Polar({POLAR_SIZE}) → "
                f"MaskLeft({POLAR_MASK_COLS})",
            ]
            ax.text(0.02, 0.98, "\n".join(lines),
                     ha="left", va="top",
                     fontsize=10, family="monospace",
                     transform=ax.transAxes)
            p = os.path.join(png_dir, "00_parameters.png")
            fig.savefig(p, dpi=140, bbox_inches="tight")
            return p
        except Exception as e:
            print(f"[report] params image failed: {e!r}", flush=True)
            return None

    def _render_report_figure(self, png_path, variant_name, cfg,
                               comp_shape, H, W, n_comp, K, err, labels):
        """Render a single variant's components + class maps to PNG."""
        Ny, Nx = (self._scan_shape if self._scan_shape
                    else (1, W.shape[0]))
        fig = Figure(figsize=(13, 8), dpi=120, facecolor="white")
        n_methods = max(1, len(labels))
        gs = fig.add_gridspec(2, max(n_comp, n_methods),
                                height_ratios=[1.0, 1.0])
        # Components
        for k in range(n_comp):
            ax = fig.add_subplot(gs[0, k])
            comp = H[k].reshape(comp_shape)
            if comp.shape[0] == 1:
                ax.plot(comp[0], color="black", lw=1.0)
                ax.set_xticks([]); ax.set_yticks([])
            else:
                ax.imshow(comp, cmap="inferno", aspect="auto")
                ax.set_xticks([]); ax.set_yticks([])
            ax.set_title(f"H[{k}]", fontsize=9)
        # Class maps
        for col, (method, lbl) in enumerate(labels.items()):
            ax = fig.add_subplot(gs[1, col])
            n_classes = int(lbl.max()) + 1
            cmap = plt.get_cmap("tab20" if n_classes > 10 else "tab10")
            palette = ListedColormap(
                [cmap(i % 20) for i in range(n_classes)])
            grid = lbl.reshape(Ny, Nx) if lbl.size == Ny * Nx \
                else lbl.reshape(1, -1)
            ax.imshow(grid, cmap=palette,
                       vmin=-0.5, vmax=n_classes - 0.5,
                       interpolation="nearest", aspect="equal")
            ax.set_title(f"{method}   K={n_classes}", fontsize=10)
            ax.set_axis_off()
            try:
                nm_per_px = float(self.app.real_res.get()) if self.app \
                              else 0.0
            except Exception:
                nm_per_px = 0.0
            if nm_per_px > 0:
                from gui_app._calib_utils import add_real_scalebar
                add_real_scalebar(ax, nm_per_px, length_nm=100,
                                    color="white")
        fig.suptitle(
            f"{self.sample}   variant: {variant_name}   "
            f"n_comp={n_comp}   K={K}   recon-err={err:.3f}",
            fontsize=11)
        fig.tight_layout()
        fig.savefig(png_path, dpi=140, bbox_inches="tight")
