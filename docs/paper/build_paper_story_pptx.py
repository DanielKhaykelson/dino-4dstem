"""Paper-story deck (materials/method audience) — distinct from the chem-PI explainer."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
FIG = os.path.join(ROOT, "docs", "explainer", "figs")
REP = os.path.join(ROOT, "docs", "interpretation_reports")
OUT = os.path.join(os.path.dirname(__file__), "paper_story.pptx")

NAVY=RGBColor(0x21,0x29,0x5C); DEEP=RGBColor(0x06,0x5A,0x82); TEAL=RGBColor(0x1C,0x72,0x93)
ICE=RGBColor(0xCA,0xDC,0xFC); LIGHT=RGBColor(0xF4,0xF7,0xFA); INK=RGBColor(0x1A,0x22,0x30)
WHITE=RGBColor(0xFF,0xFF,0xFF); GOLD=RGBColor(0xE8,0xA3,0x3D); MUTE=RGBColor(0x5A,0x66,0x78)
HEAD="Georgia"; BODY="Calibri"
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=13.333,7.5; BLANK=prs.slide_layouts[6]

def slide(bg=LIGHT):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(SW),Inches(SH))
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return s
def box(s,l,t,w,h,fill=None,line=None,rounded=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                           Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(1)
    shp.shadow.inherit=False; return shp
def text(s,txt,l,t,w,h,size=16,color=INK,bold=False,font=BODY,align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP,italic=False,ls=1.05):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=Pt(2); tf.margin_top=tf.margin_bottom=Pt(2)
    for i,ln in enumerate(txt if isinstance(txt,list) else [txt]):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=ls
        content,opts=(ln if isinstance(ln,tuple) else (ln,{}))
        r=p.add_run(); r.text=content; f=r.font
        f.size=Pt(opts.get("size",size)); f.bold=opts.get("bold",bold)
        f.italic=opts.get("italic",italic); f.name=opts.get("font",font)
        f.color.rgb=opts.get("color",color)
        if opts.get("sb"): p.space_before=Pt(opts["sb"])
    return tb
def fit(s,path,l,t,w,h,cap=None):
    im=Image.open(path); ar=im.size[0]/im.size[1]; bh=h-(0.35 if cap else 0)
    if w/bh>ar: dh=bh; dw=bh*ar
    else: dw=w; dh=w/ar
    s.shapes.add_picture(path,Inches(l+(w-dw)/2),Inches(t+(bh-dh)/2),Inches(dw),Inches(dh))
    if cap: text(s,cap,l,t+bh+0.02,w,0.33,size=10,color=MUTE,align=PP_ALIGN.CENTER,italic=True)
def num(s,n,l,t,d=0.5,fill=DEEP):
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(l),Inches(t),Inches(d),Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb=fill; c.line.fill.background(); c.shadow.inherit=False
    p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=str(n); r.font.size=Pt(20); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=HEAD

# 1 TITLE
s=slide(NAVY); box(s,0,0,0.32,SH,fill=TEAL)
text(s,"Label-free, bias-free classification of difficult 4D-STEM diffraction",
     1.0,1.9,11.3,1.8,size=38,color=WHITE,bold=True,font=HEAD,ls=1.0)
text(s,"A diffraction-tailored self-supervised model — validated against expert "
     "segmentation, and interpretable",1.0,3.9,11,0.9,size=20,color=ICE)
text(s,"IMC · NaPHI · EuInAs   |   no labels · no K · no n-components",1.0,4.7,11,0.5,size=14,color=ICE)
text(s,"Daniel Khaykelson",1.0,6.6,8,0.4,size=13,color=ICE,italic=True)

# 2 PROBLEM
s=slide(LIGHT)
text(s,"The problem",0.7,0.45,12,0.8,size=32,color=DEEP,bold=True,font=HEAD)
text(s,[("Difficult 4D-STEM — beam-sensitive, low-dose, many domains, diffuse "
         "scattering — needs a fast, unbiased first segmentation.",{"size":18}),
        ("But every standard route asks the microscopist to commit upfront:",
         {"size":17,"sb":14,"color":TEAL,"bold":True}),
        ("• k-means / HDBSCAN — choose K (or density)",{"size":16,"sb":8}),
        ("• NMF / PCA — choose n-components",{"size":16,"sb":4}),
        ("• supervised — provide labels (huge effort on hard data)",{"size":16,"sb":4}),
        ("• SAM + orientation mapping — slow, and needs user feature parameters",
         {"size":16,"sb":4}),
        ("Goal: a label-free, parameter-free first map — that we can also trust "
         "and explain physically.",{"size":17,"sb":14,"italic":True,"color":DEEP})],
     0.7,1.5,11.8,5.4)

# 3 METHOD
s=slide(LIGHT)
text(s,"The method — a model engineered for diffraction",0.7,0.45,12.5,0.8,size=30,color=DEEP,bold=True,font=HEAD)
text(s,"Self-distillation (DINO) with design choices specific to NBED — not a borrowed vision model:",
     0.7,1.25,12,0.5,size=15,color=MUTE,italic=True)
rows=[("CNN (ResNet), truncated","right inductive bias + capacity for small, sparse, low-count patterns (ViT is data-hungry, overkill here)"),
      ("Physics-chosen augmentations","azimuthal theta-roll -> in-plane-rotation invariance; flips / colour-jitter disabled (meaningless for diffraction)"),
      ("Radial-gated SupCon + 1-D loss","grouping biased toward the physically meaningful radial / ring signature"),
      ("Emergent K + optional semi-supervision","over-specify prototypes; the model uses what it needs. A few pairwise hints can steer it if wanted")]
y=1.95
for i,(h,d) in enumerate(rows):
    num(s,i+1,0.8,y+0.03); text(s,h,1.55,y,3.7,0.9,size=17,color=INK,bold=True,font=HEAD)
    text(s,d,5.4,y-0.02,7.3,1.1,size=14,color=INK); y+=1.25

# 4 VALIDATION (NaPHI)
s=slide(LIGHT)
text(s,"It gets the known answer — label-free",0.7,0.4,12.5,0.8,size=30,color=DEEP,bold=True,font=HEAD)
fit(s,os.path.join(FIG,"naphi_concordance.png"),0.5,1.2,12.4,4.2,
    cap="NaPHI: our class map vs our previously-published SAM+orientation map (same scan, aligned).")
box(s,0.7,5.55,11.9,1.5,fill=NAVY,rounded=True)
text(s,[("DINO reproduces the published 'Line' domain  —  IoU 0.68, Dice 0.81, ~75% recovered  —  "
         "with NO labels, NO SAM, and NO user line-parameters.",{"size":16,"color":WHITE,"bold":True}),
        ("...and transfers to a new (tilted) acquisition of the same flake in milliseconds.",
         {"size":14,"color":ICE,"sb":6})],0.95,5.68,11.4,1.3,anchor=MSO_ANCHOR.MIDDLE)

# 5 PAYOFF (IMC)
s=slide(LIGHT)
text(s,"It cracks the hard case",0.7,0.4,12.5,0.8,size=30,color=DEEP,bold=True,font=HEAD)
fit(s,os.path.join(REP,"IMC_SI5__m0.9700_seed42_K60","classical_vs_dino_maps.png"),0.5,1.35,12.4,3.2,
    cap="IMC (low-dose molecular): DINO (far left) vs classical clusterings of the same data.")
box(s,0.7,4.95,11.9,2.0,fill=ICE,rounded=True)
text(s,[("On the hard molecular films, classical methods (virtual-DF / radial / PCA / NMF) "
         "do NOT reproduce the DINO map (best agreement ARI ~ 0.11-0.21).",{"size":16,"color":NAVY,"bold":True}),
        ("There is no external ground truth here — classical analysis simply cannot give a coherent "
         "map. Our method is the only route, and (next slide) we prove the map is physical, not noise.",
         {"size":15,"color":INK,"sb":8})],0.95,5.1,11.4,1.7,anchor=MSO_ANCHOR.MIDDLE)

# 6 INTERPRETABILITY
s=slide(LIGHT)
text(s,"What the clusters mean (and when we add value)",0.7,0.4,12.6,0.8,size=28,color=DEEP,bold=True,font=HEAD)
fit(s,os.path.join(FIG,"cross_all.png"),0.4,1.2,12.6,3.2,
    cap="(a) what's encoded (probe R2)  (b) causal ablations  (c) classical agreement.  IMC SI3/SI4/SI5 - NaPHI - EuInAs")
fit(s,os.path.join(FIG,"ablation_demo_row.png"),0.4,4.55,8.0,2.2)
box(s,8.7,4.6,4.0,2.3,fill=NAVY,rounded=True)
text(s,[("Built-in interpretability protocol",{"size":15,"color":GOLD,"bold":True}),
        ("Clusters are driven by scattered intensity + low-q 2-D structure; rotation-invariant; "
         "NOT orientation (molecular).",{"size":13,"color":WHITE,"sb":6}),
        ("Uniqueness is material-dependent (panel c).",{"size":13,"color":ICE,"sb":6})],
     8.95,4.72,3.5,2.05,anchor=MSO_ANCHOR.MIDDLE)

# 7 GENERALITY (EuInAs)
s=slide(LIGHT)
text(s,"It generalises — crystalline, multi-phase, orientation",0.7,0.4,12.6,0.8,size=28,color=DEEP,bold=True,font=HEAD)
fit(s,os.path.join(REP,"EuInAs_B100__m0.9700_seed42_K60","classical_vs_dino_maps.png"),0.5,1.4,12.4,3.0,
    cap="EuInAs (layered, 3-phase): DINO vs classical clusterings.")
box(s,0.7,4.85,11.9,2.0,fill=ICE,rounded=True)
text(s,[("On the well-ordered EuInAs, classical methods get about halfway (ARI ~ 0.44) — DINO stays "
         "closer to classical here (it is most distinctive on the HARD molecular samples).",{"size":15,"color":NAVY,"bold":True}),
        ("And 3-phase ACOM shows DINO additionally tracks crystal orientation here (zone-axis AMI ~ 0.30) "
         "— a genuine secondary axis the molecular films don't have.",{"size":14,"color":INK,"sb":8})],
     0.95,5.0,11.4,1.7,anchor=MSO_ANCHOR.MIDDLE)

# 8 SUMMARY / CONCLUSIONS
s=slide(NAVY); box(s,0,0,0.32,SH,fill=GOLD)
text(s,"What we deliver",1.0,0.5,11,0.9,size=32,color=WHITE,bold=True,font=HEAD)
text(s,[("1.  A diffraction-tailored self-supervised classifier — label-free, no K, "
         "rotation-invariant, with optional semi-supervision.",{"size":18,"color":WHITE}),
        ("2.  Validation against our published expert segmentation (NaPHI), label-free and "
         "without SAM; ms transfer to new data.",{"size":18,"color":WHITE,"sb":12}),
        ("3.  A built-in interpretability protocol that proves the clusters are physical "
         "(scattered intensity + low-q structure) and quantifies when the model beats classical "
         "(hard molecular IMC) vs converges with it (ordered NaPHI / EuInAs).",{"size":18,"color":WHITE,"sb":12}),
        ("Most valuable exactly where existing tools fail: low-dose, diffuse, molecular crystals.",
         {"size":17,"color":GOLD,"italic":True,"sb":14})],1.0,1.7,11.5,5.2,ls=1.1)

prs.save(OUT); print("saved",OUT,"slides:",len(prs.slides._sldIdLst))
