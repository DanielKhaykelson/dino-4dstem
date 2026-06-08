"""Build the interpretation explainer deck (chemistry-PI audience)."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIG = os.path.join(os.path.dirname(__file__), "figs")
OUT = os.path.join(os.path.dirname(__file__), "interpretation_explainer.pptx")

# cropped class-means (strip the raw-filename suptitle at the top)
_cm = Image.open(os.path.join(FIG, "imc_class_means.png"))
_cm.crop((0, int(_cm.height * 0.075), _cm.width, _cm.height)).save(
    os.path.join(FIG, "imc_class_means_crop.png"))

NAVY = RGBColor(0x21, 0x29, 0x5C)
DEEP = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x1C, 0x72, 0x93)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
INK = RGBColor(0x1A, 0x22, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xE8, 0xA3, 0x3D)
MUTE = RGBColor(0x5A, 0x66, 0x78)

HEAD = "Georgia"
BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                           Inches(SW), Inches(SH))
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def box(s, l, t, w, h, fill=None, line=None, rounded=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(s, txt, l, t, w, h, size=16, color=INK, bold=False, font=BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
         line_spacing=1.05):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    lines = txt if isinstance(txt, list) else [txt]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        if isinstance(ln, tuple):
            content, opts = ln
        else:
            content, opts = ln, {}
        r = p.add_run(); r.text = content
        f = r.font
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.italic = opts.get("italic", italic)
        f.name = opts.get("font", font)
        f.color.rgb = opts.get("color", color)
        if opts.get("space_before"):
            p.space_before = Pt(opts["space_before"])
    return tb


def fit_image(s, path, l, t, w, h, caption=None, cap_color=MUTE):
    im = Image.open(path); iw, ih = im.size; ar = iw / ih
    bw, bh = w, h - (0.35 if caption else 0)
    if bw / bh > ar:          # box wider than image → fit height
        dh = bh; dw = bh * ar
    else:
        dw = bw; dh = bw / ar
    il = l + (w - dw) / 2; it = t + (bh - dh) / 2
    s.shapes.add_picture(path, Inches(il), Inches(it),
                         Inches(dw), Inches(dh))
    if caption:
        text(s, caption, l, t + bh + 0.02, w, 0.33, size=11, color=cap_color,
             align=PP_ALIGN.CENTER, italic=True)


def numbered(s, n, l, t, d=0.55, fill=DEEP):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t),
                           Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(22); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = HEAD


# ---------------------------------------------------------------- 1 TITLE
s = slide(NAVY)
box(s, 0, 0, 0.32, SH, fill=TEAL)
text(s, "What is the model actually looking at?", 1.0, 2.05, 11.3, 2.0,
     size=42, color=WHITE, bold=True, font=HEAD, line_spacing=1.0)
text(s, "Interpreting self-supervised 4D-STEM diffraction clustering",
     1.0, 4.05, 11.0, 0.7, size=21, color=ICE, font=BODY)
text(s, [("Making the DINO “grain map” explainable in physical terms "
          "— IMC & NaPHI", {})],
     1.0, 4.75, 11.0, 0.5, size=14, color=ICE)
text(s, "Daniel Khaykelson", 1.0, 6.5, 8, 0.4, size=13, color=ICE, italic=True)

# ---------------------------------------------------------------- 2 QUESTION
s = slide(LIGHT)
text(s, "The setup — and the problem", 0.7, 0.45, 12, 0.8, size=32,
     color=DEEP, bold=True, font=HEAD)
text(s, [
    ("What the model does", {"bold": True, "size": 18, "color": TEAL}),
    ("At every probe position it reads the diffraction pattern and sorts the "
     "position into one of 14 groups — with no labels and no chemistry "
     "input. The result is a “grain map.”", {"size": 16,
     "space_before": 4}),
    ("The problem", {"bold": True, "size": 18, "color": TEAL,
     "space_before": 14}),
    ("It works, but it is a black box: we do not know which physical property "
     "defines the groups — crystallinity? thickness? crystal "
     "orientation? something new?", {"size": 16, "space_before": 4}),
    ("Goal: explain the groups in physical terms a chemist can trust.",
     {"size": 16, "italic": True, "color": DEEP, "space_before": 14}),
], 0.7, 1.5, 6.0, 5.2)
box(s, 7.0, 1.5, 5.6, 4.9, fill=WHITE, line=ICE, rounded=True)
fit_image(s, os.path.join(FIG, "imc_class_means_crop.png"), 7.2, 1.7, 5.2, 4.5,
          caption="Each group's average diffraction pattern (IMC, 14 groups).")

# ---------------------------------------------------------------- 3 TOOLBOX
s = slide(LIGHT)
text(s, "Four simple tests", 0.7, 0.45, 12, 0.8, size=32, color=DEEP,
     bold=True, font=HEAD)
text(s, "Each is independent; together they triangulate what the model uses.",
     0.7, 1.25, 12, 0.5, size=15, color=MUTE, italic=True)
rows = [
    ("Probing", "Can we predict a known physical property from the model's "
     "internal fingerprint? If yes, it pays attention to that property."),
    ("Ablation", "Erase one ingredient of the diffraction pattern, re-run the "
     "model, and see whether the map survives."),
    ("Classical check", "Could a standard 4D-STEM method (virtual detector, "
     "PCA, NMF) have drawn the same map — or is this map new?"),
    ("Orientation check", "Is it merely sorting by crystal orientation? "
     "(plus Grad-CAM: which part of the pattern it looks at.)"),
]
y = 1.95
for i, (h, d) in enumerate(rows):
    numbered(s, i + 1, 0.8, y + 0.05)
    text(s, h, 1.6, y, 3.0, 0.5, size=18, color=INK, bold=True, font=HEAD)
    text(s, d, 4.7, y - 0.02, 7.9, 1.1, size=15, color=INK)
    y += 1.25

# ---------------------------------------------------------------- 4 PROBING
s = slide(LIGHT)
text(s, "Test 1 · Probing — what is encoded?", 0.7, 0.45, 12, 0.8,
     size=30, color=DEEP, bold=True, font=HEAD)
text(s, [
    ("In plain terms", {"bold": True, "size": 18, "color": TEAL}),
    ("If we can “read off” a property (say crystallinity) from the "
     "model's notes for each point, then the model is keeping track of it.",
     {"size": 16, "space_before": 4}),
    ("How it is measured", {"bold": True, "size": 18, "color": TEAL,
     "space_before": 14}),
    ("Fit a simple linear predictor from the model's 128-number fingerprint "
     "to each measured property. The score R² (0–1) says how "
     "predictable — i.e. how strongly encoded — that property is.",
     {"size": 16, "space_before": 4}),
    ("Properties tested: total scattered intensity, crystallinity, "
     "spottiness, and crystal orientation (ACOM).", {"size": 16,
     "italic": True, "color": DEEP, "space_before": 14}),
], 0.7, 1.5, 7.3, 5.2)
box(s, 8.3, 1.7, 4.3, 3.4, fill=NAVY, rounded=True)
text(s, [
    ("R² ≈ 1", {"size": 40, "bold": True, "color": GOLD,
     "font": HEAD}),
    ("strongly encoded", {"size": 16, "color": WHITE}),
    ("R² ≈ 0", {"size": 40, "bold": True, "color": ICE,
     "font": HEAD, "space_before": 18}),
    ("ignored by the model", {"size": 16, "color": WHITE}),
], 8.3, 2.1, 4.3, 2.8, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, "Result preview: scattered intensity & crystallinity score high; "
        "orientation scores ~0.", 8.3, 5.35, 4.3, 1.2, size=13, color=MUTE,
     italic=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- 5 ABLATION
s = slide(LIGHT)
text(s, "Test 2 · Ablation — what does it depend on?", 0.7, 0.45,
     12.4, 0.8, size=30, color=DEEP, bold=True, font=HEAD)
text(s, [
    ("In plain terms", {"bold": True, "size": 17, "color": TEAL}),
    ("Like covering part of a photo to see if you still recognise it.",
     {"size": 15, "space_before": 3}),
    ("How it is measured", {"bold": True, "size": 17, "color": TEAL,
     "space_before": 10}),
    ("Blank out one ingredient — overall brightness, fine angular "
     "detail, or inner vs. outer rings — re-run the model, and compare "
     "the new map to the original (agreement = ARI).", {"size": 15,
     "space_before": 3}),
    ("Map survives the edit → that ingredient was not needed. "
     "Map falls apart → the model depended on it.", {"size": 15,
     "italic": True, "color": DEEP, "space_before": 10}),
], 0.7, 1.45, 5.3, 5.4)
box(s, 6.2, 1.5, 6.5, 3.0, fill=WHITE, line=ICE, rounded=True)
fit_image(s, os.path.join(FIG, "imc_ablation_maps.png"), 6.35, 1.62, 6.2, 2.9,
          caption="Re-inference after each edit (IMC). Low agreement = the "
                  "model relied on what was removed.")
box(s, 6.2, 4.95, 6.5, 1.75, fill=ICE, rounded=True)
text(s, [
    ("What we see:", {"bold": True, "size": 15, "color": NAVY}),
    ("• Remove overall brightness → map collapses (intensity "
     "matters).", {"size": 14}),
    ("• Average away the angles → collapses (needs the 2-D pattern, "
     "not just a 1-D profile).", {"size": 14}),
    ("• Inner (low-angle) rings matter far more than outer rings.",
     {"size": 14}),
], 6.45, 5.05, 6.1, 1.6, color=INK)

# ---------------------------------------------------------------- 6 CLASSICAL
s = slide(LIGHT)
text(s, "Test 3 · Could a classical method do the same?", 0.7, 0.45,
     12.4, 0.8, size=30, color=DEEP, bold=True, font=HEAD)
text(s, [
    ("In plain terms", {"bold": True, "size": 17, "color": TEAL}),
    ("Does the AI just redraw what an old-school analysis already shows, or "
     "something genuinely new?", {"size": 15, "space_before": 3}),
    ("How it is measured", {"bold": True, "size": 17, "color": TEAL,
     "space_before": 10}),
    ("Cluster classical features — virtual dark-field, azimuthal "
     "profile, PCA / NMF of the pattern — into the same number of "
     "groups, and measure agreement with the AI map (ARI / AMI).",
     {"size": 15, "space_before": 3}),
], 0.7, 1.45, 5.3, 3.6)
box(s, 0.7, 5.15, 5.3, 1.75, fill=NAVY, rounded=True)
text(s, [
    ("Uniqueness is material-dependent", {"size": 15, "bold": True,
     "color": GOLD}),
    ("IMC: classical agreement ARI ≈ 0.15 → the AI map is distinctive "
     "(classical methods miss the fine grains).", {"size": 14, "color": WHITE,
     "space_before": 6}),
    ("NaPHI: ARI ≈ 0.5 → about half-classical.", {"size": 14,
     "color": ICE, "space_before": 6}),
], 0.95, 5.3, 4.85, 1.5, anchor=MSO_ANCHOR.MIDDLE)
box(s, 6.3, 2.05, 6.4, 3.4, fill=WHITE, line=ICE, rounded=True)
fit_image(s, os.path.join(FIG, "imc_dino_vs_classical.png"), 6.5, 2.25, 6.0,
          3.0, caption="IMC: DINO (far left) vs. classical clusterings. "
                       "Colours are arbitrary; shape match is what counts.")

# ---------------------------------------------------------------- 7 NUMBERS
s = slide(LIGHT)
text(s, "How to read the two numbers", 0.7, 0.45, 12, 0.8, size=32,
     color=DEEP, bold=True, font=HEAD)
cards = [
    ("R²", "How predictable a property is from the model's fingerprint.",
     "0 = the model ignores it     →     1 = fully encoded", DEEP),
    ("ARI / AMI", "How much two maps agree, correcting for chance.",
     "0 = no better than random     →     1 = identical maps", TEAL),
]
y = 1.9
for tag, sub, scale, col in cards:
    box(s, 0.9, y, 11.5, 2.1, fill=WHITE, line=ICE, rounded=True)
    box(s, 0.9, y, 0.22, 2.1, fill=col)
    text(s, tag, 1.3, y, 3.4, 2.1, size=34, color=col, bold=True,
         font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
    text(s, [
        (sub, {"size": 16, "bold": True, "color": INK}),
        (scale, {"size": 16, "color": MUTE, "space_before": 8}),
    ], 4.9, y, 7.3, 2.1, anchor=MSO_ANCHOR.MIDDLE)
    y += 2.45

# ---------------------------------------------------------------- 8 FINDINGS
s = slide(LIGHT)
text(s, "What we found", 0.7, 0.45, 12, 0.8, size=32, color=DEEP, bold=True,
     font=HEAD)
text(s, [
    ("• The groups are driven by the overall scattered intensity plus "
     "the inner, low-angle (low-q) 2-D diffraction structure.", {"size": 16}),
    ("• NOT by crystal orientation (agreement ≈ 0) — the model "
     "is rotation-invariant by design.", {"size": 16, "space_before": 8}),
    ("• Uniqueness is material-dependent: IMC maps are distinctive "
     "(classical ≈ 0.15); NaPHI is ~half-classical (≈ 0.5).",
     {"size": 16, "space_before": 8}),
    ("• Log-stretch preprocessing keeps the same basis but sharpens it "
     "— more, cleaner groups.", {"size": 16, "space_before": 8}),
], 0.7, 1.5, 6.1, 5.0)
box(s, 7.0, 1.6, 5.7, 4.9, fill=WHITE, line=ICE, rounded=True)
fit_image(s, os.path.join(FIG, "cross_sample.png"), 7.15, 2.6, 5.4, 2.9,
          caption="Cross-sample: (a) what's encoded, (b) ablations, "
                  "(c) classical agreement.")
text(s, "Across IMC SI3/SI4/SI5 and NaPHI", 7.15, 1.75, 5.4, 0.5, size=14,
     color=MUTE, italic=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- 9 MEANING
s = slide(NAVY)
box(s, 0, 0, 0.32, SH, fill=GOLD)
text(s, "What it means", 1.0, 0.6, 11, 0.9, size=34, color=WHITE, bold=True,
     font=HEAD)
text(s, [
    ("The model groups regions by how much, and what, they diffract at low "
     "angles — crystallinity, molecular packing, and thickness. A "
     "physically sensible axis, not an arbitrary black box.", {"size": 19,
     "color": WHITE}),
    ("For IMC it is genuinely distinctive (a standard virtual-detector / PCA "
     "map does not reproduce it); for NaPHI it stays closer to classical.",
     {"size": 17, "color": ICE, "space_before": 16}),
    ("Caveat: “scattered intensity” blends crystallinity with "
     "sample thickness — a thickness map would separate the two.",
     {"size": 15, "color": ICE, "italic": True, "space_before": 16}),
], 1.0, 1.9, 11.4, 4.8, line_spacing=1.1)

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
